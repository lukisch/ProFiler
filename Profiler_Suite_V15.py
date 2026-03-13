#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import os
import json
import uuid
import hashlib
import shutil
import sqlite3
import time
import subprocess
import traceback
from datetime import datetime, timedelta
from pathlib import Path

from PyQt6.QtWidgets import (
    QApplication, QMainWindow, QWidget, QVBoxLayout, QListWidget, QListWidgetItem,
    QMenu, QHBoxLayout, QPushButton, QProgressBar, QLabel,
    QDialog, QFormLayout, QLineEdit, QComboBox, QCheckBox, QDialogButtonBox,
    QFileDialog, QTabWidget, QSplitter, QTextEdit, QSystemTrayIcon, QStyle, QMessageBox,
    QTreeWidget, QTreeWidgetItem, QInputDialog, QGroupBox, QRadioButton, QButtonGroup,
    QFileIconProvider, QSpinBox, QScrollArea
)
from PyQt6.QtCore import (
    Qt, QThread, pyqtSignal, QObject, QTimer, 
    QSize, QFileInfo
)
from PyQt6.QtGui import QAction, QPalette, QColor, QFont, QPixmap, QIcon, QImage

# Optionale Bibliotheken
try:
    import PyPDF2
    from PyPDF2 import PdfReader, PdfWriter
    HAS_PDF = True
except ImportError:
    HAS_PDF = False

try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    import pytesseract
    from PIL import Image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    from pdf2image import convert_from_path
    HAS_PDF2IMAGE = True
except ImportError:
    HAS_PDF2IMAGE = False

try:
    import fitz  # PyMuPDF for redaction
    HAS_FITZ = True
except ImportError:
    HAS_FITZ = False

try:
    from watchdog.observers import Observer
    from watchdog.events import FileSystemEventHandler
    HAS_WATCHDOG = True
except ImportError:
    HAS_WATCHDOG = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak, Table, TableStyle
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    HAS_REPORTLAB = True
except ImportError:
    HAS_REPORTLAB = False

try:
    import pikepdf
    HAS_PIKEPDF = True
except ImportError:
    HAS_PIKEPDF = False

# ============================================================================
# ENCODING SETUP
# ============================================================================

def setup_windows_encoding():
    """
    Konfiguriert UTF-8 Encoding für Windows-Konsole.

    Verhindert Encoding-Probleme bei deutschen Umlauten und Sonderzeichen.
    """
    if sys.platform == 'win32':
        if sys.stdout.encoding != 'utf-8':
            sys.stdout.reconfigure(encoding='utf-8')
        if sys.stderr.encoding != 'utf-8':
            sys.stderr.reconfigure(encoding='utf-8')

# Encoding-Setup beim Import ausführen
setup_windows_encoding()

print("ProFiler Suite V15 startet...")
print("Encoding Check:", sys.stdout.encoding)

# ============================================================================
# 1. SHARED UTILS & CONFIG
# ============================================================================

_CONFIG_DIR = Path.home() / ".profiler_suite"
SEARCH_CONFIG_PATH = str(_CONFIG_DIR / "search_config.json")
SYNC_CONFIG_PATH = str(_CONFIG_DIR / "profiler_config.json")
SETTINGS_PATH = str(_CONFIG_DIR / "profiler_settings.json")

# Konstanten für File Processing
DEFAULT_CHUNK_SIZE = 1024 * 1024  # 1 MB für Hash-Berechnung
PDF_TEXT_MIN_CHARS = 20  # Minimale Zeichenanzahl für Text-Erkennung
PDF_PAGES_TO_CHECK = 3   # Anzahl zu prüfender Seiten in PDFs
FILE_ATTR_CLOUD_PLACEHOLDER = 0x1000  # Windows File-Attribut für Cloud-Placeholder
FILE_ATTR_RECALL_ON_ACCESS = 0x400    # Windows File-Attribut für Recall-on-Access

def sha256_file(path, chunk_size=DEFAULT_CHUNK_SIZE):
    """Berechnet den SHA256 Hash. Robust gegen leere Dateien."""
    h = hashlib.sha256()
    try:
        with open(path, "rb") as f:
            while True:
                chunk = f.read(chunk_size)
                if not chunk: break
                h.update(chunk)
    except (PermissionError, OSError):
        return None
    return h.hexdigest()

def is_cloud_placeholder(path):
    """
    Prüft ob eine Datei ein Cloud-Placeholder ist (OneDrive/Dropbox).

    Args:
        path: Dateipfad

    Returns:
        bool: True wenn Placeholder, False sonst
    """
    if os.name != 'nt':
        return False
    try:
        attrs = os.stat(path).st_file_attributes
        return (attrs & FILE_ATTR_CLOUD_PLACEHOLDER) or (attrs & FILE_ATTR_RECALL_ON_ACCESS)
    except (OSError, AttributeError):
        return False

def shorten_filename(name, max_len):
    """
    Kürzt einen Dateinamen auf maximale Länge.

    Args:
        name: Dateiname
        max_len: Maximale Länge

    Returns:
        str: Gekürzter Dateiname mit "..." in der Mitte
    """
    if len(name) <= max_len:
        return name
    root, ext = os.path.splitext(name)
    keep = max(1, max_len - len(ext) - 3)
    return root[:keep] + "..." + ext

def get_file_category(filename):
    """
    Kategorisiert Dateien nach Extension.

    Args:
        filename: Dateiname mit Extension

    Returns:
        str: Kategorie (Dokumente, Bilder, Audio, Video, Archive, Code, Tabellen, Andere)
    """
    ext = os.path.splitext(filename)[1].lower()
    if ext in ['.pdf', '.doc', '.docx', '.txt', '.md', '.rtf', '.odt']: return "Dokumente"
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.svg', '.webp']: return "Bilder"
    if ext in ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a']: return "Audio"
    if ext in ['.mp4', '.mkv', '.avi', '.mov', '.wmv']: return "Video"
    if ext in ['.zip', '.rar', '.7z', '.tar', '.gz']: return "Archive"
    if ext in ['.py', '.js', '.html', '.css', '.json', '.xml', '.sql', '.cpp', '.c', '.h']: return "Code"
    if ext in ['.xls', '.xlsx', '.csv']: return "Tabellen"
    return "Andere"

def is_pdf_encrypted(filepath):
    """
    Prüft ob eine PDF-Datei verschlüsselt ist.

    Args:
        filepath: Pfad zur PDF-Datei

    Returns:
        bool: True wenn verschlüsselt, False sonst
    """
    if not HAS_PDF:
        return False
    try:
        with open(filepath, 'rb') as f:
            pdf = PdfReader(f)
            return pdf.is_encrypted
    except (OSError, Exception):
        return False

def has_pdf_text(filepath):
    """
    Prüft ob eine PDF-Datei extrahierbaren Text enthält.

    Scannt bis zu 3 Seiten und prüft ob mindestens 20 Zeichen Text gefunden werden.
    Verhindert falsche OCR-Erkennung bei reinen Bild-PDFs.

    Args:
        filepath: Pfad zur PDF-Datei

    Returns:
        bool: True wenn Text vorhanden, False wenn nur Bilder
    """
    if not HAS_PDF:
        return False
    try:
        with open(filepath, 'rb') as f:
            pdf = PdfReader(f)
            if len(pdf.pages) == 0:
                return False
            
            # Prüfe bis zu PDF_PAGES_TO_CHECK Seiten für bessere Erkennung
            pages_to_check = min(PDF_PAGES_TO_CHECK, len(pdf.pages))

            for i in range(pages_to_check):
                text = pdf.pages[i].extract_text() or ""
                # Reduzierte Schwelle: PDF_TEXT_MIN_CHARS statt 50 Zeichen
                if len(text.strip()) > PDF_TEXT_MIN_CHARS:
                    return True
            
            return False
    except (OSError, Exception):
        return False


def find_tool_path(tool_name):
    """Sucht nach einem Tool im Projekt-Verzeichnis"""
    # 1. Im selben Verzeichnis wie das Hauptskript
    script_dir = os.path.dirname(os.path.abspath(__file__))
    tool_path = os.path.join(script_dir, tool_name)
    if os.path.exists(tool_path):
        return tool_path

    # 2. Im Parent-Verzeichnis
    parent_dir = os.path.dirname(script_dir)
    parent_path = os.path.join(parent_dir, tool_name)
    if os.path.exists(parent_path):
        return parent_path

    # Nicht gefunden
    return None

# --- CONFIG MANAGERS ---

class SearchConfigManager:
    def __init__(self):
        self.dbs = []
        self.load()
    def load(self):
        if os.path.exists(SEARCH_CONFIG_PATH):
            try:
                with open(SEARCH_CONFIG_PATH, "r", encoding="utf-8") as f:
                    self.dbs = json.load(f).get("databases", [])
            except (OSError, json.JSONDecodeError, KeyError):
                self.save()
        else:
            self.save()
    def save(self):
        os.makedirs(os.path.dirname(SEARCH_CONFIG_PATH) or ".", exist_ok=True)
        with open(SEARCH_CONFIG_PATH, "w", encoding="utf-8") as f:
            json.dump({"databases": self.dbs}, f, indent=2)
    def add_db(self, path):
        if path and path not in self.dbs: self.dbs.append(path); self.save()
    def remove_db(self, path):
        self.dbs = [d for d in self.dbs if d != path]; self.save()

class SyncConfigManager:
    def __init__(self, path):
        self.path = path
        self.data = {"connections": []}
        self.load()
    def load(self):
        if os.path.exists(self.path):
            try:
                with open(self.path, "r", encoding="utf-8") as f:
                    self.data = json.load(f)
            except (OSError, json.JSONDecodeError):
                self.save()
        else:
            self.save()
    def save(self):
        os.makedirs(os.path.dirname(self.path) or ".", exist_ok=True)
        with open(self.path, "w", encoding="utf-8") as f: json.dump(self.data, f, indent=2)
    def list_connections(self): return self.data.get("connections", [])
    def add_or_update_connection(self, conn):
        conns = self.data.get("connections", [])
        found = False
        for i, c in enumerate(conns):
            if c.get("id") == conn.get("id"):
                conns[i] = conn; found = True; break
        if not found: conns.append(conn)
        self.data["connections"] = conns; self.save()
    def remove_connection(self, conn_id):
        self.data["connections"] = [c for c in self.data.get("connections", []) if c.get("id") != conn_id]
        self.save()

class SettingsManager:
    """Verwaltet App-Einstellungen (Lösch-Verhalten, PDF-Passwrter, etc.)"""
    def __init__(self):
        self.data = {
            "delete_mode": "soft",
            "trash_retention_days": 30,
            "auto_cleanup_enabled": True,
            "pdf_master_password_open": "",  # Masterpasswort zum öffnen
            "pdf_master_password_save": "",  # Masterpasswort zum Speichern
            "ocr_language": "deu",  # Tesseract Language
            "ocr_enabled": True
        }
        self.load()
    
    def load(self):
        if os.path.exists(SETTINGS_PATH):
            try:
                with open(SETTINGS_PATH, "r", encoding="utf-8") as f:
                    loaded = json.load(f)
                    self.data.update(loaded)
            except (OSError, json.JSONDecodeError):
                self.save()
        else:
            self.save()
    
    def save(self):
        os.makedirs(os.path.dirname(SETTINGS_PATH) or ".", exist_ok=True)
        with open(SETTINGS_PATH, "w", encoding="utf-8") as f:
            json.dump(self.data, f, indent=2)
    
    def get(self, key, default=None):
        return self.data.get(key, default)
    
    def set(self, key, value):
        self.data[key] = value
        self.save()

# ============================================================================
# AUTO-SYNC WATCHDOG (V14.3)
# ============================================================================

class AutoSyncHandler(FileSystemEventHandler if HAS_WATCHDOG else object):
    """
    Überwacht einen Ordner auf neue/geänderte Dateien und synchronisiert automatisch.
    Benötigt: pip install watchdog
    """
    def __init__(self, target_folder, extensions=None, callback=None):
        if HAS_WATCHDOG:
            super().__init__()
        self.target = target_folder
        self.extensions = extensions or ['.pdf', '.docx', '.xlsx', '.txt', '.jpg', '.png']
        self.callback = callback  # UI-Callback für Status-Updates
        self.sync_count = 0
        self.last_sync = None
    
    def on_created(self, event):
        """Wird aufgerufen wenn eine neue Datei erstellt wird."""
        if event.is_directory:
            return
        if self._should_sync(event.src_path):
            self._sync_file(event.src_path, "created")
    
    def on_modified(self, event):
        """Wird aufgerufen wenn eine Datei geändert wird."""
        if event.is_directory:
            return
        if self._should_sync(event.src_path):
            self._sync_file(event.src_path, "modified")
    
    def _should_sync(self, path):
        """Prüft ob die Datei synchronisiert werden soll (basierend auf Extension)."""
        return any(path.lower().endswith(ext) for ext in self.extensions)
    
    def _sync_file(self, src_path, event_type):
        """Kopiert die Datei in den Zielordner."""
        try:
            filename = os.path.basename(src_path)
            dest_path = os.path.join(self.target, filename)
            
            # Duplikat-Handling
            if os.path.exists(dest_path):
                base, ext = os.path.splitext(filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                dest_path = os.path.join(self.target, f"{base}_{timestamp}{ext}")
            
            shutil.copy2(src_path, dest_path)
            self.sync_count += 1
            self.last_sync = datetime.now()
            
            if self.callback:
                self.callback(f"Sync: {filename} ({event_type})")
                
        except Exception as e:
            if self.callback:
                self.callback(f"Sync-Fehler: {e}")

class AutoSyncManager(QObject):
    """
    Verwaltet Watchdog-Observer für Auto-Sync.
    Kann mehrere Watch-Ordner gleichzeitig überwachen.
    """
    status_changed = pyqtSignal(str)
    file_synced = pyqtSignal(str)
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.observers = {}  # source_path -> Observer
        self.handlers = {}   # source_path -> AutoSyncHandler
        self.active = False
    
    def start_watch(self, source_folder, target_folder, extensions=None):
        """Startet die Überwachung eines Ordners."""
        if not HAS_WATCHDOG:
            self.status_changed.emit("Watchdog nicht installiert (pip install watchdog)")
            return False
        
        if source_folder in self.observers:
            self.status_changed.emit(f"Ordner wird bereits überwacht: {source_folder}")
            return False
        
        handler = AutoSyncHandler(
            target_folder, 
            extensions,
            callback=lambda msg: self.file_synced.emit(msg)
        )
        observer = Observer()
        observer.schedule(handler, source_folder, recursive=False)
        observer.start()
        
        self.observers[source_folder] = observer
        self.handlers[source_folder] = handler
        self.active = True
        
        self.status_changed.emit(f"Überwachung gestartet: {source_folder} → {target_folder}")
        return True
    
    def stop_watch(self, source_folder=None):
        """Stoppt die Überwachung (eines oder aller Ordner)."""
        if source_folder:
            if source_folder in self.observers:
                self.observers[source_folder].stop()
                self.observers[source_folder].join()
                del self.observers[source_folder]
                del self.handlers[source_folder]
                self.status_changed.emit(f"Überwachung gestoppt: {source_folder}")
        else:
            # Alle stoppen
            for path, obs in self.observers.items():
                obs.stop()
                obs.join()
            self.observers.clear()
            self.handlers.clear()
            self.status_changed.emit("Alle Überwachungen gestoppt")
        
        self.active = len(self.observers) > 0
    
    def get_status(self):
        """Gibt den aktuellen Status zurück."""
        if not self.observers:
            return "Inaktiv"
        return f"Aktiv: {len(self.observers)} Ordner überwacht"
    
    def get_stats(self):
        """Gibt Statistiken zurück."""
        total_synced = sum(h.sync_count for h in self.handlers.values())
        return {
            "watched_folders": len(self.observers),
            "total_synced": total_synced,
            "active": self.active
        }

# ============================================================================
# CONNECTIONS DATABASE (from ProFiler V4 + multi-folder support)
# ============================================================================

DDL_SCHEMA = """
CREATE TABLE IF NOT EXISTS files(
    id INTEGER PRIMARY KEY, 
    content_hash TEXT UNIQUE, 
    size INTEGER, 
    mime TEXT, 
    first_seen TEXT
);
CREATE TABLE IF NOT EXISTS versions(
    id INTEGER PRIMARY KEY, 
    file_id INTEGER, 
    name TEXT, 
    path TEXT, 
    mtime TEXT, 
    ctime TEXT, 
    version_index INTEGER, 
    source_folder TEXT
);
CREATE TABLE IF NOT EXISTS tags(
    id INTEGER PRIMARY KEY, 
    file_id INTEGER, 
    tag TEXT
);
CREATE TABLE IF NOT EXISTS events(
    id INTEGER PRIMARY KEY, 
    file_id INTEGER, 
    event_type TEXT, 
    details TEXT, 
    ts TEXT
);
"""


# ============================================================================
# 2. PDF UTILITY FUNCTIONS
# ============================================================================

class PDFUtils:
    """PDF-spezifische Utility-Funktionen"""
    
    @staticmethod
    def encrypt_pdf(input_path, output_path, password):
        """Verschlüsselt ein PDF mit Passwort"""
        if not HAS_PDF:
            raise Exception("PyPDF2 nicht installiert")
        
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()
            
            # Alle Seiten kopieren
            for page in reader.pages:
                writer.add_page(page)
            
            # Verschlüsseln
            writer.encrypt(password)
            
            # Speichern
            with open(output_path, 'wb') as f:
                writer.write(f)
            
            return True
        except Exception as e:
            raise Exception(f"Verschlüsselung fehlgeschlagen: {str(e)}")
    
    @staticmethod
    def decrypt_pdf(input_path, output_path, password):
        """Entschlüsselt ein PDF"""
        if not HAS_PDF:
            raise Exception("PyPDF2 nicht installiert")
        
        try:
            reader = PdfReader(input_path)
            
            if reader.is_encrypted:
                # Passwort versuchen
                if not reader.decrypt(password):
                    raise Exception("Falsches Passwort")
            
            writer = PdfWriter()
            
            # Alle Seiten kopieren
            for page in reader.pages:
                writer.add_page(page)
            
            # OHNE Verschlüsselung speichern
            with open(output_path, 'wb') as f:
                writer.write(f)
            
            return True
        except Exception as e:
            raise Exception(f"Entschlüsselung fehlgeschlagen: {str(e)}")
    
    @staticmethod
    def extract_pages(input_path, output_path, page_indices):
        """Erstellt PDF-Auszug mit ausgewählten Seiten"""
        if not HAS_PDF:
            raise Exception("PyPDF2 nicht installiert")
        
        try:
            reader = PdfReader(input_path)
            writer = PdfWriter()
            
            # Nur ausgewählte Seiten
            for idx in page_indices:
                if 0 <= idx < len(reader.pages):
                    writer.add_page(reader.pages[idx])
            
            with open(output_path, 'wb') as f:
                writer.write(f)
            
            return True
        except Exception as e:
            raise Exception(f"Auszug-Erstellung fehlgeschlagen: {str(e)}")
    
    @staticmethod
    def remove_text_from_pdf(input_path, output_path):
        """Entfernt Text aus PDF, behlt nur Bilder"""
        if not HAS_PDF:
            raise Exception("PyPDF2 nicht installiert")
        
        try:
            # Dies ist komplex - vereinfachte Version:
            # Konvertiere zu Bildern und zurück zu PDF
            if not HAS_PDF2IMAGE:
                raise Exception("pdf2image nicht installiert")
            
            images = convert_from_path(input_path)
            
            if images:
                images[0].save(output_path, "PDF", save_all=True, 
                             append_images=images[1:] if len(images) > 1 else [])
                return True
            
            return False
        except Exception as e:
            raise Exception(f"Text-Entfernung fehlgeschlagen: {str(e)}")
    
    @staticmethod
    def apply_ocr_to_pdf(input_path, output_path, lang='deu'):
        """Wendet OCR auf PDF an"""
        if not HAS_OCR or not HAS_PDF2IMAGE:
            raise Exception("pytesseract oder pdf2image nicht installiert")
        
        try:
            # PDF zu Bildern
            images = convert_from_path(input_path)
            
            # OCR auf jedes Bild
            writer = PdfWriter()
            
            for img in images:
                # OCR durchführen
                text = pytesseract.image_to_string(img, lang=lang)
                
                # Hier müsste man eigentlich ein searchable PDF erstellen
                # Vereinfachte Version: Nur Text extrahieren
                # Für production: pdf2pdfocr oder ocrmypdf verwenden
                pass
            
            # Für jetzt: Kopiere Original
            shutil.copy(input_path, output_path)
            return True
            
        except Exception as e:
            raise Exception(f"OCR fehlgeschlagen: {str(e)}")

# ============================================================================
# ANONYMIZATION ENGINE
# ============================================================================

class AnonymizationWorker(QThread):
    """Worker-Thread für Anonymisierung und Schwärzung"""
    progress = pyqtSignal(int, int)  # current, total
    log_message = pyqtSignal(str)
    finished = pyqtSignal()
    
    def __init__(self, file_paths, blacklist, whitelist, placeholder="[-----]", mode="anonymize"):
        super().__init__()
        self.file_paths = file_paths
        self.blacklist = blacklist
        self.whitelist = whitelist
        self.placeholder = placeholder
        self.mode = mode  # "anonymize" or "redact"
        self.is_running = True
    
    def run(self):
        """Verarbeitet Dateien"""
        total = len(self.file_paths)
        
        for idx, file_path in enumerate(self.file_paths):
            if not self.is_running:
                break
            
            self.progress.emit(idx + 1, total)
            
            try:
                ext = os.path.splitext(file_path)[1].lower()
                folder = os.path.dirname(file_path)
                basename = os.path.splitext(os.path.basename(file_path))[0]
                
                if self.mode == "anonymize":
                    # Textdatei-Anonymisierung
                    if ext in ['.txt', '.log', '.py', '.md']:
                        output_path = os.path.join(folder, f"{basename}_anonymisiert{ext}")
                        self.anonymize_text_file(file_path, output_path)
                        self.log_message.emit(f"✅ Anonymisiert: {os.path.basename(output_path)}")
                    
                    elif ext == '.docx':
                        output_path = os.path.join(folder, f"{basename}_anonymisiert{ext}")
                        self.anonymize_docx_file(file_path, output_path)
                        self.log_message.emit(f"✅ Anonymisiert: {os.path.basename(output_path)}")
                    
                    elif ext == '.pdf':
                        output_path = os.path.join(folder, f"{basename}_geschwrzt.pdf")
                        self.redact_pdf(file_path, output_path)
                        self.log_message.emit(f"✅ Geschwärzt: {os.path.basename(output_path)}")
                    
                    else:
                        self.log_message.emit(f"⚠️ Format nicht unterstützt: {os.path.basename(file_path)}")
                
                elif self.mode == "redact":
                    # PDF-Schwärzung
                    if ext != '.pdf':
                        # Konvertiere zu PDF zuerst
                        temp_pdf = self.convert_to_pdf(file_path)
                        if temp_pdf:
                            output_path = os.path.join(folder, f"{basename}_geschwrzt.pdf")
                            self.redact_pdf(temp_pdf, output_path)
                            os.remove(temp_pdf)
                            self.log_message.emit(f"✅ Geschwärzt: {os.path.basename(output_path)}")
                        else:
                            self.log_message.emit(f"❌ Konvertierung fehlgeschlagen: {os.path.basename(file_path)}")
                    else:
                        output_path = os.path.join(folder, f"{basename}_geschwrzt.pdf")
                        self.redact_pdf(file_path, output_path)
                        self.log_message.emit(f" Geschwrzt: {os.path.basename(output_path)}")
            
            except Exception as e:
                self.log_message.emit(f"❌ Fehler bei {os.path.basename(file_path)}: {str(e)}")
        
        self.finished.emit()
    
    def anonymize_text_file(self, input_path, output_path):
        """Anonymisiert Textdatei durch Platzhalter-Ersetzung"""
        with open(input_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        for word in self.blacklist:
            if self.is_whitelisted(word):
                continue
            # Case-insensitive replace
            import re
            pattern = re.compile(re.escape(word), re.IGNORECASE)
            content = pattern.sub(self.placeholder, content)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(content)
    
    def anonymize_docx_file(self, input_path, output_path):
        """Anonymisiert Word-Dokument"""
        if not HAS_DOCX:
            raise Exception("python-docx nicht installiert")
        
        doc = docx.Document(input_path)
        
        # Absätze
        for paragraph in doc.paragraphs:
            for word in self.blacklist:
                if self.is_whitelisted(word):
                    continue
                if word.lower() in paragraph.text.lower():
                    # Vereinfachte Ersetzung
                    import re
                    pattern = re.compile(re.escape(word), re.IGNORECASE)
                    paragraph.text = pattern.sub(self.placeholder, paragraph.text)
        
        # Tabellen
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for word in self.blacklist:
                            if self.is_whitelisted(word):
                                continue
                            if word.lower() in paragraph.text.lower():
                                import re
                                pattern = re.compile(re.escape(word), re.IGNORECASE)
                                paragraph.text = pattern.sub(self.placeholder, paragraph.text)
        
        doc.save(output_path)
    
    def redact_pdf(self, input_path, output_path):
        """Schwärzt PDF mit schwarzen Balken"""
        if not HAS_FITZ:
            raise Exception("PyMuPDF (fitz) nicht installiert")
        
        doc = fitz.open(input_path)
        
        for page in doc:
            for word in self.blacklist:
                if self.is_whitelisted(word):
                    continue
                
                # Suche Wort im PDF
                hits = page.search_for(word)
                
                for rect in hits:
                    # Fge Schwrzungs-Annotation hinzu
                    page.add_redact_annot(rect, fill=(0, 0, 0))
            
            # Wende Schwrzungen an
            page.apply_redactions()
        
        doc.save(output_path)
        doc.close()
    
    def convert_to_pdf(self, path):
        """Konvertiert Datei zu PDF"""
        import tempfile
        ext = os.path.splitext(path)[1].lower()
        temp_pdf = os.path.join(tempfile.gettempdir(), f"temp_{int(time.time())}.pdf")
        
        # Bild -> PDF
        if ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff']:
            try:
                from PIL import Image
                img = Image.open(path)
                img.convert("RGB").save(temp_pdf)
                return temp_pdf
            except (OSError, IOError):
                return None
        
        # TXT -> PDF
        if ext in ['.txt', '.log', '.py', '.md']:
            try:
                doc = fitz.open()
                page = doc.new_page()
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
                page.insert_text((50, 50), text, fontsize=10)
                doc.save(temp_pdf)
                doc.close()
                return temp_pdf
            except (OSError, IOError):
                return None
        
        return None
    
    def is_whitelisted(self, word):
        """Prüft ob Wort auf Whitelist steht"""
        norm_word = word.strip().lower()
        for white in self.whitelist:
            if norm_word == white.strip().lower():
                return True
        return False
    
    def stop(self):
        """Stoppt Worker"""
        self.is_running = False

# ============================================================================
# 3. ENHANCED DATABASE WITH PDF METADATA
# ============================================================================

DDL_BASE = """
PRAGMA journal_mode=WAL;
CREATE TABLE IF NOT EXISTS files(
    id INTEGER PRIMARY KEY, 
    content_hash TEXT UNIQUE, 
    size INTEGER, 
    mime TEXT, 
    first_seen TEXT
);

CREATE TABLE IF NOT EXISTS versions(
    id INTEGER PRIMARY KEY, 
    file_id INTEGER, 
    name TEXT, 
    path TEXT, 
    mtime TEXT, 
    ctime TEXT, 
    version_index INTEGER, 
    source_side TEXT
);

CREATE TABLE IF NOT EXISTS collections(
    id INTEGER PRIMARY KEY, 
    name TEXT UNIQUE
);

CREATE TABLE IF NOT EXISTS collection_items(
    collection_id INTEGER, 
    version_id INTEGER, 
    PRIMARY KEY(collection_id, version_id)
);

CREATE TABLE IF NOT EXISTS tags(
    id INTEGER PRIMARY KEY, 
    file_id INTEGER, 
    tag TEXT
);

CREATE INDEX IF NOT EXISTS idx_versions_path ON versions(path);
CREATE INDEX IF NOT EXISTS idx_versions_mtime ON versions(mtime);
"""



# ============================================================================
# AUTO-UPDATE SYSTEM mit Watchdog
# ============================================================================

if HAS_WATCHDOG:
    class ConnectionWatcher(FileSystemEventHandler):
        """überwacht Verbindungsordner auf änderungen"""
        
        def __init__(self, connection_config, callback):
            super().__init__()
            self.conn_config = connection_config
            self.callback = callback
            self.pending_changes = set()
            self.last_trigger = time.time()
            self.cooldown_seconds = 5
        
        def on_any_event(self, event):
            """Sammelt änderungen mit Cooldown"""
            if event.is_directory:
                return
            
            # Ignoriere temporre Dateien
            if any(event.src_path.endswith(ext) for ext in ['.tmp', '.temp', '~']):
                return
            
            self.pending_changes.add(event.src_path)
            
            # Check if cooldown expired
            now = time.time()
            if now - self.last_trigger >= self.cooldown_seconds:
                self._trigger_update()
        
        def _trigger_update(self):
            """Triggered Update nach Cooldown"""
            if self.pending_changes:
                print(f"Auto-Update: {len(self.pending_changes)} änderungen in '{self.conn_config['name']}'")
                self.callback(self.conn_config, list(self.pending_changes))
                self.pending_changes.clear()
                self.last_trigger = time.time()
    
    
    class ConnectionAutoUpdater:
        """Verwaltet Watchdog Observer für alle enabled Verbindungen"""
        
        def __init__(self, conn_manager):
            self.conn_manager = conn_manager
            self.observers = {}  # conn_id -> Observer
        
        def start_watching(self, conn_id):
            """Startet Watchdog für eine Verbindung"""
            conns = self.conn_manager.list_connections()
            conn = next((c for c in conns if c.get("id") == conn_id), None)
            
            if not conn or not conn.get("enabled", True) or not conn.get("auto_update", False):
                return False
            
            # Stop existing observer
            self.stop_watching(conn_id)
            
            # Create new observer
            observer = Observer()
            handler = ConnectionWatcher(conn, self._on_change_detected)
            
            # Watch all source folders
            for source in conn.get("sources", []):
                path = source.get("path")
                if path and os.path.exists(path):
                    try:
                        observer.schedule(handler, path, recursive=True)
                        print(f"👁️ Watching: {path}")
                    except Exception as e:
                        print(f"⚠️ Cannot watch {path}: {e}")
            
            observer.start()
            self.observers[conn_id] = observer
            return True
        
        def stop_watching(self, conn_id):
            """Stoppt Watchdog für eine Verbindung"""
            if conn_id in self.observers:
                self.observers[conn_id].stop()
                self.observers[conn_id].join(timeout=2)
                del self.observers[conn_id]
        
        def stop_all(self):
            """Stoppt alle Observer"""
            for conn_id in list(self.observers.keys()):
                self.stop_watching(conn_id)
        
        def _on_change_detected(self, conn_config, changed_files):
            """Callback wenn änderungen erkannt wurden"""
            # Hier knnte man eine partielle Re-Indizierung triggern
            # Vorerst nur Logging
            print(f"{len(changed_files)} Datei(en) geändert:")
            for f in changed_files[:5]:  # Max 5 anzeigen
                print(f"  - {os.path.basename(f)}")
            if len(changed_files) > 5:
                print(f"  ... und {len(changed_files) - 5} weitere")


class ConnectionDB:
    def __init__(self, db_path):
        self.db_path = db_path
        os.makedirs(os.path.dirname(db_path) or ".", exist_ok=True)
        self.conn = sqlite3.connect(db_path, check_same_thread=False)
        self.conn.executescript(DDL_BASE)
        self._migrate_v9()

    def _migrate_v9(self):
        """Migration für V9 Features (PDF-Metadaten)"""
        cur = self.conn.cursor()
        
        # Files Tabelle
        cur.execute("PRAGMA table_info(files)")
        files_columns = [row[1] for row in cur.fetchall()]
        
        if 'pdf_encrypted' not in files_columns:
            try: 
                self.conn.execute("ALTER TABLE files ADD COLUMN pdf_encrypted INTEGER DEFAULT 0")
                print("✓ Added pdf_encrypted column")
            except Exception as e: 
                print(f"Could not add pdf_encrypted: {e}")
        
        if 'pdf_has_text' not in files_columns:
            try: 
                self.conn.execute("ALTER TABLE files ADD COLUMN pdf_has_text INTEGER DEFAULT 0")
                print("✓ Added pdf_has_text column")
            except Exception as e:
                print(f"Could not add pdf_has_text: {e}")
        
        if 'pdf_was_encrypted' not in files_columns:
            try: 
                self.conn.execute("ALTER TABLE files ADD COLUMN pdf_was_encrypted INTEGER DEFAULT 0")
                print("✓ Added pdf_was_encrypted column")
            except Exception as e:
                print(f"Could not add pdf_was_encrypted: {e}")
        
        # Versions Tabelle
        cur.execute("PRAGMA table_info(versions)")
        versions_columns = [row[1] for row in cur.fetchall()]
        
        if 'is_favorite' not in versions_columns:
            try: 
                self.conn.execute("ALTER TABLE versions ADD COLUMN is_favorite INTEGER DEFAULT 0")
                print("✓ Added is_favorite column")
            except Exception as e:
                print(f"Could not add is_favorite: {e}")
        
        if 'version_label' not in versions_columns:
            try: 
                self.conn.execute("ALTER TABLE versions ADD COLUMN version_label TEXT")
                print("✓ Added version_label column")
            except Exception as e:
                print(f"Could not add version_label: {e}")
        
        if 'is_deleted' not in versions_columns:
            try: 
                self.conn.execute("ALTER TABLE versions ADD COLUMN is_deleted INTEGER DEFAULT 0")
                self.conn.execute("ALTER TABLE versions ADD COLUMN deleted_at TEXT")
                print("✓ Added is_deleted and deleted_at columns")
            except Exception as e:
                print(f"Could not add is_deleted: {e}")
        
        # Safety-Mode: is_hidden (NEU V13.2!)
        if 'is_hidden' not in versions_columns:
            try: 
                self.conn.execute("ALTER TABLE versions ADD COLUMN is_hidden INTEGER DEFAULT 0")
                self.conn.execute("ALTER TABLE versions ADD COLUMN hidden_at TEXT")
                print(" Added is_hidden and hidden_at columns (Safety-Mode)")
            except Exception as e:
                print(f"Could not add is_hidden: {e}")
        
        # Collections
        cur.execute("PRAGMA table_info(collections)")
        coll_columns = [row[1] for row in cur.fetchall()]
        
        if 'description' not in coll_columns:
            try: 
                self.conn.execute("ALTER TABLE collections ADD COLUMN description TEXT")
                print("✓ Added description column")
            except Exception as e:
                print(f"Could not add description: {e}")
        
        if 'created_at' not in coll_columns:
            try: 
                self.conn.execute("ALTER TABLE collections ADD COLUMN created_at TEXT")
                print("✓ Added created_at column")
            except Exception as e:
                print(f"Could not add created_at: {e}")
        
        # Collection Items
        cur.execute("PRAGMA table_info(collection_items)")
        ci_columns = [row[1] for row in cur.fetchall()]
        
        if 'added_at' not in ci_columns:
            try: 
                self.conn.execute("ALTER TABLE collection_items ADD COLUMN added_at TEXT")
                print("✓ Added added_at column")
            except Exception as e:
                print(f"Could not add added_at: {e}")
        
        # Indices erstellen
        try:
            self.conn.execute("CREATE INDEX IF NOT EXISTS idx_versions_deleted ON versions(is_deleted)")
            self.conn.execute("CREATE INDEX IF NOT EXISTS idx_versions_favorite ON versions(is_favorite)")
            self.conn.execute("CREATE INDEX IF NOT EXISTS idx_versions_hidden ON versions(is_hidden)")
            self.conn.execute("CREATE INDEX IF NOT EXISTS idx_files_encrypted ON files(pdf_encrypted)")
            print("✓ Created indices")
        except Exception as e:
            print(f"Could not create indices: {e}")
        
        self.conn.commit()
        print("✓ Migration V9 completed")

    def close(self):
        try: self.conn.close()
        except (sqlite3.Error, OSError): pass

    def get_file_id_by_hash(self, content_hash):
        cur = self.conn.cursor()
        cur.execute("SELECT id FROM files WHERE content_hash=?", (content_hash,))
        row = cur.fetchone()
        return row[0] if row else None

    def upsert_file(self, content_hash, size, mime=None, pdf_encrypted=False, pdf_has_text=False, pdf_was_encrypted=False):
        fid = self.get_file_id_by_hash(content_hash)
        if fid: 
            # Update PDF metadata if columns exist
            try:
                self.conn.execute("""UPDATE files SET pdf_encrypted=?, pdf_has_text=?, pdf_was_encrypted=? 
                                    WHERE id=?""", 
                                (pdf_encrypted, pdf_has_text, pdf_was_encrypted, fid))
                self.conn.commit()
            except sqlite3.OperationalError:
                # Columns don't exist yet - skip update
                pass
            return fid
        
        ts = datetime.utcnow().isoformat()
        try:
            cur = self.conn.cursor()
            # Try with PDF columns
            try:
                cur.execute("""INSERT INTO files(content_hash,size,mime,first_seen,pdf_encrypted,pdf_has_text,pdf_was_encrypted) 
                              VALUES (?,?,?,?,?,?,?)""", 
                           (content_hash, size, mime, ts, pdf_encrypted, pdf_has_text, pdf_was_encrypted))
            except sqlite3.OperationalError:
                # Fallback: Without PDF columns (old schema)
                cur.execute("""INSERT INTO files(content_hash,size,mime,first_seen) 
                              VALUES (?,?,?,?)""", 
                           (content_hash, size, mime, ts))
            
            self.conn.commit()
            return cur.lastrowid
        except sqlite3.IntegrityError:
            return self.get_file_id_by_hash(content_hash)

    def get_file_pdf_status(self, file_id):
        """Gibt PDF-Status zurück (encrypted, has_text, was_encrypted)"""
        try:
            cur = self.conn.cursor()
            cur.execute("SELECT pdf_encrypted, pdf_has_text, pdf_was_encrypted FROM files WHERE id=?", (file_id,))
            row = cur.fetchone()
            if row:
                return {'encrypted': bool(row[0]), 'has_text': bool(row[1]), 'was_encrypted': bool(row[2])}
        except sqlite3.OperationalError:
            # Columns don't exist yet
            pass
        
        return {'encrypted': False, 'has_text': False, 'was_encrypted': False}

    def upsert_version(self, file_id, name, path, mtime, ctime, idx, side):
        """Update existing or insert new version"""
        cur = self.conn.cursor()
        cur.execute("SELECT id FROM versions WHERE path=?", (path,))
        row = cur.fetchone()
        
        if row:
            vid = row[0]
            cur.execute("""UPDATE versions SET file_id=?, mtime=?, ctime=?, name=?, 
                          source_side=?, is_deleted=0, deleted_at=NULL WHERE id=?""",
                       (file_id, mtime, ctime, name, side, vid))
        else:
            cur.execute("""INSERT INTO versions(file_id,name,path,mtime,ctime,version_index,source_side) 
                          VALUES (?,?,?,?,?,?,?)""",
                       (file_id, name, path, mtime, ctime, idx, side))
            vid = cur.lastrowid
        
        self.conn.commit()
        return vid

    # SOFT/HARD DELETE
    def soft_delete_version(self, version_id):
        ts = datetime.utcnow().isoformat()
        self.conn.execute("UPDATE versions SET is_deleted=1, deleted_at=? WHERE id=?", (ts, version_id))
        self.conn.commit()
    
    def hard_delete_version(self, version_id):
        cur = self.conn.cursor()
        cur.execute("SELECT path FROM versions WHERE id=?", (version_id,))
        row = cur.fetchone()
        
        if row:
            path = row[0]
            try:
                if os.path.exists(path):
                    os.remove(path)
            except Exception as e:
                print(f"Fehler beim physischen Löschen von {path}: {e}")
        
        self.conn.execute("DELETE FROM collection_items WHERE version_id=?", (version_id,))
        self.conn.execute("DELETE FROM versions WHERE id=?", (version_id,))
        self.conn.commit()
    
    def restore_version(self, version_id):
        self.conn.execute("UPDATE versions SET is_deleted=0, deleted_at=NULL WHERE id=?", (version_id,))
        self.conn.commit()
    
    def safety_hide_version(self, version_id):
        """Versteckt Version (Safety-Mode - KEINE Dateisystem-nderung!) - NEU V13.2!"""
        ts = datetime.utcnow().isoformat()
        self.conn.execute("UPDATE versions SET is_hidden=1, hidden_at=? WHERE id=?", (ts, version_id))
        self.conn.commit()
    
    def safety_unhide_version(self, version_id):
        """Zeigt versteckte Version wieder an (Safety-Mode) - NEU V13.2!"""
        self.conn.execute("UPDATE versions SET is_hidden=0, hidden_at=NULL WHERE id=?", (version_id,))
        self.conn.commit()
    
    def cleanup_old_deleted(self, days):
        if days <= 0:
            return 0
        
        cutoff = (datetime.utcnow() - timedelta(days=days)).isoformat()
        cur = self.conn.cursor()
        cur.execute("SELECT id FROM versions WHERE is_deleted=1 AND deleted_at < ?", (cutoff,))
        old_ids = [row[0] for row in cur.fetchall()]
        
        for vid in old_ids:
            self.hard_delete_version(vid)
        
        return len(old_ids)

    def prune_deleted_files(self, active_version_ids, source_root):
        if not active_version_ids: return 0
        
        placeholders = ','.join('?' for _ in active_version_ids)
        query = f"""
            SELECT id FROM versions 
            WHERE path LIKE ? || '%' 
            AND id NOT IN ({placeholders})
            AND is_deleted=0
        """
        
        try:
            cur = self.conn.cursor()
            cur.execute(query, [source_root] + list(active_version_ids))
            to_delete = [row[0] for row in cur.fetchall()]
            
            ts = datetime.utcnow().isoformat()
            for vid in to_delete:
                self.conn.execute("UPDATE versions SET is_deleted=1, deleted_at=? WHERE id=?", (ts, vid))
            
            self.conn.commit()
            return len(to_delete)
        except Exception as e:
            print(f"Prune error: {e}")
            return 0

    def get_latest_version_by_path(self, path):
        cur = self.conn.cursor()
        cur.execute("SELECT mtime, file_id, id FROM versions WHERE path=? LIMIT 1", (path,))
        return cur.fetchone()

    # VERSION SWAPPING
    def swap_version_index(self, version_id, direction):
        cur = self.conn.cursor()
        cur.execute("SELECT version_index, file_id FROM versions WHERE id=?", (version_id,))
        row = cur.fetchone()
        if not row: return
        
        current_idx, file_id = row
        if current_idx is None: current_idx = 1
        
        target_idx = current_idx + direction
        if target_idx < 1: return
        
        cur.execute("SELECT id FROM versions WHERE file_id=? AND version_index=?", (file_id, target_idx))
        collision = cur.fetchone()
        
        if collision:
            other_id = collision[0]
            self.conn.execute("UPDATE versions SET version_index=? WHERE id=?", (current_idx, other_id))
            self.conn.execute("UPDATE versions SET version_index=? WHERE id=?", (target_idx, version_id))
        else:
            self.conn.execute("UPDATE versions SET version_index=? WHERE id=?", (target_idx, version_id))
        
        self.conn.commit()

    # COLLECTIONS
    def get_collections(self):
        cur = self.conn.cursor()
        cur.execute("SELECT id, name, description FROM collections ORDER BY name")
        return cur.fetchall()
    
    def add_collection(self, name, description=""):
        ts = datetime.utcnow().isoformat()
        try:
            self.conn.execute("INSERT INTO collections(name, description, created_at) VALUES (?,?,?)", 
                            (name, description, ts))
            self.conn.commit()
            return True
        except sqlite3.IntegrityError:
            return False

    def remove_collection(self, collection_id):
        self.conn.execute("DELETE FROM collection_items WHERE collection_id=?", (collection_id,))
        self.conn.execute("DELETE FROM collections WHERE id=?", (collection_id,))
        self.conn.commit()

    def add_to_collection(self, col_id, version_id):
        ts = datetime.utcnow().isoformat()
        try:
            self.conn.execute("INSERT OR IGNORE INTO collection_items(collection_id, version_id, added_at) VALUES (?,?,?)",
                            (col_id, version_id, ts))
            self.conn.commit()
            return True
        except sqlite3.Error:
            return False

    def remove_from_collection(self, col_id, version_id):
        self.conn.execute("DELETE FROM collection_items WHERE collection_id=? AND version_id=?", (col_id, version_id))
        self.conn.commit()

    def set_favorite(self, version_id, is_fav):
        self.conn.execute("UPDATE versions SET is_favorite=? WHERE id=?", (1 if is_fav else 0, version_id))
        self.conn.commit()
    
    def set_version_label(self, version_id, label):
        self.conn.execute("UPDATE versions SET version_label=? WHERE id=?", (label, version_id))
        self.conn.commit()

# ============================================================================
# 4. SYNC WORKER (mit PDF-Metadaten)
# ============================================================================

class SyncSignals(QObject):
    progress = pyqtSignal(int)
    status = pyqtSignal(str)
    finished = pyqtSignal()
    error = pyqtSignal(str)

class SyncWorker(QThread):
    def __init__(self, cfg, mode="sync"):
        super().__init__()
        self.cfg = cfg
        self.mode = mode
        self.signals = SyncSignals()
        self.is_paused = False
        self.is_killed = False
        self.max_len = cfg.get("max_filename_length") or 120

    def pause(self): self.is_paused = True
    def resume(self): self.is_paused = False
    def kill(self): self.is_killed = True

    def run(self):
        db = None
        try:
            db = ConnectionDB(self.cfg["db_path"])
            if self.mode == "sync": 
                self._run_fast_scan(db)
        except Exception as e:
            import traceback
            traceback.print_exc()
            self.signals.error.emit(str(e))
        finally:
            if db: db.close()
            self.signals.finished.emit()

    def _run_fast_scan(self, db):
        self.signals.status.emit("Sync gestartet...")
        mode = self.cfg["direction"]
        found_ids_src = set()
        found_ids_tgt = set()

        all_files = []
        sources = [(self.cfg["source"], "source")]
        if mode == "two_way": 
            sources.append((self.cfg["target"], "target"))
        
        for base, side in sources:
            if not base or not os.path.exists(base): continue
            for root, _, fs in os.walk(base):
                for f in fs:
                    if self.is_killed: return
                    all_files.append((side, os.path.join(root, f)))
        
        total = len(all_files)
        done = 0
        last_emit = 0
        
        for side, path in all_files:
            while self.is_paused: time.sleep(0.5)
            if self.is_killed: return
            
            done += 1
            if done - last_emit > 10 or done == total:
                self.signals.progress.emit(int(done * 100 / max(1, total)))
                last_emit = done

            try:
                try: 
                    stat = os.stat(path)
                except FileNotFoundError: 
                    continue 

                mtime_iso = datetime.utcfromtimestamp(stat.st_mtime).isoformat()
                ctime = datetime.utcfromtimestamp(stat.st_ctime).isoformat()
                size = stat.st_size
                name = os.path.basename(path)
                is_cloud = is_cloud_placeholder(path)
                
                # PDF-Metadaten extrahieren
                pdf_encrypted = False
                pdf_has_text = False
                pdf_was_encrypted = False
                
                if path.lower().endswith('.pdf'):
                    pdf_encrypted = is_pdf_encrypted(path)
                    if not pdf_encrypted:
                        pdf_has_text = has_pdf_text(path)
                        # TODO: Detect if was previously encrypted
                
                latest = db.get_latest_version_by_path(path)
                file_id, version_id = None, None
                
                if latest and latest[0] == mtime_iso:
                    file_id, version_id = latest[1], latest[2]
                else:
                    content_hash = f"CLOUD:{size}:{mtime_iso}" if is_cloud else sha256_file(path)
                    if content_hash: 
                        file_id = db.upsert_file(content_hash, size, mime=None,
                                                pdf_encrypted=pdf_encrypted,
                                                pdf_has_text=pdf_has_text,
                                                pdf_was_encrypted=pdf_was_encrypted)
                        version_id = db.upsert_version(file_id, name, path, mtime_iso, ctime, 1, side)

                if version_id:
                    if side == "source": found_ids_src.add(version_id)
                    else: found_ids_tgt.add(version_id)

                if mode != "index_only" and file_id:
                    self._handle_copy(db, side, path, name, file_id, mtime_iso, ctime, found_ids_tgt)

            except Exception as e: 
                print(f"Error {path}: {e}")

        self.signals.status.emit("Bereinige Datenbank...")
        if self.cfg["source"]: 
            deleted = db.prune_deleted_files(found_ids_src, self.cfg["source"])
            if deleted > 0:
                self.signals.status.emit(f"{deleted} Dateien als gelöscht markiert")
        
        if mode == "two_way" and self.cfg["target"]: 
            db.prune_deleted_files(found_ids_tgt, self.cfg["target"])

        self.signals.status.emit("Fertig.")
        self.signals.progress.emit(100)

    def _handle_copy(self, db, side, path, name, file_id, mtime, ctime, tracker):
        tgt_root = self.cfg["target"]
        if not tgt_root or side != "source": return
        
        tgt_path = self._calc_tgt(path, self.cfg["source"], tgt_root)
        if self._should_copy(path, tgt_path):
            try:
                os.makedirs(os.path.dirname(tgt_path), exist_ok=True)
                shutil.copy2(path, tgt_path)
                vid = db.upsert_version(file_id, name, tgt_path, mtime, ctime, 1, "target_copy")
                if vid: tracker.add(vid)
            except Exception as e:
                print(f"Copy error: {e}")
        
    def _calc_tgt(self, src_path, src_root, tgt_root):
        rel = os.path.relpath(src_path, src_root)
        if self.cfg.get("structure_handling", {}).get("flatten_to_single_folder"):
            return os.path.join(tgt_root, shorten_filename(os.path.basename(rel), self.max_len))
        return os.path.join(tgt_root, rel)

    def _should_copy(self, src, tgt):
        if not os.path.exists(tgt): return True
        try:
            return (os.stat(src).st_size != os.stat(tgt).st_size) or \
                   (os.stat(src).st_mtime > os.stat(tgt).st_mtime + 2)
        except OSError:
            return False

# ============================================================================
# 5. SEARCH WORKER
# ============================================================================

class SearchWorker(QThread):
    results_found = pyqtSignal(list)
    finished = pyqtSignal()
    
    def __init__(self, manager, params, settings):
        super().__init__()
        self.manager = manager
        self.params = params
        self.settings = settings

    def run(self):
        results = []
        LIMIT = 500
        term = self.params.get("term", "").lower()
        wildcard = f"%{term}%"
        file_types = self.params.get("types", [])
        only_fav = self.params.get("favorites", False)
        collection_id = self.params.get("collection_id", None)
        show_deleted = self.params.get("show_deleted", False)
        show_hidden = self.params.get("show_hidden", False)
        
        for db_path in self.manager.dbs:
            if not os.path.exists(db_path): continue
            conn = None
            try:
                conn = sqlite3.connect(db_path)
                conn.row_factory = sqlite3.Row
                
                sql = """
                    SELECT v.id, v.name, v.path, COALESCE(v.display_name, v.name) as display_name, v.mtime, 
                           COALESCE(v.is_favorite, 0) as is_favorite, 
                           v.version_index,
                           v.version_label, 
                           COALESCE(v.is_deleted, 0) as is_deleted,
                           COALESCE(v.is_hidden, 0) as is_hidden,
                           v.hidden_at, 
                           v.deleted_at,
                           f.content_hash,
                           f.size
                    FROM versions v
                    JOIN files f ON v.file_id = f.id
                    LEFT JOIN collection_items ci ON ci.version_id = v.id
                    WHERE (lower(v.name) LIKE ? OR lower(v.path) LIKE ?)
                """
                args = [wildcard, wildcard]
                
                if not show_deleted:
                    sql += " AND v.is_deleted = 0"
                
                if not show_hidden:
                    sql += " AND v.is_deleted = 0"
                
                if only_fav: 
                    sql += " AND v.is_favorite = 1"
                
                if collection_id: 
                    sql += " AND ci.collection_id = ?"
                    args.append(collection_id)
                
                sql += f" ORDER BY v.is_favorite DESC, v.name LIMIT {LIMIT}"
                
                rows = conn.execute(sql, args).fetchall()
                for row in rows:
                    cat = get_file_category(row["name"])
                    if file_types and cat not in file_types: continue
                    
                    # Get PDF status separately (safe for old schemas)
                    pdf_status = {'encrypted': False, 'has_text': False, 'was_encrypted': False}
                    try:
                        pdf_row = conn.execute(
                            "SELECT pdf_encrypted, pdf_has_text, pdf_was_encrypted FROM files WHERE id=(SELECT file_id FROM versions WHERE id=?)",
                            (row["id"],)
                        ).fetchone()
                        if pdf_row:
                            pdf_status = {
                                'encrypted': bool(pdf_row[0]) if pdf_row[0] is not None else False,
                                'has_text': bool(pdf_row[1]) if pdf_row[1] is not None else False,
                                'was_encrypted': bool(pdf_row[2]) if pdf_row[2] is not None else False
                            }
                    except sqlite3.OperationalError:
                        # Columns don't exist - keep defaults
                        pass
                    
                    results.append({
                        "id": row["id"],
                        "name": row["name"], 
                        "path": row["path"], 
                        "mtime": row["mtime"],
                        "is_favorite": bool(row["is_favorite"]),
                        "version_index": row["version_index"] if row["version_index"] else 1,
                        "version_label": row["version_label"] or "",
                        "category": cat,
                        "is_deleted": bool(row["is_deleted"]),
                        "deleted_at": row["deleted_at"] or "",
                        "pdf_encrypted": pdf_status['encrypted'],
                        "pdf_has_text": pdf_status['has_text'],
                        "pdf_was_encrypted": pdf_status['was_encrypted'],
                        "content_hash": row["content_hash"],
                        "size": row["size"],
                        "db": db_path
                    })
            except Exception as e:
                print(f"Search error in {db_path}: {e}")
            finally:
                if conn: conn.close()
            
            if len(results) >= LIMIT: break
            
        self.results_found.emit(results)
        self.finished.emit()


class DuplicateWorker(QThread):
    """Findet Duplikate basierend auf Content-Hash"""
    results_ready = pyqtSignal(dict)
    progress = pyqtSignal(int, str)
    finished = pyqtSignal()
    
    def __init__(self, db_paths, criteria):
        super().__init__()
        self.db_paths = db_paths
        self.criteria = criteria

    def run(self):
        self.progress.emit(0, "Suche nach Duplikaten...")
        
        all_files = {}
        
        total_dbs = len(self.db_paths)
        for idx, db_path in enumerate(self.db_paths):
            if not os.path.exists(db_path): 
                continue
            
            self.progress.emit(int((idx / total_dbs) * 50), f"Analysiere Datenbank {idx+1}/{total_dbs}...")
            
            try:
                conn = sqlite3.connect(db_path)
                conn.row_factory = sqlite3.Row
                
                sql = """
                    SELECT v.id, v.name, v.path, COALESCE(v.display_name, v.name) as display_name, v.mtime, v.ctime, f.content_hash, f.size
                    FROM versions v 
                    JOIN files f ON v.file_id = f.id
                    WHERE v.is_deleted = 0
                    AND f.content_hash NOT LIKE 'CLOUD:%'
                """
                
                rows = conn.execute(sql).fetchall()
                
                for row in rows:
                    h = row['content_hash']
                    if not h:
                        continue
                    
                    file_info = {
                        'id': row['id'],
                        'name': row['name'],
                        'path': row['path'],
                        'mtime': row['mtime'],
                        'ctime': row['ctime'],
                        'size': row['size'],
                        'hash': h,
                        'db': db_path
                    }
                    
                    if h not in all_files:
                        all_files[h] = []
                    
                    all_files[h].append(file_info)
                
                conn.close()
                
            except Exception as e:
                print(f"Error scanning {db_path}: {e}")
        
        self.progress.emit(60, "Analysiere Duplikate...")
        
        duplicates = {}
        
        for h, files in all_files.items():
            if len(files) < 2:
                continue
            
            if self.criteria == "hash":
                duplicates[h] = files
            
            elif self.criteria == "hash_name":
                name_groups = {}
                for f in files:
                    name = f['name']
                    if name not in name_groups:
                        name_groups[name] = []
                    name_groups[name].append(f)
                
                for name, group in name_groups.items():
                    if len(group) > 1:
                        key = f"{h}_{name}"
                        duplicates[key] = group
            
            elif self.criteria == "hash_similar_name":
                base_groups = {}
                for f in files:
                    base = os.path.splitext(f['name'])[0].lower()
                    if base not in base_groups:
                        base_groups[base] = []
                    base_groups[base].append(f)
                
                for base, group in base_groups.items():
                    if len(group) > 1:
                        key = f"{h}_{base}"
                        duplicates[key] = group
        
        self.progress.emit(100, f"Gefunden: {len(duplicates)} Duplikat-Gruppen")
        self.results_ready.emit(duplicates)
        self.finished.emit()


# ============================================================================
# 5A. BATCH PROCESSING SYSTEM (NEU V14.2)
# ============================================================================

class BatchProcessor(QThread):
    """Worker-Thread für Batch-Operationen (NEU V14.2)
    
    Verarbeitet mehrere Dateien im Hintergrund mit Fortschrittsrückmeldung.
    Unterstützte Operationen: copy, pdf_encrypt, pdf_decrypt, pdf_extract_text
    """
    progress = pyqtSignal(int, int, str)  # current, total, current_filename
    file_completed = pyqtSignal(str, bool, str)  # filename, success, message
    all_completed = pyqtSignal(int, int)  # success_count, error_count
    log_message = pyqtSignal(str)  # log entry
    
    def __init__(self, files: list, operation: str, params: dict = None):
        super().__init__()
        self.files = files
        self.operation = operation
        self.params = params or {}
        self._cancelled = False
    
    def cancel(self):
        """Bricht die Verarbeitung ab"""
        self._cancelled = True
    
    def run(self):
        """Hauptverarbeitung"""
        total = len(self.files)
        success_count = 0
        error_count = 0
        
        for idx, filepath in enumerate(self.files):
            if self._cancelled:
                self.log_message.emit("⚠️ Verarbeitung abgebrochen")
                break
            
            filename = os.path.basename(filepath)
            self.progress.emit(idx + 1, total, filename)
            
            try:
                if self.operation == "copy":
                    result = self._copy_file(filepath)
                elif self.operation == "pdf_encrypt":
                    result = self._encrypt_pdf(filepath)
                elif self.operation == "pdf_decrypt":
                    result = self._decrypt_pdf(filepath)
                elif self.operation == "pdf_extract_text":
                    result = self._extract_pdf_text(filepath)
                else:
                    result = (False, f"Unbekannte Operation: {self.operation}")
                
                success, message = result
                if success:
                    success_count += 1
                    self.log_message.emit(f"✓ {filename} - OK")
                else:
                    error_count += 1
                    self.log_message.emit(f"✗ {filename} - {message}")
                
                self.file_completed.emit(filename, success, message)
                
            except Exception as e:
                error_count += 1
                self.log_message.emit(f"✗ {filename} - Fehler: {str(e)}")
                self.file_completed.emit(filename, False, str(e))
        
        self.all_completed.emit(success_count, error_count)
    
    def _copy_file(self, filepath):
        """Kopiert eine Datei in den Zielordner"""
        target_dir = self.params.get("target_dir")
        if not target_dir:
            return (False, "Kein Zielordner angegeben")
        
        filename = os.path.basename(filepath)
        target_path = os.path.join(target_dir, filename)
        
        # Duplikat-Handling
        if os.path.exists(target_path):
            base, ext = os.path.splitext(filename)
            counter = 1
            while os.path.exists(target_path):
                target_path = os.path.join(target_dir, f"{base}_{counter}{ext}")
                counter += 1
        
        shutil.copy2(filepath, target_path)
        return (True, f"Kopiert nach {os.path.basename(target_path)}")
    
    def _encrypt_pdf(self, filepath):
        """Verschlüsselt eine PDF-Datei"""
        if not HAS_PDF:
            return (False, "PyPDF2 nicht installiert")
        
        password = self.params.get("password")
        if not password:
            return (False, "Kein Passwort angegeben")
        
        keep_original = self.params.get("keep_original", True)
        output_dir = self.params.get("output_dir", os.path.dirname(filepath))
        
        base, ext = os.path.splitext(os.path.basename(filepath))
        output_path = os.path.join(output_dir, f"{base}_encrypted{ext}")
        
        try:
            reader = PdfReader(filepath)
            if reader.is_encrypted:
                return (False, "Bereits verschlüsselt")
            
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            
            writer.encrypt(password)
            
            with open(output_path, "wb") as f:
                writer.write(f)
            
            if not keep_original:
                os.remove(filepath)
            
            return (True, "Verschlüsselt")
        except Exception as e:
            return (False, str(e))
    
    def _decrypt_pdf(self, filepath):
        """Entschlüsselt eine PDF-Datei"""
        if not HAS_PDF:
            return (False, "PyPDF2 nicht installiert")
        
        password = self.params.get("password")
        if not password:
            return (False, "Kein Passwort angegeben")
        
        output_dir = self.params.get("output_dir", os.path.dirname(filepath))
        base, ext = os.path.splitext(os.path.basename(filepath))
        output_path = os.path.join(output_dir, f"{base}_decrypted{ext}")
        
        try:
            reader = PdfReader(filepath)
            if not reader.is_encrypted:
                return (False, "Nicht verschlüsselt")
            
            if not reader.decrypt(password):
                return (False, "Falsches Passwort")
            
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            
            with open(output_path, "wb") as f:
                writer.write(f)
            
            return (True, "Entschlüsselt")
        except Exception as e:
            return (False, str(e))
    
    def _extract_pdf_text(self, filepath):
        """Extrahiert Text aus einer PDF-Datei"""
        if not HAS_PDF:
            return (False, "PyPDF2 nicht installiert")
        
        output_dir = self.params.get("output_dir", os.path.dirname(filepath))
        base = os.path.splitext(os.path.basename(filepath))[0]
        output_path = os.path.join(output_dir, f"{base}.txt")
        
        try:
            reader = PdfReader(filepath)
            text_parts = []
            
            for i, page in enumerate(reader.pages):
                text = page.extract_text() or ""
                if text.strip():
                    text_parts.append(f"--- Seite {i+1} ---\n{text}")
            
            if not text_parts:
                return (False, "Kein Text gefunden")
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write("\n\n".join(text_parts))
            
            return (True, f"Text extrahiert ({len(text_parts)} Seiten)")
        except Exception as e:
            return (False, str(e))


class BatchDialog(QDialog):
    """Dialog für Batch-Operationen mit Fortschritt (NEU V14.2)
    
    Zeigt Fortschritt, Log und ermöglicht Abbrechen.
    """
    
    def __init__(self, files: list, operation: str, operation_title: str, parent=None):
        super().__init__(parent)
        self.files = files
        self.operation = operation
        self.operation_title = operation_title
        self.processor = None
        self.params = {}
        
        self.setWindowTitle(f"Batch: {operation_title}")
        self.setMinimumSize(500, 400)
        self.setModal(True)
        
        self._setup_ui()
    
    def _setup_ui(self):
        """Erstellt die UI-Komponenten"""
        layout = QVBoxLayout(self)
        
        # Header mit Dateianzahl
        header_label = QLabel(f"📁 {len(self.files)} Dateien ausgewählt")
        header_label.setStyleSheet("font-weight: bold; font-size: 14px;")
        layout.addWidget(header_label)
        
        # Parameter-Bereich (wird je nach Operation gefüllt)
        self.param_group = QGroupBox("Parameter")
        self.param_layout = QFormLayout(self.param_group)
        layout.addWidget(self.param_group)
        
        # Operation-spezifische Parameter
        self._setup_params()
        
        # Optionen
        self.skip_errors_cb = QCheckBox("Fehler überspringen und fortfahren")
        self.skip_errors_cb.setChecked(True)
        layout.addWidget(self.skip_errors_cb)
        
        # Fortschritts-Bereich
        progress_group = QGroupBox("Fortschritt")
        progress_layout = QVBoxLayout(progress_group)
        
        self.progress_bar = QProgressBar()
        self.progress_bar.setMinimum(0)
        self.progress_bar.setMaximum(len(self.files))
        self.progress_bar.setValue(0)
        progress_layout.addWidget(self.progress_bar)
        
        self.status_label = QLabel("Bereit")
        progress_layout.addWidget(self.status_label)
        
        layout.addWidget(progress_group)
        
        # Log-Bereich
        log_group = QGroupBox("Log")
        log_layout = QVBoxLayout(log_group)
        
        self.log_text = QTextEdit()
        self.log_text.setReadOnly(True)
        self.log_text.setMaximumHeight(150)
        log_layout.addWidget(self.log_text)
        
        layout.addWidget(log_group)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        self.start_btn = QPushButton("▶ Starten")
        self.start_btn.clicked.connect(self._start_processing)
        btn_layout.addWidget(self.start_btn)
        
        self.cancel_btn = QPushButton("Abbrechen")
        self.cancel_btn.clicked.connect(self._cancel_or_close)
        btn_layout.addWidget(self.cancel_btn)
        
        layout.addLayout(btn_layout)
    
    def _setup_params(self):
        """Richtet operation-spezifische Parameter ein"""
        if self.operation == "copy":
            self.target_dir_edit = QLineEdit()
            self.target_dir_edit.setPlaceholderText("Zielordner auswählen...")
            self.target_dir_btn = QPushButton("...")
            self.target_dir_btn.setMaximumWidth(30)
            self.target_dir_btn.clicked.connect(self._select_target_dir)
            
            dir_layout = QHBoxLayout()
            dir_layout.addWidget(self.target_dir_edit)
            dir_layout.addWidget(self.target_dir_btn)
            
            self.param_layout.addRow("Zielordner:", dir_layout)
        
        elif self.operation in ("pdf_encrypt", "pdf_decrypt"):
            self.password_edit = QLineEdit()
            self.password_edit.setEchoMode(QLineEdit.EchoMode.Password)
            
            self.show_pw_btn = QPushButton("👁")
            self.show_pw_btn.setMaximumWidth(30)
            self.show_pw_btn.pressed.connect(lambda: self.password_edit.setEchoMode(QLineEdit.EchoMode.Normal))
            self.show_pw_btn.released.connect(lambda: self.password_edit.setEchoMode(QLineEdit.EchoMode.Password))
            
            pw_layout = QHBoxLayout()
            pw_layout.addWidget(self.password_edit)
            pw_layout.addWidget(self.show_pw_btn)
            
            self.param_layout.addRow("Passwort:", pw_layout)
            
            if self.operation == "pdf_encrypt":
                self.password_confirm = QLineEdit()
                self.password_confirm.setEchoMode(QLineEdit.EchoMode.Password)
                self.param_layout.addRow("Bestätigen:", self.password_confirm)
                
                self.keep_original_cb = QCheckBox("Originale behalten")
                self.keep_original_cb.setChecked(True)
                self.param_layout.addRow("", self.keep_original_cb)
        
        elif self.operation == "pdf_extract_text":
            self.output_dir_edit = QLineEdit()
            self.output_dir_edit.setPlaceholderText("Ausgabeordner (leer = neben Original)")
            self.output_dir_btn = QPushButton("...")
            self.output_dir_btn.setMaximumWidth(30)
            self.output_dir_btn.clicked.connect(self._select_output_dir)
            
            dir_layout = QHBoxLayout()
            dir_layout.addWidget(self.output_dir_edit)
            dir_layout.addWidget(self.output_dir_btn)
            
            self.param_layout.addRow("Ausgabeordner:", dir_layout)
    
    def _select_target_dir(self):
        """Öffnet Dialog zur Ordnerauswahl"""
        folder = QFileDialog.getExistingDirectory(self, "Zielordner wählen")
        if folder:
            self.target_dir_edit.setText(folder)
    
    def _select_output_dir(self):
        """Öffnet Dialog zur Ausgabeordner-Auswahl"""
        folder = QFileDialog.getExistingDirectory(self, "Ausgabeordner wählen")
        if folder:
            self.output_dir_edit.setText(folder)
    
    def _validate_params(self) -> bool:
        """Prüft Parameter und gibt True bei Erfolg zurück"""
        if self.operation == "copy":
            target = self.target_dir_edit.text().strip()
            if not target or not os.path.isdir(target):
                QMessageBox.warning(self, "Fehler", "Bitte wählen Sie einen gültigen Zielordner!")
                return False
            self.params["target_dir"] = target
        
        elif self.operation == "pdf_encrypt":
            pw = self.password_edit.text()
            pw_confirm = self.password_confirm.text()
            if not pw:
                QMessageBox.warning(self, "Fehler", "Bitte geben Sie ein Passwort ein!")
                return False
            if pw != pw_confirm:
                QMessageBox.warning(self, "Fehler", "Passwörter stimmen nicht überein!")
                return False
            self.params["password"] = pw
            self.params["keep_original"] = self.keep_original_cb.isChecked()
        
        elif self.operation == "pdf_decrypt":
            pw = self.password_edit.text()
            if not pw:
                QMessageBox.warning(self, "Fehler", "Bitte geben Sie das Passwort ein!")
                return False
            self.params["password"] = pw
        
        elif self.operation == "pdf_extract_text":
            output_dir = self.output_dir_edit.text().strip()
            if output_dir and not os.path.isdir(output_dir):
                QMessageBox.warning(self, "Fehler", "Ungültiger Ausgabeordner!")
                return False
            if output_dir:
                self.params["output_dir"] = output_dir
        
        return True
    
    def _start_processing(self):
        """Startet die Batch-Verarbeitung"""
        if not self._validate_params():
            return
        
        # UI anpassen
        self.start_btn.setEnabled(False)
        self.param_group.setEnabled(False)
        self.cancel_btn.setText("Abbrechen")
        self.log_text.clear()
        
        # Processor starten
        self.processor = BatchProcessor(self.files, self.operation, self.params)
        self.processor.progress.connect(self._on_progress)
        self.processor.log_message.connect(self._on_log)
        self.processor.all_completed.connect(self._on_completed)
        self.processor.start()
    
    def _on_progress(self, current, total, filename):
        """Aktualisiert Fortschrittsanzeige"""
        self.progress_bar.setValue(current)
        self.status_label.setText(f"Verarbeite {current}/{total}: {filename}")
    
    def _on_log(self, message):
        """Fügt Log-Eintrag hinzu"""
        self.log_text.append(message)
    
    def _on_completed(self, success_count, error_count):
        """Wird aufgerufen wenn alle Dateien verarbeitet wurden"""
        self.status_label.setText(f"Fertig: {success_count} OK, {error_count} Fehler")
        self.start_btn.setText("Fertig")
        self.cancel_btn.setText("Schließen")
        
        # Abschluss-Meldung
        if error_count == 0:
            QMessageBox.information(
                self, "Abgeschlossen",
                f"Alle {success_count} Dateien erfolgreich verarbeitet!"
            )
        else:
            QMessageBox.warning(
                self, "Abgeschlossen mit Fehlern",
                f"Verarbeitet: {success_count} OK, {error_count} Fehler\n\n"
                "Details siehe Log."
            )
    
    def _cancel_or_close(self):
        """Bricht ab oder schließt den Dialog"""
        if self.processor and self.processor.isRunning():
            reply = QMessageBox.question(
                self, "Abbrechen",
                "Verarbeitung wirklich abbrechen?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.Yes:
                self.processor.cancel()
                self.processor.wait()
                self.reject()
        else:
            self.accept()
    
    def closeEvent(self, event):
        """Verhindert Schließen während Verarbeitung"""
        if self.processor and self.processor.isRunning():
            event.ignore()
            self._cancel_or_close()
        else:
            event.accept()


# ============================================================================
# 5B. VERBINDUNGEN-SYSTEM (aus ProFiler V4 + Multi-Source Enhancement)
# ============================================================================

class ConnectionConfigManager:
    """Verwaltet Verbindungen/Themen in JSON-Datei"""
    def __init__(self, path):
        self.path = path
        self.data = {"app": {}, "connections": []}
        self.load()
        
        # SearchConfigManager für Auto-Sync
        self.search_mgr = SearchConfigManager()
        
        # Migriere alte Verbindungen
        self.migrate_connections()
        
        # Initial-Sync aller enabled Verbindungen
        self.sync_all_to_search()
    
    def load(self):
        if os.path.exists(self.path):
            try:
                with open(self.path, "r", encoding="utf-8") as f:
                    self.data = json.load(f)
            except (OSError, json.JSONDecodeError):
                self.save()
        else:
            self.save()
    
    def save(self):
        os.makedirs(os.path.dirname(self.path) or ".", exist_ok=True)
        with open(self.path, "w", encoding="utf-8") as f:
            json.dump(self.data, f, indent=2, ensure_ascii=False)
    
    def list_connections(self):
        return self.data.get("connections", [])
    
    def add_or_update_connection(self, conn):
        conns = self.data.get("connections", [])
        for i, c in enumerate(conns):
            if c.get("id") == conn.get("id"):
                conns[i] = conn
                self.save()
                return
        conns.append(conn)
        self.data["connections"] = conns
        self.save()
    
    def remove_connection(self, conn_id):
        self.data["connections"] = [c for c in self.data.get("connections", []) 
                                     if c.get("id") != conn_id]
        self.save()
    
    def app_settings(self):
        return self.data.get("app", {})
    
    def update_app_settings(self, settings):
        self.data["app"] = settings
        self.save()
    
    def migrate_connections(self):
        """Migriert alte Verbindungen und fügt enabled/auto_update hinzu"""
        migrated = False
        for conn in self.data.get("connections", []):
            if "enabled" not in conn:
                conn["enabled"] = True
                migrated = True
            if "auto_update" not in conn:
                conn["auto_update"] = False
                migrated = True

        if migrated:
            self.save()
            print(" Verbindungen migriert (enabled/auto_update Flags hinzugefügt)")

    def sync_all_to_search(self):
        """Synchronisiert alle enabled Verbindungen mit SearchConfigManager"""
        # Leere Search-DB-Liste
        self.search_mgr.dbs = []

        # Füge alle enabled Verbindungen hinzu
        for conn in self.data.get("connections", []):
            if conn.get("enabled", True):
                db_path = conn.get("db_path")
                if db_path and os.path.exists(db_path):
                    self.search_mgr.add_db(db_path)

        print(f" {len(self.search_mgr.dbs)} Verbindungen für Suche aktiviert")

    def get_enabled_connections(self):
        """Gibt nur enabled Verbindungen zurück"""
        return [c for c in self.data.get("connections", []) if c.get("enabled", True)]

    def toggle_connection(self, conn_id, enabled):
        """Aktiviert/Deaktiviert eine Verbindung"""
        conns = self.data.get("connections", [])
        for conn in conns:
            if conn.get("id") == conn_id:
                conn["enabled"] = enabled
                self.save()

                # Sync mit SearchConfigManager
                if enabled:
                    db_path = conn.get("db_path")
                    if db_path:
                        self.search_mgr.add_db(db_path)
                else:
                    db_path = conn.get("db_path")
                    if db_path:
                        self.search_mgr.remove_db(db_path)

                return True
        return False

    def toggle_auto_update(self, conn_id, auto_update):
        """Schaltet Auto-Update für eine Verbindung ein/aus"""
        conns = self.data.get("connections", [])
        for conn in conns:
            if conn.get("id") == conn_id:
                conn["auto_update"] = auto_update
                self.save()
                return True
        return False




# Utility Functions für Verbindungen-System



def is_cloud_placeholder(path):
    """Prüft ob Datei nur in Cloud liegt (Windows-spezifisch)"""
    if os.name != 'nt':
        return False
    try:
        attrs = os.stat(path).st_file_attributes
        return (attrs & 0x1000) or (attrs & 0x400)
    except (OSError, AttributeError):
        return False


def path_to_tags(path, root):
    """Generiert Tags aus Ordnerstruktur"""
    try:
        rel = os.path.relpath(path, root)
        parts = rel.split(os.sep)
        return [p for p in parts[:-1] if p and p != "."]
    except (ValueError, OSError):
        return []


def shorten_filename(name, max_len):
    """Krzt Dateinamen wenn ntig"""
    root, ext = os.path.splitext(name)
    if len(name) <= max_len:
        return name
    keep = max(1, max_len - len(ext) - 1)
    return root[:keep] + "_" + ext



class SyncWorkerSignals(QObject):
    """Signals für SyncWorker"""
    status = pyqtSignal(str)
    progress = pyqtSignal(int, str)
    finished = pyqtSignal()


class SyncWorker(QThread):
    """Background Worker für Indizierung und Synchronisation"""
    
    def __init__(self, connection_config, mode="index"):
        super().__init__()
        self.cfg = connection_config
        self.mode = mode  # "index" or "sync"
        self.signals = SyncWorkerSignals()
        self.is_killed = False
        self.is_paused = False
    
    def run(self):
        try:
            db = ConnectionDB(self.cfg["db_path"])
            
            if self.mode == "index":
                self._run_indexing(db)
            elif self.mode == "sync":
                self._run_sync(db)
            
            db.close()
        except Exception as e:
            self.signals.status.emit(f"Fehler: {str(e)}")
        finally:
            self.signals.finished.emit()
    
    def _run_indexing(self, db):
        """Indiziert alle Source-Ordner (MULTI-SOURCE!)"""
        self.signals.status.emit("Starte Indizierung...")
        
        # NDERUNG: sources ist jetzt ein Array!
        source_folders = self.cfg.get("sources", [])
        if not source_folders:
            self.signals.status.emit("Keine Quellordner definiert!")
            return
        
        all_files = []
        
        # Alle Dateien aus allen Source-Ordnern sammeln
        for source_idx, source in enumerate(source_folders):
            if not os.path.exists(source):
                continue
            
            self.signals.status.emit(f"Durchsuche: {source} ({source_idx+1}/{len(source_folders)})")
            
            for root, _, files in os.walk(source):
                if self.is_killed:
                    return
                for f in files:
                    all_files.append((source, os.path.join(root, f)))
        
        total = len(all_files)
        done = 0
        
        self.signals.status.emit(f"Indiziere {total} Dateien...")
        
        for source_root, path in all_files:
            while self.is_paused:
                time.sleep(0.5)
                if self.is_killed:
                    return
            
            if self.is_killed:
                return
            
            try:
                stat = os.stat(path)
                mtime_iso = datetime.utcfromtimestamp(stat.st_mtime).isoformat()
                ctime_iso = datetime.utcfromtimestamp(stat.st_ctime).isoformat()
                size = stat.st_size
                name = os.path.basename(path)
                
                # Cloud-Placeholder Check
                is_cloud = is_cloud_placeholder(path)
                
                # Cache Check
                latest = db.get_latest_version_by_path(path)
                content_hash = None
                file_id = None
                
                if latest and latest[0] == mtime_iso:
                    file_id = latest[1]
                
                if not file_id:
                    if is_cloud:
                        content_hash = f"CLOUD:{size}:{mtime_iso}"
                    else:
                        content_hash = sha256_file(path)
                    
                    file_id = db.upsert_file(content_hash, size)
                
                # Tags aus Pfad
                if self.cfg.get("structure_handling", {}).get("tags_from_path", True):
                    for tag in path_to_tags(path, source_root):
                        db.add_tag(file_id, tag)
                
                # Version hinzufügen
                if not content_hash:
                    if is_cloud:
                        content_hash = f"CLOUD:{size}:{mtime_iso}"
                    else:
                        content_hash = sha256_file(path)
                
                existing = db.get_versions_by_hash(content_hash)
                v_idx = (existing[0][4] + 1) if existing else 1
                
                db.add_version(file_id, name, path, mtime_iso, ctime_iso, v_idx, "source")
                
                done += 1
                pct = int(done * 100 / max(1, total))
                self.signals.progress.emit(pct, f"{done}/{total} Dateien")
                
            except Exception as e:
                try:
                    db.add_event(file_id if file_id else -1, "error", str(e))
                except sqlite3.Error:
                    pass
        
        self.signals.status.emit(f"Indizierung abgeschlossen! {done} Dateien verarbeitet.")
        self.signals.progress.emit(100, "Fertig")
    
    def _run_sync(self, db):
        """Synchronisiert zum Ziel (falls konfiguriert)"""
        target = self.cfg.get("target", "")
        if not target:
            self.signals.status.emit("Kein Ziel-Ordner konfiguriert!")
            return
        
        self.signals.status.emit("Synchronisation noch nicht implementiert in V12")
        # TODO: Implementierung bei Bedarf

# ============================================================================
# 6. PDF DIALOGS
# ============================================================================

class PDFPasswordDialog(QDialog):
    """Dialog zum Verschlüsseln/Entschlüsseln von PDFs"""
    
    def __init__(self, file_paths, mode="encrypt", settings=None, parent=None):
        super().__init__(parent)
        self.file_paths = file_paths
        self.mode = mode  # "encrypt" oder "decrypt"
        self.settings = settings
        self.password = None
        self.use_master = False
        
        self.setWindowTitle("PDF Verschlüsselung" if mode == "encrypt" else "PDF Entschlüsselung")
        self.resize(500, 250)
        
        layout = QVBoxLayout(self)
        
        # Info
        if mode == "encrypt":
            info_text = f"{len(file_paths)} Datei(en) ausgewählt\n\nPasswort zum Verschlüsseln festlegen:"
        else:
            info_text = f" {len(file_paths)} verschlüsselte Datei(en)\n\nPasswort zum Entschlüsseln eingeben:"
        
        info = QLabel(info_text)
        info.setStyleSheet("font-size: 12px; padding: 10px; background-color: #2b2b2b; border-radius: 4px;")
        layout.addWidget(info)
        
        # Passwort-Gruppe
        pwd_group = QGroupBox("Passwort")
        pwd_layout = QVBoxLayout()
        
        # Individuelles Passwort
        self.radio_individual = QRadioButton("Individuelles Passwort eingeben")
        self.radio_individual.setChecked(True)
        self.radio_individual.toggled.connect(self.on_mode_changed)
        pwd_layout.addWidget(self.radio_individual)
        
        self.password_input = QLineEdit()
        self.password_input.setEchoMode(QLineEdit.EchoMode.Password)
        self.password_input.setPlaceholderText("Passwort eingeben...")
        pwd_layout.addWidget(self.password_input)
        
        # Masterpasswort
        if mode == "encrypt":
            master_saved = settings.get("pdf_master_password_save", "") if settings else ""
        else:
            master_saved = settings.get("pdf_master_password_open", "") if settings else ""
        
        has_master = bool(master_saved)
        
        self.radio_master = QRadioButton("Masterpasswort verwenden")
        self.radio_master.setEnabled(has_master)
        if not has_master:
            self.radio_master.setToolTip("Kein Masterpasswort hinterlegt. Bitte in Einstellungen konfigurieren.")
        self.radio_master.toggled.connect(self.on_mode_changed)
        pwd_layout.addWidget(self.radio_master)
        
        pwd_group.setLayout(pwd_layout)
        layout.addWidget(pwd_group)
        
        # Checkbox: Passwort anzeigen
        self.cb_show_password = QCheckBox("Passwort anzeigen")
        self.cb_show_password.toggled.connect(self.toggle_password_visibility)
        layout.addWidget(self.cb_show_password)
        
        layout.addStretch()
        
        # Buttons
        buttons = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        buttons.accepted.connect(self.accept_and_validate)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)
    
    def on_mode_changed(self):
        """Aktiviert/Deaktiviert Passwort-Eingabe"""
        self.password_input.setEnabled(self.radio_individual.isChecked())
    
    def toggle_password_visibility(self, checked):
        """Zeigt/Versteckt Passwort"""
        if checked:
            self.password_input.setEchoMode(QLineEdit.EchoMode.Normal)
        else:
            self.password_input.setEchoMode(QLineEdit.EchoMode.Password)
    
    def accept_and_validate(self):
        """Validiert Eingabe vor Akzeptieren"""
        if self.radio_master.isChecked():
            # Masterpasswort verwenden
            if self.mode == "encrypt":
                self.password = self.settings.get("pdf_master_password_save", "")
            else:
                self.password = self.settings.get("pdf_master_password_open", "")
            
            self.use_master = True
            
            if not self.password:
                QMessageBox.warning(self, "Fehler", "Kein Masterpasswort hinterlegt!")
                return
        else:
            # Individuelles Passwort
            pwd = self.password_input.text().strip()
            if not pwd:
                QMessageBox.warning(self, "Fehler", "Bitte Passwort eingeben!")
                return
            
            self.password = pwd
            self.use_master = False
        
        self.accept()


class PDFExcerptDialog(QDialog):
    """Dialog zum Erstellen von PDF-Auszgen"""
    
    def __init__(self, pdf_path, parent=None):
        super().__init__(parent)
        self.pdf_path = pdf_path
        self.selected_pages = set()
        self.current_page = 0
        self.total_pages = 0
        self.pdf_reader = None
        
        self.setWindowTitle(f"PDF-Auszug: {os.path.basename(pdf_path)}")
        self.resize(900, 700)
        
        if not self.load_pdf():
            QMessageBox.critical(self, "Fehler", "PDF konnte nicht geladen werden.")
            self.reject()
            return
        
        layout = QVBoxLayout(self)
        
        # Top: Info & Controls
        top_layout = QHBoxLayout()
        
        self.lbl_info = QLabel(f"Seite 1 von {self.total_pages}")
        self.lbl_info.setStyleSheet("font-size: 13px; font-weight: bold;")
        top_layout.addWidget(self.lbl_info)
        
        top_layout.addStretch()
        
        self.lbl_selected = QLabel("Ausgewhlt: 0 Seiten")
        self.lbl_selected.setStyleSheet("color: #90ee90; font-weight: bold;")
        top_layout.addWidget(self.lbl_selected)
        
        layout.addLayout(top_layout)
        
        # Middle: PDF Vorschau
        scroll = QScrollArea()
        scroll.setWidgetResizable(True)
        scroll.setStyleSheet("background-color: #1e1e1e; border: 1px solid #444;")
        
        self.preview_label = QLabel()
        self.preview_label.setAlignment(Qt.AlignmentFlag.AlignCenter)
        self.preview_label.setStyleSheet("background-color: white;")
        scroll.setWidget(self.preview_label)
        
        layout.addWidget(scroll)
        
        # Navigation & Selection
        nav_layout = QHBoxLayout()
        
        btn_prev = QPushButton("Vorherige")
        btn_prev.clicked.connect(self.prev_page)
        nav_layout.addWidget(btn_prev)
        
        self.btn_toggle = QPushButton(" Auswählen [SPACE]")
        self.btn_toggle.setStyleSheet("background-color: #2a82da; font-weight: bold; padding: 10px;")
        self.btn_toggle.clicked.connect(self.toggle_current_page)
        nav_layout.addWidget(self.btn_toggle)
        
        btn_next = QPushButton("Nchste ")
        btn_next.clicked.connect(self.next_page)
        nav_layout.addWidget(btn_next)
        
        layout.addLayout(nav_layout)
        
        # Shortcuts Info
        shortcuts = QLabel("Shortcuts: [SPACE] Auswahl umschalten | [][] Navigation | [A] Alle | [N] Keine")
        shortcuts.setStyleSheet("color: #888; font-size: 11px; padding: 5px;")
        layout.addWidget(shortcuts)
        
        # Bottom Buttons
        btn_layout = QHBoxLayout()
        
        btn_all = QPushButton("Alle auswählen [A]")
        btn_all.clicked.connect(self.select_all)
        btn_layout.addWidget(btn_all)
        
        btn_none = QPushButton("Keine [N]")
        btn_none.clicked.connect(self.select_none)
        btn_layout.addWidget(btn_none)
        
        btn_layout.addStretch()
        
        btn_save = QPushButton("Auszug speichern")
        btn_save.setStyleSheet("background-color: #28a745; color: white; font-weight: bold; padding: 10px;")
        btn_save.clicked.connect(self.save_excerpt)
        btn_layout.addWidget(btn_save)
        
        btn_cancel = QPushButton("Abbrechen")
        btn_cancel.clicked.connect(self.reject)
        btn_layout.addWidget(btn_cancel)
        
        layout.addLayout(btn_layout)
        
        # Initial page anzeigen
        self.show_current_page()
    
    def load_pdf(self):
        """Ldt PDF"""
        if not HAS_PDF:
            return False
        
        try:
            self.pdf_reader = PdfReader(self.pdf_path)
            self.total_pages = len(self.pdf_reader.pages)
            return self.total_pages > 0
        except Exception as e:
            print(f"PDF Load Error: {e}")
            return False
    
    def show_current_page(self):
        """Zeigt aktuelle Seite an"""
        if not self.pdf_reader or self.current_page >= self.total_pages:
            return
        
        # Update Info
        self.lbl_info.setText(f"Seite {self.current_page + 1} von {self.total_pages}")
        
        # Update Toggle Button
        if self.current_page in self.selected_pages:
            self.btn_toggle.setText(" Abwhlen [SPACE]")
            self.btn_toggle.setStyleSheet("background-color: #c9302c; font-weight: bold; padding: 10px;")
        else:
            self.btn_toggle.setText(" Auswählen [SPACE]")
            self.btn_toggle.setStyleSheet("background-color: #2a82da; font-weight: bold; padding: 10px;")
        
        # Update Selected Count
        self.lbl_selected.setText(f"Ausgewhlt: {len(self.selected_pages)} Seiten")
        
        # Vereinfachte Vorschau (nur Text)
        try:
            page = self.pdf_reader.pages[self.current_page]
            text = page.extract_text() or "Keine Textvorschau verfgbar"
            
            # Zeige Text in Label
            preview_text = f"<h3>Seite {self.current_page + 1}</h3><pre>{text[:500]}...</pre>"
            
            if self.current_page in self.selected_pages:
                preview_text = f"<div style='background-color: #d4edda; padding: 10px;'>{preview_text}<br><b> Diese Seite ist ausgewählt</b></div>"
            
            self.preview_label.setText(preview_text)
            self.preview_label.setWordWrap(True)
            
        except Exception as e:
            self.preview_label.setText(f"Vorschau-Fehler: {str(e)}")
    
    def toggle_current_page(self):
        """Whlt/Abwhlt aktuelle Seite"""
        if self.current_page in self.selected_pages:
            self.selected_pages.remove(self.current_page)
        else:
            self.selected_pages.add(self.current_page)
        
        self.show_current_page()
    
    def prev_page(self):
        """Vorherige Seite"""
        if self.current_page > 0:
            self.current_page -= 1
            self.show_current_page()
    
    def next_page(self):
        """Nchste Seite"""
        if self.current_page < self.total_pages - 1:
            self.current_page += 1
            self.show_current_page()
    
    def select_all(self):
        """Whlt alle Seiten"""
        self.selected_pages = set(range(self.total_pages))
        self.show_current_page()
    
    def select_none(self):
        """Whlt keine Seiten"""
        self.selected_pages.clear()
        self.show_current_page()
    
    def keyPressEvent(self, event):
        """Tastatur-Shortcuts"""
        key = event.key()
        
        if key == Qt.Key.Key_Space:
            self.toggle_current_page()
        elif key == Qt.Key.Key_Left:
            self.prev_page()
        elif key == Qt.Key.Key_Right:
            self.next_page()
        elif key == Qt.Key.Key_A:
            self.select_all()
        elif key == Qt.Key.Key_N:
            self.select_none()
        else:
            super().keyPressEvent(event)
    
    def save_excerpt(self):
        """Speichert Auszug"""
        if not self.selected_pages:
            QMessageBox.warning(self, "Keine Auswahl", "Bitte mindestens eine Seite auswählen!")
            return
        
        # Dateiname generieren
        base_dir = os.path.dirname(self.pdf_path)
        base_name = os.path.splitext(os.path.basename(self.pdf_path))[0]
        output_path = os.path.join(base_dir, f"{base_name}_Auszug.pdf")
        
        # Bei Konflikt: Suffix hinzufügen
        counter = 1
        while os.path.exists(output_path):
            output_path = os.path.join(base_dir, f"{base_name}_Auszug_{counter}.pdf")
            counter += 1
        
        try:
            # Sortierte Seitenliste
            sorted_pages = sorted(list(self.selected_pages))
            
            # PDF erstellen
            PDFUtils.extract_pages(self.pdf_path, output_path, sorted_pages)
            
            QMessageBox.information(
                self,
                "Erfolg",
                f" Auszug erstellt!\n\n{len(sorted_pages)} Seiten gespeichert:\n{output_path}"
            )
            
            self.accept()
            
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Auszug-Erstellung fehlgeschlagen:\n\n{str(e)}")


# ============================================================================
# PROFILER SUITE V9 - PHASE 1 CONTINUATION (Part 2/2)
# ============================================================================
# Diese Datei wird an Profiler_Suite_V9_Phase1.py angehngt
# ============================================================================

# 7. SETTINGS DIALOG MIT PDF-BEREICH
# ============================================================================

class SettingsDialog(QDialog):
    """Erweiterte Einstellungen mit PDF-Bereich"""
    
    def __init__(self, settings_manager, parent=None):
        super().__init__(parent)
        self.settings = settings_manager
        
        self.setWindowTitle("⚙️ Einstellungen")
        self.resize(600, 500)
        
        layout = QVBoxLayout(self)
        
        # Tab Widget
        tabs = QTabWidget()
        
        # Tab 1: Lösch-Verhalten
        delete_tab = self.create_delete_tab()
        tabs.addTab(delete_tab, "Löschen")
        
        # Tab 2: PDF-Einstellungen (NEU!)
        pdf_tab = self.create_pdf_tab()
        tabs.addTab(pdf_tab, " PDF")
        
        # Tab 3: Externe Tools (NEU V13!)
        tools_tab = self.create_tools_tab()
        tabs.addTab(tools_tab, " Externe Tools")
        
        layout.addWidget(tabs)
        
        # Buttons
        buttons = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Save | QDialogButtonBox.StandardButton.Cancel
        )
        buttons.accepted.connect(self.save_and_close)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)
    
    def create_delete_tab(self):
        """Tab für Lösch-Einstellungen"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        # Delete Mode
        mode_group = QGroupBox("Lösch-Modus")
        mode_layout = QVBoxLayout()
        
        self.radio_soft = QRadioButton("Soft-Delete (Papierkorb)")
        self.radio_soft.setToolTip("Dateien werden markiert, können wiederhergestellt werden")
        
        self.radio_hard = QRadioButton("Hard-Delete (Permanent)")
        self.radio_hard.setToolTip("Dateien werden sofort permanent gelöscht")
        
        # HIER WAR DER FEHLER: Definition muss VOR dem addWidget kommen
        self.radio_safety = QRadioButton("Safety-Mode (Nur Ausblenden)")
        self.radio_safety.setToolTip("Dateien werden NICHT gelöscht, nur in der Ansicht ausgeblendet.\nKeine Änderung auf Festplatte.")
        
        # Aktuellen Modus setzen
        mode = self.settings.get("delete_mode", "soft")
        if mode == "soft":
            self.radio_soft.setChecked(True)
        elif mode == "hard":
            self.radio_hard.setChecked(True)
        else:  # safety
            self.radio_safety.setChecked(True)
        
        # Jetzt erst zum Layout hinzufügen
        mode_layout.addWidget(self.radio_soft)
        mode_layout.addWidget(self.radio_hard)
        mode_layout.addWidget(self.radio_safety)
        mode_group.setLayout(mode_layout)
        layout.addWidget(mode_group)
        
        # Retention
        retention_group = QGroupBox("Aufbewahrungsdauer (Soft-Delete)")
        retention_layout = QFormLayout()
        
        self.spin_retention = QSpinBox()
        self.spin_retention.setRange(0, 365)
        self.spin_retention.setSuffix(" Tage")
        self.spin_retention.setValue(self.settings.get("trash_retention_days", 30))
        self.spin_retention.setSpecialValueText("Unbegrenzt")
        
        retention_layout.addRow("Gelöschte Dateien behalten:", self.spin_retention)
        retention_group.setLayout(retention_layout)
        layout.addWidget(retention_group)
        
        # Auto-Cleanup
        self.cb_auto_cleanup = QCheckBox("Automatisches Aufräumen beim Start")
        self.cb_auto_cleanup.setChecked(self.settings.get("auto_cleanup_enabled", True))
        layout.addWidget(self.cb_auto_cleanup)
        
        # Spawn-Einstellungen (NEU V13.2!)
        spawn_group = QGroupBox("Zwischenablage Spawning")
        spawn_layout = QFormLayout()
        
        self.combo_spawn_format = QComboBox()
        self.combo_spawn_format.addItems(["pdf", "txt", "docx", "rtf", "odt"])
        current_format = self.settings.get("default_spawn_format", "txt")
        index = self.combo_spawn_format.findText(current_format)
        if index >= 0:
            self.combo_spawn_format.setCurrentIndex(index)
        
        spawn_layout.addRow("Standard-Spawn-Format:", self.combo_spawn_format)
        spawn_group.setLayout(spawn_layout)
        layout.addWidget(spawn_group)

        # Umbenennen-Einstellungen (NEU V13.2!)
        rename_group = QGroupBox("Umbenennen")
        rename_layout = QVBoxLayout()
        
        self.cb_rename_filesystem = QCheckBox("Umbenennung wirkt im Dateisystem")
        self.cb_rename_filesystem.setToolTip(
            " AN: Datei wird auf Festplatte umbenannt\n"
            " AUS: Nur Anzeigename in ProFiler (Datei bleibt unverändert)"
        )
        self.cb_rename_filesystem.setChecked(self.settings.get("rename_in_filesystem", True))
        
        rename_layout.addWidget(self.cb_rename_filesystem)
        rename_group.setLayout(rename_layout)
        layout.addWidget(rename_group)
        
        layout.addStretch()
        return widget
    
    def create_pdf_tab(self):
        """Tab für PDF-Einstellungen (NEU!)"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        # Masterpasswrter
        pwd_group = QGroupBox("Masterpasswrter")
        pwd_layout = QFormLayout()
        
        info = QLabel("Masterpasswrter für schnelle Batch-Operationen")
        info.setStyleSheet("color: #888; font-size: 11px; padding: 5px;")
        pwd_layout.addRow(info)
        
        # Masterpasswort 1 (öffnen)
        self.master_pwd1 = QLineEdit()
        self.master_pwd1.setEchoMode(QLineEdit.EchoMode.Password)
        self.master_pwd1.setText(self.settings.get("pdf_master_password_open", ""))
        self.master_pwd1.setPlaceholderText("Leer = nicht gesetzt")
        
        pwd1_layout = QHBoxLayout()
        pwd1_layout.addWidget(self.master_pwd1)
        
        self.cb_show_pwd1 = QCheckBox("")
        self.cb_show_pwd1.setMaximumWidth(40)
        self.cb_show_pwd1.toggled.connect(
            lambda checked: self.master_pwd1.setEchoMode(
                QLineEdit.EchoMode.Normal if checked else QLineEdit.EchoMode.Password
            )
        )
        pwd1_layout.addWidget(self.cb_show_pwd1)
        
        pwd_layout.addRow("Masterpasswort 1 (öffnen):", pwd1_layout)
        
        # Masterpasswort 2 (Speichern)
        self.master_pwd2 = QLineEdit()
        self.master_pwd2.setEchoMode(QLineEdit.EchoMode.Password)
        self.master_pwd2.setText(self.settings.get("pdf_master_password_save", ""))
        self.master_pwd2.setPlaceholderText("Leer = nicht gesetzt")
        
        pwd2_layout = QHBoxLayout()
        pwd2_layout.addWidget(self.master_pwd2)
        
        self.cb_show_pwd2 = QCheckBox("")
        self.cb_show_pwd2.setMaximumWidth(40)
        self.cb_show_pwd2.toggled.connect(
            lambda checked: self.master_pwd2.setEchoMode(
                QLineEdit.EchoMode.Normal if checked else QLineEdit.EchoMode.Password
            )
        )
        pwd2_layout.addWidget(self.cb_show_pwd2)
        
        pwd_layout.addRow("Masterpasswort 2 (Speichern):", pwd2_layout)
        
        pwd_group.setLayout(pwd_layout)
        layout.addWidget(pwd_group)
        
        # OCR-Einstellungen
        ocr_group = QGroupBox("OCR (Texterkennung)")
        ocr_layout = QFormLayout()
        
        self.cb_ocr_enabled = QCheckBox("OCR aktiviert")
        self.cb_ocr_enabled.setChecked(self.settings.get("ocr_enabled", True))
        ocr_layout.addRow(self.cb_ocr_enabled)
        
        self.combo_ocr_lang = QComboBox()
        self.combo_ocr_lang.addItems(["deu", "eng", "fra", "spa", "ita"])
        current_lang = self.settings.get("ocr_language", "deu")
        idx = self.combo_ocr_lang.findText(current_lang)
        if idx >= 0:
            self.combo_ocr_lang.setCurrentIndex(idx)
        
        ocr_layout.addRow("OCR-Sprache:", self.combo_ocr_lang)
        
        ocr_info = QLabel("Benötigt: Tesseract-OCR installiert")
        ocr_info.setStyleSheet("color: #888; font-size: 10px;")
        ocr_layout.addRow(ocr_info)
        
        ocr_group.setLayout(ocr_layout)
        layout.addWidget(ocr_group)
        
        layout.addStretch()
        return widget
    
    
    def create_tools_tab(self):
        """Tab für Externe Tools (NEU!)"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        # PythonBox
        pythonbox_group = QGroupBox(" PythonBox")
        pythonbox_layout = QFormLayout()
        
        info = QLabel("Python-Entwicklungsumgebung öffnen")
        info.setStyleSheet("color: #888; font-size: 11px; padding: 5px;")
        pythonbox_layout.addRow(info)
        
        # PythonBox Pfad
        path_layout = QHBoxLayout()
        self.pythonbox_path = QLineEdit()
        self.pythonbox_path.setText(self.settings.get("pythonbox_path", ""))
        self.pythonbox_path.setPlaceholderText("Pfad zu PythonBox.py")
        path_layout.addWidget(self.pythonbox_path)
        
        btn_browse = QPushButton(" Durchsuchen")
        btn_browse.clicked.connect(self.browse_pythonbox)
        path_layout.addWidget(btn_browse)
        
        pythonbox_layout.addRow("PythonBox Pfad:", path_layout)
        
        # Test Button
        btn_test = QPushButton(" Test")
        btn_test.clicked.connect(self.test_pythonbox)
        pythonbox_layout.addRow(btn_test)
        
        pythonbox_group.setLayout(pythonbox_layout)
        layout.addWidget(pythonbox_group)
        
        # SQLiteViewer
        sqlite_group = QGroupBox(" SQLite Viewer")
        sqlite_layout = QFormLayout()
        
        path_layout2 = QHBoxLayout()
        self.sqlite_path = QLineEdit()
        self.sqlite_path.setText(self.settings.get("sqlite_viewer_path", ""))
        self.sqlite_path.setPlaceholderText("Pfad zu SQLiteViewer.py")
        path_layout2.addWidget(self.sqlite_path)
        
        btn_browse2 = QPushButton(" Durchsuchen")
        btn_browse2.clicked.connect(self.browse_sqlite)
        path_layout2.addWidget(btn_browse2)
        
        sqlite_layout.addRow("SQLiteViewer Pfad:", path_layout2)
        sqlite_group.setLayout(sqlite_layout)
        layout.addWidget(sqlite_group)
        
        # FormConstructor
        form_group = QGroupBox(" Form Constructor")
        form_layout = QFormLayout()
        
        path_layout3 = QHBoxLayout()
        self.formconstr_path = QLineEdit()
        self.formconstr_path.setText(self.settings.get("formconstructor_path", ""))
        self.formconstr_path.setPlaceholderText("Pfad zu FormConstructor_V1_5.py")
        path_layout3.addWidget(self.formconstr_path)
        
        btn_browse3 = QPushButton(" Durchsuchen")
        btn_browse3.clicked.connect(self.browse_formconstructor)
        path_layout3.addWidget(btn_browse3)
        
        form_layout.addRow("FormConstructor Pfad:", path_layout3)
        form_group.setLayout(form_layout)
        layout.addWidget(form_group)
        
        layout.addStretch()
        return widget
    
    def browse_pythonbox(self):
        """Durchsuche nach PythonBox.py"""
        path, _ = QFileDialog.getOpenFileName(
            self, "PythonBox.py auswählen", "", "Python Files (*.py)"
        )
        if path:
            self.pythonbox_path.setText(path)
    
    def browse_sqlite(self):
        """Durchsuche nach SQLiteViewer.py"""
        path, _ = QFileDialog.getOpenFileName(
            self, "SQLiteViewer.py auswählen", "", "Python Files (*.py)"
        )
        if path:
            self.sqlite_path.setText(path)
    
    def browse_formconstructor(self):
        """Durchsuche nach FormConstructor_V1_5.py"""
        path, _ = QFileDialog.getOpenFileName(
            self, "FormConstructor_V1_5.py auswählen", "", "Python Files (*.py)"
        )
        if path:
            self.formconstr_path.setText(path)
    
    def test_pythonbox(self):
        """Teste PythonBox-Pfad"""
        path = self.pythonbox_path.text().strip()
        if not path:
            QMessageBox.warning(self, "Fehler", "Kein Pfad angegeben!")
            return
        
        if not os.path.exists(path):
            QMessageBox.warning(self, "Fehler", "Datei nicht gefunden!")
            return
        
        try:
            subprocess.Popen([sys.executable, path])
            QMessageBox.information(self, "Erfolg", "PythonBox wurde gestartet!")
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")


    def save_and_close(self):
        """Speichert Einstellungen"""
        # Delete settings
        if self.radio_soft.isChecked():
            delete_mode = "soft"
        elif self.radio_hard.isChecked():
            delete_mode = "hard"
        else:  # Safety-Mode
            delete_mode = "safety"
        self.settings.set("delete_mode", delete_mode)
        self.settings.set("trash_retention_days", self.spin_retention.value())
        self.settings.set("auto_cleanup_enabled", self.cb_auto_cleanup.isChecked())
        self.settings.set("default_spawn_format", self.combo_spawn_format.currentText())
        self.settings.set("rename_in_filesystem", self.cb_rename_filesystem.isChecked())
        
        # PDF settings
        self.settings.set("pdf_master_password_open", self.master_pwd1.text().strip())
        self.settings.set("pdf_master_password_save", self.master_pwd2.text().strip())
        self.settings.set("ocr_enabled", self.cb_ocr_enabled.isChecked())
        self.settings.set("ocr_language", self.combo_ocr_lang.currentText())
        
        # External Tools settings (NEU V13!)
        self.settings.set("pythonbox_path", self.pythonbox_path.text().strip())
        self.settings.set("sqlite_viewer_path", self.sqlite_path.text().strip())
        self.settings.set("formconstructor_path", self.formconstr_path.text().strip())
        
        self.accept()


# ============================================================================
# 7A. NEUE DOKUMENTTYP-DIALOGE (PHASE 8)
# ============================================================================

class MaterialReferenceDialog(QDialog):
    """Dialog zum Anlegen von Materialverweisen"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Neuer Materialverweis")
        self.resize(600, 500)
        self.result_data = None
        self.init_ui()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Info Label
        info = QLabel("Erstellen Sie einen Verweis auf physisch vorhandenes Material:")
        info.setWordWrap(True)
        layout.addWidget(info)
        
        # Formular
        form = QFormLayout()
        
        self.bezeichnung_input = QLineEdit()
        self.bezeichnung_input.setPlaceholderText("z.B. Ordner mit Dokumenten")
        form.addRow("Bezeichnung*:", self.bezeichnung_input)
        
        self.raum_input = QLineEdit()
        self.raum_input.setPlaceholderText("z.B. Bro, Keller, Raum 204")
        form.addRow("Raum:", self.raum_input)
        
        self.regal_input = QLineEdit()
        self.regal_input.setPlaceholderText("z.B. Regal A, Schrank 3")
        form.addRow("Regal:", self.regal_input)
        
        self.bereich_input = QLineEdit()
        self.bereich_input.setPlaceholderText("z.B. Fach 2, links oben")
        form.addRow("Bereich:", self.bereich_input)
        
        self.inventar_input = QLineEdit()
        self.inventar_input.setPlaceholderText("z.B. INV-2024-001")
        form.addRow("Inventarnummer:", self.inventar_input)
        
        self.tags_input = QLineEdit()
        self.tags_input.setPlaceholderText("Tags mit Komma getrennt")
        form.addRow("Tags:", self.tags_input)
        
        self.beschreibung_input = QTextEdit()
        self.beschreibung_input.setPlaceholderText("Weitere Details zum Material...")
        self.beschreibung_input.setMaximumHeight(150)
        form.addRow("Beschreibung:", self.beschreibung_input)
        
        layout.addLayout(form)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        btn_save = QPushButton("Speichern")
        btn_save.clicked.connect(self.save_material)
        btn_save.setDefault(True)
        
        btn_cancel = QPushButton(" Abbrechen")
        btn_cancel.clicked.connect(self.reject)
        
        btn_layout.addStretch()
        btn_layout.addWidget(btn_save)
        btn_layout.addWidget(btn_cancel)
        
        layout.addLayout(btn_layout)
    
    def save_material(self):
        """Speichert Materialverweis"""
        bezeichnung = self.bezeichnung_input.text().strip()
        
        if not bezeichnung:
            QMessageBox.warning(self, "Fehler", "Bezeichnung ist erforderlich!")
            return
        
        self.result_data = {
            'bezeichnung': bezeichnung,
            'raum': self.raum_input.text().strip(),
            'regal': self.regal_input.text().strip(),
            'bereich': self.bereich_input.text().strip(),
            'inventarnummer': self.inventar_input.text().strip(),
            'tags': self.tags_input.text().strip(),
            'beschreibung': self.beschreibung_input.toPlainText().strip()
        }
        
        self.accept()


class PromptFileDialog(QDialog):
    """Dialog zum Anlegen/Bearbeiten von Prompt-Dateien"""
    
    def __init__(self, parent=None, existing_data=None):
        super().__init__(parent)
        self.setWindowTitle("Prompt-Datei" + (" bearbeiten" if existing_data else " erstellen"))
        self.resize(800, 700)
        self.result_data = None
        self.existing_data = existing_data
        self.init_ui()
        
        if existing_data:
            self.load_existing_data()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Info
        if not self.existing_data:
            info = QLabel("Erstellen Sie eine neue Prompt-Datei mit Versionsverwaltung:")
        else:
            info = QLabel("Neue Version wird in derselben Datei gespeichert:")
        info.setWordWrap(True)
        layout.addWidget(info)
        
        # Formular
        form = QFormLayout()
        
        self.name_input = QLineEdit()
        self.name_input.setPlaceholderText("z.B. Blog-Artikel Generator")
        if self.existing_data:
            self.name_input.setEnabled(False)
        form.addRow("Name*:", self.name_input)
        
        # Version (auto oder manuell)
        version_layout = QHBoxLayout()
        self.version_input = QLineEdit()
        self.version_input.setPlaceholderText("Auto-generiert")
        self.version_input.setEnabled(False)
        version_layout.addWidget(self.version_input)
        form.addRow("Version:", version_layout)
        
        self.zweck_input = QTextEdit()
        self.zweck_input.setPlaceholderText("Wofr ist dieser Prompt gedacht?")
        self.zweck_input.setMaximumHeight(80)
        form.addRow("Zweck/Beschreibung*:", self.zweck_input)
        
        self.prompt_input = QTextEdit()
        self.prompt_input.setPlaceholderText("Der eigentliche Prompt-Text...")
        self.prompt_input.setMinimumHeight(200)
        form.addRow("Prompt-Text*:", self.prompt_input)
        
        self.resultate_input = QTextEdit()
        self.resultate_input.setPlaceholderText("Ergebnisse/Outputs (optional)")
        self.resultate_input.setMaximumHeight(100)
        form.addRow("Resultate:", self.resultate_input)
        
        self.verbesserungen_input = QTextEdit()
        self.verbesserungen_input.setPlaceholderText("Ideen für Verbesserungen (optional)")
        self.verbesserungen_input.setMaximumHeight(80)
        form.addRow("Verbesserungsvorschlge:", self.verbesserungen_input)
        
        layout.addLayout(form)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        btn_save = QPushButton("Speichern")
        btn_save.clicked.connect(self.save_prompt)
        btn_save.setDefault(True)
        
        btn_cancel = QPushButton(" Abbrechen")
        btn_cancel.clicked.connect(self.reject)
        
        btn_layout.addStretch()
        btn_layout.addWidget(btn_save)
        btn_layout.addWidget(btn_cancel)
        
        layout.addLayout(btn_layout)
    
    def load_existing_data(self):
        """Ldt existierende Daten"""
        if not self.existing_data:
            return
        
        self.name_input.setText(self.existing_data.get('name', ''))
        # Zweck, Prompt etc. knnen leer sein für neue Version
    
    def save_prompt(self):
        """Speichert Prompt-Datei"""
        name = self.name_input.text().strip()
        zweck = self.zweck_input.toPlainText().strip()
        prompt_text = self.prompt_input.toPlainText().strip()
        
        if not name or not zweck or not prompt_text:
            QMessageBox.warning(self, "Fehler", "Name, Zweck und Prompt-Text sind erforderlich!")
            return
        
        # Version berechnen
        if self.existing_data:
            version = self.existing_data.get('version', 1) + 1
        else:
            version = 1
        
        self.result_data = {
            'name': name,
            'version': version,
            'zweck': zweck,
            'prompt_text': prompt_text,
            'resultate': self.resultate_input.toPlainText().strip(),
            'verbesserungen': self.verbesserungen_input.toPlainText().strip(),
            'created_at': datetime.now().isoformat(),
            'modified_at': datetime.now().isoformat()
        }
        
        self.accept()


class InternetResourceDialog(QDialog):
    """Dialog zum Anlegen von Internetressourcen"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle(" Neue Internetressource")
        self.resize(600, 450)
        self.result_data = None
        self.init_ui()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Info
        info = QLabel("Erstellen Sie einen Verweis auf eine Internetressource:")
        info.setWordWrap(True)
        layout.addWidget(info)
        
        # Formular
        form = QFormLayout()
        
        self.bezeichnung_input = QLineEdit()
        self.bezeichnung_input.setPlaceholderText("z.B. Anthropic Documentation")
        form.addRow("Bezeichnung*:", self.bezeichnung_input)
        
        self.anbieter_input = QLineEdit()
        self.anbieter_input.setPlaceholderText("z.B. Anthropic, GitHub, Stack Overflow")
        form.addRow("Anbieter:", self.anbieter_input)
        
        self.adresse_input = QLineEdit()
        self.adresse_input.setPlaceholderText("https://...")
        form.addRow("Adresse (URL)*:", self.adresse_input)
        
        self.tags_input = QLineEdit()
        self.tags_input.setPlaceholderText("Tags mit Komma getrennt")
        form.addRow("Tags:", self.tags_input)
        
        self.beschreibung_input = QTextEdit()
        self.beschreibung_input.setPlaceholderText("Beschreibung der Ressource...")
        self.beschreibung_input.setMaximumHeight(150)
        form.addRow("Beschreibung:", self.beschreibung_input)
        
        layout.addLayout(form)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        btn_save = QPushButton("Speichern")
        btn_save.clicked.connect(self.save_resource)
        btn_save.setDefault(True)
        
        btn_cancel = QPushButton(" Abbrechen")
        btn_cancel.clicked.connect(self.reject)
        
        btn_layout.addStretch()
        btn_layout.addWidget(btn_save)
        btn_layout.addWidget(btn_cancel)
        
        layout.addLayout(btn_layout)
    
    def save_resource(self):
        """Speichert Internetressource"""
        bezeichnung = self.bezeichnung_input.text().strip()
        adresse = self.adresse_input.text().strip()
        
        if not bezeichnung or not adresse:
            QMessageBox.warning(self, "Fehler", "Bezeichnung und Adresse sind erforderlich!")
            return
        
        # URL validieren
        if not adresse.startswith(('http://', 'https://', 'ftp://')):
            QMessageBox.warning(self, "Fehler", "Adresse muss mit http://, https:// oder ftp:// beginnen!")
            return
        
        self.result_data = {
            'bezeichnung': bezeichnung,
            'anbieter': self.anbieter_input.text().strip(),
            'adresse': adresse,
            'tags': self.tags_input.text().strip(),
            'beschreibung': self.beschreibung_input.toPlainText().strip()
        }
        
        self.accept()


class LiteratureReferenceDialog(QDialog):
    """Dialog zum Anlegen von Literaturverweisen"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setWindowTitle(" Neuer Literaturverweis")
        self.resize(650, 550)
        self.result_data = None
        self.init_ui()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Info
        info = QLabel("Erstellen Sie einen Literaturverweis:")
        info.setWordWrap(True)
        layout.addWidget(info)
        
        # Formular
        form = QFormLayout()
        
        self.titel_input = QLineEdit()
        self.titel_input.setPlaceholderText("z.B. Knstliche Intelligenz - Eine Einführung")
        form.addRow("Titel*:", self.titel_input)
        
        self.jahr_input = QLineEdit()
        self.jahr_input.setPlaceholderText("z.B. 2024")
        form.addRow("Jahr:", self.jahr_input)
        
        self.autoren_input = QLineEdit()
        self.autoren_input.setPlaceholderText("z.B. Max Mustermann, Anna Schmidt")
        form.addRow("Autoren:", self.autoren_input)
        
        self.verlag_input = QLineEdit()
        self.verlag_input.setPlaceholderText("z.B. Springer Verlag")
        form.addRow("Verlag:", self.verlag_input)
        
        self.isbn_input = QLineEdit()
        self.isbn_input.setPlaceholderText("z.B. 978-3-16-148410-0")
        form.addRow("ISBN:", self.isbn_input)
        
        # Physische Standortangaben
        self.raum_input = QLineEdit()
        self.raum_input.setPlaceholderText("z.B. Arbeitszimmer, Bibliothek")
        form.addRow("Raum:", self.raum_input)
        
        self.regal_input = QLineEdit()
        self.regal_input.setPlaceholderText("z.B. Regal B, links")
        form.addRow("Regal:", self.regal_input)
        
        # Online
        self.internet_input = QLineEdit()
        self.internet_input.setPlaceholderText("https://... (optional)")
        form.addRow("Internetlink:", self.internet_input)
        
        self.tags_input = QLineEdit()
        self.tags_input.setPlaceholderText("Tags mit Komma getrennt")
        form.addRow("Tags:", self.tags_input)
        
        self.beschreibung_input = QTextEdit()
        self.beschreibung_input.setPlaceholderText("Notizen, Zusammenfassung...")
        self.beschreibung_input.setMaximumHeight(120)
        form.addRow("Beschreibung:", self.beschreibung_input)
        
        layout.addLayout(form)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        btn_save = QPushButton("Speichern")
        btn_save.clicked.connect(self.save_reference)
        btn_save.setDefault(True)
        
        btn_cancel = QPushButton(" Abbrechen")
        btn_cancel.clicked.connect(self.reject)
        
        btn_layout.addStretch()
        btn_layout.addWidget(btn_save)
        btn_layout.addWidget(btn_cancel)
        
        layout.addLayout(btn_layout)
    
    def save_reference(self):
        """Speichert Literaturverweis"""
        titel = self.titel_input.text().strip()
        
        if not titel:
            QMessageBox.warning(self, "Fehler", "Titel ist erforderlich!")
            return
        
        self.result_data = {
            'titel': titel,
            'jahr': self.jahr_input.text().strip(),
            'autoren': self.autoren_input.text().strip(),
            'verlag': self.verlag_input.text().strip(),
            'isbn': self.isbn_input.text().strip(),
            'raum': self.raum_input.text().strip(),
            'regal': self.regal_input.text().strip(),
            'internet': self.internet_input.text().strip(),
            'tags': self.tags_input.text().strip(),
            'beschreibung': self.beschreibung_input.toPlainText().strip()
        }
        
        self.accept()


# ============================================================================
# ANONYMIZATION SETTINGS DIALOG
# ============================================================================

class AnonymizationSettingsDialog(QDialog):
    """Dialog für Blacklist/Whitelist Verwaltung"""
    
    def __init__(self, settings, parent=None):
        super().__init__(parent)
        self.settings = settings
        self.setWindowTitle("Anonymisierungs-Einstellungen")
        self.resize(700, 600)
        
        self.blacklist = self.settings.get("anonymization_blacklist", [])
        self.whitelist = self.settings.get("anonymization_whitelist", [])
        
        self.init_ui()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Tabs für Blacklist/Whitelist
        tabs = QTabWidget()
        
        # === BLACKLIST TAB ===
        blacklist_widget = QWidget()
        bl_layout = QVBoxLayout(blacklist_widget)
        
        # Input
        input_layout = QHBoxLayout()
        self.blacklist_input = QLineEdit()
        self.blacklist_input.setPlaceholderText("Wort zur Blacklist hinzufügen...")
        self.blacklist_input.returnPressed.connect(self.add_to_blacklist)
        btn_add_bl = QPushButton(" Hinzufügen")
        btn_add_bl.clicked.connect(self.add_to_blacklist)
        input_layout.addWidget(self.blacklist_input)
        input_layout.addWidget(btn_add_bl)
        
        # Liste
        self.blacklist_widget = QListWidget()
        self.blacklist_widget.setSelectionMode(QListWidget.SelectionMode.ExtendedSelection)
        self.update_blacklist_display()
        
        # Buttons
        btn_layout = QHBoxLayout()
        btn_import_bl = QPushButton(" Importieren...")
        btn_import_bl.clicked.connect(lambda: self.import_list("blacklist"))
        btn_export_bl = QPushButton("Exportieren...")
        btn_export_bl.clicked.connect(lambda: self.export_list("blacklist"))
        btn_remove_bl = QPushButton("Ausgewhlte entfernen")
        btn_remove_bl.clicked.connect(self.remove_from_blacklist)
        btn_clear_bl = QPushButton(" Liste leeren")
        btn_clear_bl.clicked.connect(lambda: self.clear_list("blacklist"))
        btn_clear_bl.setStyleSheet("background-color: #d9534f; color: white;")
        
        btn_layout.addWidget(btn_import_bl)
        btn_layout.addWidget(btn_export_bl)
        btn_layout.addWidget(btn_remove_bl)
        btn_layout.addWidget(btn_clear_bl)
        
        bl_layout.addLayout(input_layout)
        bl_layout.addWidget(self.blacklist_widget)
        bl_layout.addLayout(btn_layout)
        
        # === WHITELIST TAB ===
        whitelist_widget = QWidget()
        wl_layout = QVBoxLayout(whitelist_widget)
        
        # Input
        input_layout2 = QHBoxLayout()
        self.whitelist_input = QLineEdit()
        self.whitelist_input.setPlaceholderText("Wort zur Whitelist hinzufügen...")
        self.whitelist_input.returnPressed.connect(self.add_to_whitelist)
        btn_add_wl = QPushButton(" Hinzufügen")
        btn_add_wl.clicked.connect(self.add_to_whitelist)
        input_layout2.addWidget(self.whitelist_input)
        input_layout2.addWidget(btn_add_wl)
        
        # Liste
        self.whitelist_widget = QListWidget()
        self.whitelist_widget.setSelectionMode(QListWidget.SelectionMode.ExtendedSelection)
        self.update_whitelist_display()
        
        # Buttons
        btn_layout2 = QHBoxLayout()
        btn_import_wl = QPushButton(" Importieren...")
        btn_import_wl.clicked.connect(lambda: self.import_list("whitelist"))
        btn_export_wl = QPushButton("Exportieren...")
        btn_export_wl.clicked.connect(lambda: self.export_list("whitelist"))
        btn_remove_wl = QPushButton("Ausgewhlte entfernen")
        btn_remove_wl.clicked.connect(self.remove_from_whitelist)
        btn_clear_wl = QPushButton(" Liste leeren")
        btn_clear_wl.clicked.connect(lambda: self.clear_list("whitelist"))
        btn_clear_wl.setStyleSheet("background-color: #d9534f; color: white;")
        
        btn_layout2.addWidget(btn_import_wl)
        btn_layout2.addWidget(btn_export_wl)
        btn_layout2.addWidget(btn_remove_wl)
        btn_layout2.addWidget(btn_clear_wl)
        
        wl_layout.addLayout(input_layout2)
        wl_layout.addWidget(self.whitelist_widget)
        wl_layout.addLayout(btn_layout2)
        
        tabs.addTab(blacklist_widget, "Blacklist")
        tabs.addTab(whitelist_widget, " Whitelist")
        
        # === EINSTELLUNGEN ===
        settings_group = QGroupBox("⚙️ Einstellungen")
        settings_layout = QFormLayout()
        
        self.placeholder_input = QLineEdit()
        self.placeholder_input.setText(self.settings.get("anonymization_placeholder", "[-----]"))
        settings_layout.addRow("Platzhalter:", self.placeholder_input)
        
        settings_group.setLayout(settings_layout)
        
        # === BUTTONS ===
        buttons = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        buttons.accepted.connect(self.save_and_close)
        buttons.rejected.connect(self.reject)
        
        layout.addWidget(tabs)
        layout.addWidget(settings_group)
        layout.addWidget(buttons)
    
    def add_to_blacklist(self):
        word = self.blacklist_input.text().strip()
        if word and word not in self.blacklist:
            self.blacklist.append(word)
            self.update_blacklist_display()
            self.blacklist_input.clear()
    
    def add_to_whitelist(self):
        word = self.whitelist_input.text().strip()
        if word and word not in self.whitelist:
            self.whitelist.append(word)
            self.update_whitelist_display()
            self.whitelist_input.clear()
    
    def remove_from_blacklist(self):
        for item in self.blacklist_widget.selectedItems():
            self.blacklist.remove(item.text())
        self.update_blacklist_display()
    
    def remove_from_whitelist(self):
        for item in self.whitelist_widget.selectedItems():
            self.whitelist.remove(item.text())
        self.update_whitelist_display()
    
    def update_blacklist_display(self):
        self.blacklist_widget.clear()
        for word in sorted(self.blacklist):
            self.blacklist_widget.addItem(word)
    
    def update_whitelist_display(self):
        self.whitelist_widget.clear()
        for word in sorted(self.whitelist):
            self.whitelist_widget.addItem(word)
    
    def import_list(self, list_type):
        """Importiert Liste aus TXT, Excel, Word oder PDF"""
        file_path, _ = QFileDialog.getOpenFileName(
            self,
            f"{list_type.capitalize()} importieren",
            "",
            "Alle Dateien (*.txt *.xlsx *.docx *.pdf);;Text (*.txt);;Excel (*.xlsx);;Word (*.docx);;PDF (*.pdf)"
        )
        
        if not file_path:
            return
        
        try:
            words = []
            ext = os.path.splitext(file_path)[1].lower()
            
            if ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    words = [line.strip() for line in f if line.strip()]
            
            elif ext == '.xlsx':
                try:
                    import pandas as pd
                    df = pd.read_excel(file_path)
                    words = df.values.flatten().astype(str).tolist()
                    words = [w.strip() for w in words if str(w).strip() and str(w).strip() != 'nan']
                except ImportError:
                    QMessageBox.warning(self, "Fehler", "pandas nicht installiert für Excel-Import")
                    return
            
            elif ext == '.docx':
                if HAS_DOCX:
                    doc = docx.Document(file_path)
                    for para in doc.paragraphs:
                        if para.text.strip():
                            words.extend([w.strip() for w in para.text.split() if w.strip()])
                else:
                    QMessageBox.warning(self, "Fehler", "python-docx nicht installiert")
                    return
            
            elif ext == '.pdf':
                if HAS_PDF:
                    reader = PdfReader(file_path)
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            words.extend([w.strip() for w in text.split() if w.strip()])
                else:
                    QMessageBox.warning(self, "Fehler", "PyPDF2 nicht installiert")
                    return
            
            # Fge Wrter hinzu
            target_list = self.blacklist if list_type == "blacklist" else self.whitelist
            count = 0
            for word in words:
                if word and word not in target_list:
                    target_list.append(word)
                    count += 1
            
            if list_type == "blacklist":
                self.update_blacklist_display()
            else:
                self.update_whitelist_display()
            
            QMessageBox.information(self, "Erfolg", f"{count} Wrter importiert")
        
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Import fehlgeschlagen: {str(e)}")
    
    def export_list(self, list_type):
        """Exportiert Liste als TXT"""
        file_path, _ = QFileDialog.getSaveFileName(
            self,
            f"{list_type.capitalize()} exportieren",
            f"{list_type}.txt",
            "Text (*.txt)"
        )
        
        if not file_path:
            return
        
        try:
            target_list = self.blacklist if list_type == "blacklist" else self.whitelist
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write('\n'.join(sorted(target_list)))
            
            QMessageBox.information(self, "Erfolg", f"{len(target_list)} Wrter exportiert")
        
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Export fehlgeschlagen: {str(e)}")
    
    def clear_list(self, list_type):
        """Leert Liste"""
        reply = QMessageBox.question(
            self,
            "Liste leeren",
            f"{list_type.capitalize()} wirklich leeren?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            if list_type == "blacklist":
                self.blacklist.clear()
                self.update_blacklist_display()
            else:
                self.whitelist.clear()
                self.update_whitelist_display()
    
    def save_and_close(self):
        """Speichert Einstellungen"""
        self.settings.set("anonymization_blacklist", self.blacklist)
        self.settings.set("anonymization_whitelist", self.whitelist)
        self.settings.set("anonymization_placeholder", self.placeholder_input.text())
        self.accept()



# ============================================================================
# 7B. PYTHON-TOOLS FUNKTIONEN (PHASE 9)
# ============================================================================

def pycutter_split_classes(filepath, output_format='txt'):
    """
    Zerlegt Python-Datei in Klassen (aus pyCuttertxt.py recycelt)
    
    Args:
        filepath: Pfad zur .py Datei
        output_format: 'txt' oder 'py'
    
    Returns:
        Tuple (success, output_dir, error_message)
    """
    import ast
    from datetime import datetime
    
    try:
        # Ausgabeordner erzeugen
        base = os.path.basename(filepath)
        name, _ = os.path.splitext(base)
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        outdir = os.path.join(os.path.dirname(filepath), f"pyCutter_{name}_{timestamp}")
        os.makedirs(outdir, exist_ok=True)
        
        with open(filepath, "r", encoding="utf-8") as f:
            source = f.read()
        
        tree = ast.parse(source)
        
        # Klassen und Funktionen sammeln
        classes = [node for node in tree.body if isinstance(node, ast.ClassDef)]
        functions = [node for node in tree.body if isinstance(node, ast.FunctionDef)]
        imports = [node for node in tree.body if isinstance(node, (ast.Import, ast.ImportFrom))]
        
        # Klassen speichern
        for cls in classes:
            start_line = cls.lineno - 1
            end_line = max(getattr(cls, "end_lineno", start_line), start_line)
            code = "\n".join(source.splitlines()[start_line:end_line])
            
            ext = '.txt' if output_format == 'txt' else '.py'
            with open(os.path.join(outdir, f"{cls.name}{ext}"), "w", encoding="utf-8") as out:
                out.write(code)
        
        # Hilfsfunktionen + Imports + globaler Text
        helper_lines = []
        lines = source.splitlines()
        
        # Imports
        for imp in imports:
            start, end = imp.lineno - 1, getattr(imp, "end_lineno", imp.lineno) - 1
            helper_lines.extend(lines[start:end+1])
        
        # Funktionen
        for func in functions:
            start, end = func.lineno - 1, getattr(func, "end_lineno", func.lineno) - 1
            helper_lines.extend(lines[start:end+1])
        
        # Restlicher Top-Level-Code
        occupied = set()
        for node in classes + functions + imports:
            occupied.update(range(node.lineno - 1, getattr(node, "end_lineno", node.lineno)))
        for i, line in enumerate(lines):
            if i not in occupied and line.strip():
                helper_lines.append(line)
        
        if helper_lines:
            ext = '.txt' if output_format == 'txt' else '.py'
            with open(os.path.join(outdir, f"Hilfsfunktionen{ext}"), "w", encoding="utf-8") as out:
                out.write("\n".join(helper_lines))
        
        return True, outdir, None
    
    except Exception as e:
        return False, None, str(e)


def encoding_fix_file(filepath, use_ftfy=False):
    """
    Repariert Encoding-Probleme in Python-Datei (aus EncodingFixxer.py)
    
    Args:
        filepath: Pfad zur .py Datei
        use_ftfy: ftfy verwenden wenn verfgbar
    
    Returns:
        Tuple (success, backup_path, error_message)
    """
    try:
        # Lese Datei
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()
        
        # Versuche ftfy zu verwenden
        if use_ftfy:
            try:
                from ftfy import fix_text
                fixed_content = fix_text(content)
            except ImportError:
                # Fallback: Einfache Encoding-Korrektur
                fixed_content = content.encode('utf-8', errors='ignore').decode('utf-8')
        else:
            # Einfache Korrektur
            fixed_content = content.encode('utf-8', errors='ignore').decode('utf-8')
        
        # Backup erstellen
        backup_path = filepath + ".bak"
        shutil.copy2(filepath, backup_path)
        
        # Korrigierte Version speichern
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(fixed_content)
        
        return True, backup_path, None
    
    except Exception as e:
        return False, None, str(e)


def check_indentation_errors(filepath):
    """
    Prüft Python-Datei auf Einrckungsfehler (aus indent_gui_checker.py)
    
    Returns:
        List of error messages
    """
    import re
    
    errors = []
    
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            lines = f.readlines()
        
        for i, line in enumerate(lines):
            stripped = line.strip()
            indent_level = len(line) - len(line.lstrip())
            
            # Strukturen ohne ':'
            if re.match(r"^(def|if|elif|else|for|while|try|except|class)\b", stripped) and not stripped.endswith(":"):
                errors.append(f"Zeile {i+1}: Struktur ohne ':'")
            
            # return außerhalb von Block
            if stripped.startswith("return") and indent_level == 0:
                errors.append(f"Zeile {i+1}: 'return' außerhalb von Block?")
            
            # Mischung Tab/Leerzeichen
            if "\t" in line and " " in line[:line.find("\t")]:
                errors.append(f"Zeile {i+1}: Mischung aus Tab & Leerzeichen")
        
        return errors
    
    except Exception as e:
        return [f"Fehler beim Lesen: {str(e)}"]


def parse_browser_bookmarks(html_path):
    """
    Parst Browser-Favoriten aus HTML (aus FavoritenExtraktor.py recycelt)
    
    Returns:
        List of dicts with 'title', 'url', 'folder'
    """
    from html.parser import HTMLParser
    
    class SimpleBookmarkParser(HTMLParser):
        def __init__(self):
            super().__init__()
            self.bookmarks = []
            self.current_folder = []
            self.in_a = False
            self.current_link = {}
        
        def handle_starttag(self, tag, attrs):
            if tag.lower() == 'h3':
                # Neuer Ordner
                pass
            elif tag.lower() == 'a':
                href = None
                for k, v in attrs:
                    if k.lower() == 'href':
                        href = v
                        break
                self.current_link = {'url': href, 'title': '', 'folder': '/'.join(self.current_folder)}
                self.in_a = True
        
        def handle_data(self, data):
            text = data.strip()
            if self.in_a and text:
                self.current_link['title'] = text
        
        def handle_endtag(self, tag):
            if tag.lower() == 'a' and self.in_a:
                if self.current_link.get('url'):
                    self.bookmarks.append(self.current_link)
                self.current_link = {}
                self.in_a = False
    
    try:
        # Versuche verschiedene Encodings
        content = None
        for encoding in ['utf-8-sig', 'utf-8', 'cp1252', 'latin-1']:
            try:
                with open(html_path, 'r', encoding=encoding) as f:
                    content = f.read()
                break
            except:
                continue
        
        if not content:
            return []
        
        parser = SimpleBookmarkParser()
        parser.feed(content)
        return parser.bookmarks
    
    except Exception as e:
        return []

# ============================================================================
# 7B. CONNECTION DIALOG (Multi-Source Support)
# ============================================================================

class PathLineEdit(QLineEdit):
    """Helper Widget für Pfad-Eingabe mit Drag&Drop"""
    def __init__(self, text=""):
        super().__init__(text)
        self.setAcceptDrops(True)
    
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
    
    def dropEvent(self, event):
        if event.mimeData().hasUrls():
            url = event.mimeData().urls()[0]
            path = url.toLocalFile()
            if os.path.isdir(path):
                self.setText(path)


class ConnectionDialog(QDialog):
    """Dialog zum Erstellen/Bearbeiten einer Verbindung (Thema)"""
    
    def __init__(self, parent=None, existing=None):
        super().__init__(parent)
        self.setWindowTitle("Verbindung/Thema konfigurieren")
        self.existing = existing
        self.resize(650, 700)
        
        layout = QVBoxLayout(self)
        form = QFormLayout()
        
        # Name (= Thema)
        self.name = QLineEdit(existing["name"] if existing else "")
        self.name.setPlaceholderText("z.B. Forschungsprojekt, Privat, Arbeitsordner...")
        form.addRow(" Name/Thema:", self.name)
        
        # === MULTI-SOURCE LISTE ===
        source_group = QGroupBox("Quellordner (mehrere möglich)")
        source_layout = QVBoxLayout()
        
        self.source_list = QListWidget()
        self.source_list.setMaximumHeight(150)
        
        # Bestehende Sources laden
        if existing and "sources" in existing:
            for src in existing["sources"]:
                self.source_list.addItem(src)
        
        source_buttons = QHBoxLayout()
        btn_add_source = QPushButton("+ Ordner hinzufügen")
        btn_remove_source = QPushButton("- Entfernen")
        btn_add_source.clicked.connect(self.add_source_folder)
        btn_remove_source.clicked.connect(self.remove_source_folder)
        
        source_buttons.addWidget(btn_add_source)
        source_buttons.addWidget(btn_remove_source)
        source_buttons.addStretch()
        
        source_layout.addWidget(self.source_list)
        source_layout.addLayout(source_buttons)
        source_group.setLayout(source_layout)
        
        # Ziel (optional)
        self.target = PathLineEdit(existing.get("target", "") if existing else "")
        self.target.setPlaceholderText("Optional: Backup/Sync-Ziel")
        btn_tgt = QPushButton("...")
        btn_tgt.clicked.connect(lambda: self.pick_folder(self.target))
        h_tgt = QHBoxLayout()
        h_tgt.addWidget(self.target)
        h_tgt.addWidget(btn_tgt)
        
        form.addRow(source_group)
        form.addRow(" Ziel (optional):", h_tgt)
        
        # Modus
        self.direction = QComboBox()
        self.direction.addItems(["index_only", "one_way", "two_way"])
        if existing:
            self.direction.setCurrentText(existing.get("direction", "index_only"))
        form.addRow(" Modus:", self.direction)
        
        # Optionen
        self.physical = QCheckBox("Physische Struktur behalten")
        self.physical.setChecked(True)
        self.tags_from_path = QCheckBox("Tags aus Ordnern generieren")
        self.tags_from_path.setChecked(True)
        
        form.addRow(self.physical)
        form.addRow(self.tags_from_path)
        
        # Datenbank
        self.db_path = QLineEdit(existing.get("db_path", "") if existing else "")
        self.db_path.setPlaceholderText("Wird automatisch generiert, falls leer")
        btn_db = QPushButton("...")
        btn_db.clicked.connect(lambda: self.pick_db_file())
        h_db = QHBoxLayout()
        h_db.addWidget(self.db_path)
        h_db.addWidget(btn_db)
        
        form.addRow("Datenbank:", h_db)
        
        # Info
        info = QLabel("Die Datenbank speichert Hashes, Versionen und Tags aller Dateien.\n"
                      "    Mehrere Quellordner werden in einer gemeinsamen DB indiziert.")
        info.setStyleSheet("color: gray; font-size: 10px;")
        form.addRow(info)
        
        layout.addLayout(form)
        
        # Buttons
        buttons = QDialogButtonBox(
            QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel
        )
        buttons.accepted.connect(self.accept)
        buttons.rejected.connect(self.reject)
        layout.addWidget(buttons)
    
    def add_source_folder(self):
        """Fgt einen neuen Quellordner hinzu"""
        folder = QFileDialog.getExistingDirectory(self, "Quellordner wählen")
        if folder:
            # Prfen ob schon vorhanden
            existing = [self.source_list.item(i).text() 
                       for i in range(self.source_list.count())]
            if folder not in existing:
                self.source_list.addItem(folder)
    
    def remove_source_folder(self):
        """Entfernt ausgewählten Quellordner"""
        for item in self.source_list.selectedItems():
            self.source_list.takeItem(self.source_list.row(item))
    
    def pick_folder(self, line_edit):
        """Ordner-Auswahl für Target"""
        folder = QFileDialog.getExistingDirectory(self, "Ordner wählen")
        if folder:
            line_edit.setText(folder)
    
    def pick_db_file(self):
        """Datei-Auswahl für Datenbank"""
        file_path, _ = QFileDialog.getSaveFileName(
            self, "Datenbank speichern", "", "SQLite (*.db)"
        )
        if file_path:
            self.db_path.setText(file_path)
    
    def get_result(self):
        """Gibt Verbindungs-Dictionary zurück"""
        
        # Sources sammeln
        sources = []
        for i in range(self.source_list.count()):
            sources.append(self.source_list.item(i).text())
        
        if not sources:
            QMessageBox.warning(self, "Fehler", "Mindestens ein Quellordner muss angegeben werden!")
            return None
        
        # DB-Pfad generieren falls leer
        db_p = self.db_path.text().strip()
        if not db_p:
            # Fallback: Neben erstem Source-Ordner
            first_source = sources[0]
            parent_dir = os.path.dirname(first_source)
            safe_name = "".join(c for c in self.name.text() if c.isalnum() or c in (' ', '_'))
            db_p = os.path.join(parent_dir, f"profiler_{safe_name}.db")
        
        return {
            "id": self.existing["id"] if self.existing else f"conn-{uuid.uuid4().hex[:8]}",
            "name": self.name.text(),
            "sources": sources,  # ARRAY statt single string!
            "target": self.target.text(),
            "direction": self.direction.currentText(),
            "structure_handling": {
                "physical_copy": self.physical.isChecked(),
                "tags_from_path": self.tags_from_path.isChecked()
            },
            "versioning": {"enabled": True, "policy": "count", "max_count": 10},
            "db_path": db_p,
            "max_filename_length": 120
        }

# ============================================================================
# CONNECTION INDEXING WORKER (from ProFiler V4)
# ============================================================================

class IndexWorkerSignals(QObject):
    """Signals for IndexWorker"""
    status = pyqtSignal(str)
    progress = pyqtSignal(int, str)  # (percentage, phase_name)
    finished = pyqtSignal()
    error = pyqtSignal(str)

class IndexWorker(QThread):
    """Worker thread for indexing files from multiple source folders"""
    
    def __init__(self, connection_config):
        super().__init__()
        self.cfg = connection_config
        self.signals = IndexWorkerSignals()
        self.is_killed = False
        self.is_paused = False
    
    def run(self):
        try:
            db = ConnectionDB(self.cfg["db_path"])
            self._run_indexing(db)
            db.close()
            self.signals.status.emit("Indexierung abgeschlossen!")
            self.signals.finished.emit()
        except Exception as e:
            self.signals.error.emit(f"Fehler: {str(e)}")
    
    def _run_indexing(self, db):
        """Index all files from all source folders"""
        self.signals.status.emit("Starte Indizierung...")
        self.signals.progress.emit(0, "indexing")
        
        # Collect all files from all sources
        all_files = []
        sources = self.cfg.get("sources", [])
        
        for source_folder in sources:
            if not source_folder or not os.path.exists(source_folder):
                continue
            
            for root, _, files in os.walk(source_folder):
                if self.is_killed:
                    return
                for filename in files:
                    filepath = os.path.join(root, filename)
                    all_files.append((source_folder, filepath))
        
        total = len(all_files)
        done = 0
        
        self.signals.status.emit(f"Gefunden: {total} Dateien")
        
        # Index each file
        for source_folder, filepath in all_files:
            while self.is_paused:
                time.sleep(0.5)
                if self.is_killed:
                    return
            
            if self.is_killed:
                return
            
            try:
                # Get file metadata
                stat = os.stat(filepath)
                mtime_iso = datetime.fromtimestamp(stat.st_mtime).isoformat()
                ctime_iso = datetime.fromtimestamp(stat.st_ctime).isoformat()
                size = stat.st_size
                name = os.path.basename(filepath)
                
                # Check if already indexed with same mtime
                latest = db.get_latest_version_by_path(filepath)
                
                if latest and latest[0] == mtime_iso:
                    # Already indexed, skip
                    done += 1
                    continue
                
                # Calculate hash
                content_hash = sha256_file(filepath)
                if not content_hash:
                    done += 1
                    continue
                
                # Add to database
                file_id = db.upsert_file(content_hash, size)
                
                # Add version
                existing_versions = []  # Could query for version index
                v_idx = len(existing_versions) + 1
                
                db.add_version(file_id, name, filepath, mtime_iso, ctime_iso, v_idx, source_folder)
                
                # Add tags from path
                if self.cfg.get("structure_handling", {}).get("tags_from_path"):
                    rel_path = os.path.relpath(filepath, source_folder)
                    folders = os.path.dirname(rel_path).split(os.sep)
                    for folder in folders:
                        if folder and folder != ".":
                            db.add_tag(file_id, folder)
                
                done += 1
                pct = int(done * 100 / max(1, total))
                self.signals.progress.emit(pct, "indexing")
                
                if done % 50 == 0:
                    self.signals.status.emit(f"Indiziert: {done}/{total} Dateien")
            
            except Exception as e:
                # Log error but continue
                try:
                    if 'file_id' in locals():
                        db.add_event(file_id, "error", str(e))
                except:
                    pass
                done += 1
                continue
        
        self.signals.status.emit(f"Fertig! {done} Dateien indiziert.")
        self.signals.progress.emit(100, "indexing")
    
    def stop(self):
        """Stop the worker"""
        self.is_killed = True
    
    def toggle_pause(self):
        """Toggle pause state"""
        self.is_paused = not self.is_paused

# ============================================================================
# 8. ENHANCED SEARCH WIDGET MIT PDF-ICONS UND ERWEITERTEN MENS
# ============================================================================

class SearchWidgetHybrid(QWidget):
    """Haupt-Suchwidget mit PDF-Features"""
    
    def __init__(self, manager, settings):
        super().__init__()
        self.manager = manager
        self.settings = settings
        self.current_results = []
        self.icon_provider = QFileIconProvider()
        
        # Timer für Auto-Search
        self.search_timer = QTimer()
        self.search_timer.setSingleShot(True)
        self.search_timer.timeout.connect(self.perform_search)
        
        self.init_ui()
    
    def init_ui(self):
        layout = QHBoxLayout(self)
        layout.setContentsMargins(0, 0, 0, 0)
        
        # LEFT SIDEBAR (220px)
        left_panel = self.create_left_panel()
        layout.addWidget(left_panel)
        
        # CENTER PANEL
        center_panel = self.create_center_panel()
        layout.addWidget(center_panel, stretch=3)
        
        # RIGHT PANEL
        right_panel = self.create_right_panel()
        layout.addWidget(right_panel, stretch=2)
    
    def create_left_panel(self):
        """Linke Sidebar mit Filtern"""
        widget = QWidget()
        widget.setMaximumWidth(220)
        layout = QVBoxLayout(widget)
        
        lbl = QLabel("Filter")
        lbl.setStyleSheet("font-weight: bold; font-size: 13px; padding: 5px;")
        layout.addWidget(lbl)
        
        # Favoriten
        self.cb_favorites = QCheckBox(" Nur Favoriten")
        self.cb_favorites.toggled.connect(self.on_filter_changed)
        layout.addWidget(self.cb_favorites)
        
        # Dateitypen
        type_group = QGroupBox("Dateitypen")
        type_layout = QVBoxLayout()
        
        self.type_checks = {}
        categories = ["Dokumente", "Bilder", "Audio", "Video", "Archive", "Code", "Tabellen"]
        
        for cat in categories:
            cb = QCheckBox(cat)
            cb.toggled.connect(self.on_filter_changed)
            self.type_checks[cat] = cb
            type_layout.addWidget(cb)
        
        type_group.setLayout(type_layout)
        layout.addWidget(type_group)
        
        # Collections
        coll_label = QLabel(" Sammlungen")
        coll_label.setStyleSheet("font-weight: bold; margin-top: 10px;")
        layout.addWidget(coll_label)
        
        self.collection_list = QListWidget()
        self.collection_list.setMaximumHeight(150)
        self.collection_list.itemClicked.connect(self.on_collection_selected)
        self.collection_list.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.collection_list.customContextMenuRequested.connect(self.show_collection_menu)
        layout.addWidget(self.collection_list)
        
        btn_new_coll = QPushButton("Neue Sammlung")
        btn_new_coll.clicked.connect(self.create_new_collection)
        layout.addWidget(btn_new_coll)
        
        layout.addStretch()
        
        self.load_collections()
        return widget
    
    def create_center_panel(self):
        """Zentrale Ergebnisliste"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        layout.setContentsMargins(5, 5, 5, 5)
        
        # Search Input
        search_layout = QHBoxLayout()
        
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("Suche... (Auto-Suche nach 300ms)")
        self.search_input.textChanged.connect(self.on_search_text_changed)
        search_layout.addWidget(self.search_input)
        
        layout.addLayout(search_layout)
        
        # Toolbar
        toolbar = QHBoxLayout()
        
        self.cb_show_deleted = QCheckBox("Gelschte anzeigen")
        self.cb_show_deleted.toggled.connect(self.perform_search)
        toolbar.addWidget(self.cb_show_deleted)
        
        self.cb_show_hidden = QCheckBox(" Versteckte anzeigen")
        self.cb_show_hidden.setToolTip("Safety-Mode: Zeigt ausgeblendete Dateien an")
        self.cb_show_hidden.toggled.connect(self.perform_search)
        toolbar.addWidget(self.cb_show_hidden)
        
        
        # "Neues Dokument" Button (NEU V13.3!)
        btn_new_doc = QPushButton(" Neues Dokument ")
        btn_new_doc_menu = QMenu(btn_new_doc)
        btn_new_doc_menu.addAction(" Materialverweis", self.create_material_reference)
        btn_new_doc_menu.addAction(" Prompt-Datei", self.create_prompt_file)
        btn_new_doc_menu.addAction(" Internetressource", self.create_internet_resource)
        btn_new_doc_menu.addAction(" Literaturverweis", self.create_literature_reference)
        btn_new_doc.setMenu(btn_new_doc_menu)
        toolbar.addWidget(btn_new_doc)
        toolbar.addStretch()
        
        self.status_label = QLabel("Bereit")
        self.status_label.setStyleSheet("color: #888; font-size: 11px;")
        toolbar.addWidget(self.status_label)
        
        layout.addLayout(toolbar)
        
        # Results Tree (NEU: für gruppierte Anzeige)
        self.result_tree = QTreeWidget()
        # Drag & Drop aktivieren (NEU V13.1!)
        self.result_tree.setDragEnabled(True)
        self.result_tree.setAcceptDrops(True)
        self.result_tree.setDragDropMode(QTreeWidget.DragDropMode.DragDrop)
        self.result_tree.setHeaderLabels(["Datei", "Typ", "Größe", "Datum"])
        self.result_tree.setColumnWidth(0, 400)
        self.result_tree.setColumnWidth(1, 100)
        self.result_tree.setColumnWidth(2, 100)
        self.result_tree.setColumnWidth(3, 150)
        self.result_tree.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.result_tree.customContextMenuRequested.connect(self.show_context_menu)
        self.result_tree.itemSelectionChanged.connect(self.on_selection_changed)
        self.result_tree.itemDoubleClicked.connect(self.open_on_double_click)
        self.result_tree.setSelectionMode(QTreeWidget.SelectionMode.ExtendedSelection)
        self.result_tree.setAlternatingRowColors(True)
        self.result_tree.setAnimated(True)
        layout.addWidget(self.result_tree)
        
        return widget
    
    def create_right_panel(self):
        """Rechte Vorschau"""
        widget = QWidget()
        layout = QVBoxLayout(widget)
        
        lbl = QLabel("Vorschau")
        lbl.setStyleSheet("font-weight: bold; padding: 5px;")
        layout.addWidget(lbl)
        
        self.preview_text = QTextEdit()
        self.preview_text.setReadOnly(True)
        layout.addWidget(self.preview_text)
        
        # Buttons
        btn_layout = QHBoxLayout()
        
        self.btn_open = QPushButton(" öffnen")
        self.btn_open.clicked.connect(self.open_selected_file)
        btn_layout.addWidget(self.btn_open)
        
        self.btn_show = QPushButton("Im Explorer")
        self.btn_show.clicked.connect(self.show_in_explorer)
        btn_layout.addWidget(self.btn_show)
        
        layout.addLayout(btn_layout)
        
        return widget
    
    def dragEnterEvent(self, event):
        """Drag Enter Event für Pooling (NEU V13.1!)"""
        if event.mimeData().hasUrls():
            event.acceptProposedAction()
    
    def dropEvent(self, event):
        """Drop Event für Pooling (NEU V13.1!)"""
        if not event.mimeData().hasUrls():
            return
        
        # Hole Quelldateien (die gedropped werden)
        source_files = []
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if os.path.isfile(path):
                source_files.append(path)
        
        if not source_files:
            return
        
        # Hole Zieldatei (Item unter Cursor)
        pos = event.pos()
        item = self.result_tree.itemAt(pos)
        
        if not item:
            QMessageBox.information(
                self,
                "Kein Ziel",
                "Bitte Dateien auf eine Zieldatei ziehen zum Poolen."
            )
            return
        
        # Wenn Gruppen-Header -> nimm erstes Kind
        if item.childCount() > 0:
            item = item.child(0)
        
        result = item.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        target_file = result['path']
        
        # Starte Pooling
        try:
            self.pool_files(source_files, target_file)
        except Exception as e:
            QMessageBox.critical(
                self,
                "Pooling fehlgeschlagen",
                f"Fehler beim Poolen:\n{str(e)}"
            )
    
    def get_pdf_icon_text(self, result):
        """Gibt PDF-Icon-Text zurück"""
        if not result['name'].lower().endswith('.pdf'):
            return ""
        
        if result.get('pdf_encrypted', False):
            return " 🔒"
        elif result.get('pdf_was_encrypted', False):
            return " 🔓"
        elif result.get('pdf_has_text', False):
            return " ✒️"
        else:
            return " 📷"
    
    def on_search_text_changed(self):
        """Startet Auto-Search Timer"""
        self.search_timer.stop()
        self.search_timer.start(300)  # 300ms Verzögerung
    
    def on_filter_changed(self):
        """Filter geändert -> Neue Suche"""
        self.perform_search()
    
    def perform_search(self):
        """Fhrt Suche durch"""
        term = self.search_input.text().strip()
        
        # Aktive Dateitypen
        types = [cat for cat, cb in self.type_checks.items() if cb.isChecked()]
        
        # Collection
        coll_item = self.collection_list.currentItem()
        coll_id = coll_item.data(Qt.ItemDataRole.UserRole) if coll_item else None
        
        params = {
            "term": term,
            "types": types,
            "favorites": self.cb_favorites.isChecked(),
            "collection_id": coll_id,
            "show_deleted": self.cb_show_deleted.isChecked(),
            "show_hidden": self.cb_show_hidden.isChecked()
        }
        
        self.status_label.setText("Suche läuft...")
        
        self.worker = SearchWorker(self.manager, params, self.settings)
        self.worker.results_found.connect(self.display_results)
        self.worker.finished.connect(lambda: self.status_label.setText("Fertig"))
        self.worker.start()
    
    def display_results(self, results):
        """Zeigt Suchergebnisse mit Gruppierung"""
        self.current_results = results
        self.result_tree.clear()
        
        if not results:
            self.status_label.setText("Keine Ergebnisse")
            return
        
        # Gruppierung nach content_hash
        groups = {}
        for r in results:
            h = r.get('content_hash', '')
            if h and not h.startswith('CLOUD:'):
                if h not in groups:
                    groups[h] = []
                groups[h].append(r)
            else:
                # Einzeldatei ohne Hash (Cloud-Placeholder)
                if 'singles' not in groups:
                    groups['singles'] = []
                groups['singles'].append(r)
        
        # Sortiere Gruppen: Erst groe Gruppen (Duplikate), dann kleine
        sorted_groups = sorted(groups.items(), key=lambda x: len(x[1]), reverse=True)
        
        total_files = 0
        duplicate_groups = 0
        
        for hash_key, files in sorted_groups:
            total_files += len(files)
            
            if len(files) > 1:
                # GRUPPE (Duplikate/Versionen)
                duplicate_groups += 1
                
                # Sortiere innerhalb Gruppe nach Version-Index
                files_sorted = sorted(files, key=lambda f: f.get('version_index', 1), reverse=True)
                
                # Gruppen-Header erstellen
                first_file = files_sorted[0]
                group_name = first_file['name']
                
                # Total size berechnen
                total_size = sum(f.get('size', 0) for f in files_sorted)
                size_mb = total_size / (1024 * 1024)
                
                # Gruppen-Item
                group_item = QTreeWidgetItem(self.result_tree)
                group_item.setText(0, f"{group_name} ({len(files)} Versionen)")
                group_item.setText(1, first_file['category'])
                group_item.setText(2, f"{size_mb:.1f} MB")
                group_item.setText(3, "")
                
                # Gruppen-Farbe
                group_item.setBackground(0, QColor("#2a4a5a"))
                group_item.setForeground(0, QColor("#ffffff"))
                group_item.setExpanded(False)  # Collapsed by default
                
                # Font bold für Gruppe
                font = group_item.font(0)
                font.setBold(True)
                group_item.setFont(0, font)
                
                # Kinder hinzufügen
                for idx, f in enumerate(files_sorted):
                    child = self.create_tree_item(f, is_child=True, child_index=idx)
                    group_item.addChild(child)
                
            else:
                # EINZELDATEI (kein Duplikat)
                f = files[0]
                item = self.create_tree_item(f, is_child=False)
                self.result_tree.addTopLevelItem(item)
        
        self.status_label.setText(f"{total_files} Dateien | {duplicate_groups} Gruppen")
    
    def create_tree_item(self, result, is_child=False, child_index=0):
        """Erstellt ein Tree-Item für eine Datei"""
        # Icon + PDF-Status
        icon = self.icon_provider.icon(QFileInfo(result['path']))
        pdf_icon = self.get_pdf_icon_text(result)
        
        # Version Label
        label = f"[{result['version_label']}] " if result['version_label'] else ""
        if not label and result.get('version_index', 1) > 1:
            label = f"[V{result['version_index']}] "
        
        # Prefix für Kinder
        prefix = ""
        if is_child:
            if child_index == 0:
                prefix = "    "  # Neueste Version
            else:
                prefix = "   "
        
        # Filename
        # Verwende display_name falls vorhanden (Umbenennen-Feature)
        file_name = result.get('display_name', result['name'])
        display_name = f"{prefix}{label}{file_name}{pdf_icon}"
        
        # Item erstellen
        item = QTreeWidgetItem()
        item.setText(0, display_name)
        item.setText(1, result['category'])
        
        # Size
        size_mb = result.get('size', 0) / (1024 * 1024)
        item.setText(2, f"{size_mb:.1f} MB")
        
        # Date
        try:
            date_str = result['mtime'][:10] if result.get('mtime') else ""
            item.setText(3, date_str)
        except:
            item.setText(3, "")
        
        # Icon
        item.setIcon(0, icon)
        
        # Store data
        item.setData(0, Qt.ItemDataRole.UserRole, result)
        item.setToolTip(0, result['path'])
        
        # Color coding
        if result.get('is_deleted', False):
            item.setForeground(0, QColor("#888888"))
            deleted_date = result.get('deleted_at', '')[:10]
            item.setToolTip(0, f"Gelscht am {deleted_date}\n{result['path']}")
        elif result.get('is_favorite', False):
            item.setForeground(0, QColor("#FFD700"))
        
        # Einrckung für Kinder
        if is_child:
            if child_index == 0:
                # Neueste Version - grner Hintergrund
                item.setBackground(0, QColor("#1a3a2a"))
            else:
                # ltere Versionen - leicht ausgegraut
                item.setForeground(0, QColor("#aaaaaa"))
        
        return item
    
    def on_selection_changed(self):
        """Vorschau aktualisieren"""
        items = self.result_tree.selectedItems()
        if not items:
            self.preview_text.clear()
            return
        
        item = items[0]
        
        # Wenn Gruppen-Header ausgewählt -> erstes Kind nehmen
        if item.childCount() > 0:
            item = item.child(0)
        
        result = item.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        
        preview = f"<b>{result['name']}</b><br>"
        preview += f"<i>{path}</i><br><br>"
        
        # Version Info
        if result.get('version_label'):
            preview += f"<b>Version:</b> {result['version_label']}<br>"
        if result.get('version_index', 1) > 1:
            preview += f"<b>Version-Index:</b> {result['version_index']}<br>"
        
        # Size
        size_mb = result.get('size', 0) / (1024 * 1024)
        preview += f"<b>Größe:</b> {size_mb:.2f} MB<br>"
        
        preview += "<br>"
        
        if result.get('is_deleted'):
            preview += "<b>Diese Datei wurde gelöscht</b><br><br>"
        
        # PDF-Status
        if path.lower().endswith('.pdf'):
            preview += "<b>PDF-Status:</b><br>"
            if result.get('pdf_encrypted'):
                preview += " Verschlüsselt<br>"
            if result.get('pdf_has_text'):
                preview += "Enthält Text<br>"
            else:
                preview += "Nur Bilder<br>"
            if result.get('pdf_was_encrypted'):
                preview += " War verschlüsselt<br>"
            preview += "<br>"
        
        # Versuche Dateivorschau
        try:
            if os.path.exists(path):
                if path.lower().endswith(('.txt', '.md', '.py', '.log')):
                    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read(500)
                        preview += f"<pre>{content}...</pre>"
        except:
            preview += "<i>Vorschau nicht verfgbar</i>"
        
        self.preview_text.setHtml(preview)
    
    def show_context_menu(self, pos):
        """Zeigt Rechtsklick-Menü"""
        item = self.result_tree.itemAt(pos)
        
        # NEU V14.1: Batch-Menü bei Mehrfachauswahl
        selected_items = self.result_tree.selectedItems()
        if len(selected_items) > 1:
            self.show_batch_context_menu(pos, selected_items)
            return
        
        # Kein Item -> Spawn-Menü anzeigen (NEU V13.2!)
        if not item:
            self.show_spawn_menu(pos)
            return
        
        # Wenn Gruppen-Header -> Gruppen-Menü
        if item.childCount() > 0:
            self.show_group_context_menu(pos, item)
            return
        
        result = item.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        is_deleted = result.get('is_deleted', False)
        is_hidden = result.get('is_hidden', False)
        is_pdf = path.lower().endswith('.pdf')
        
        menu = QMenu(self)
        
        if is_hidden:
            # Versteckte Datei - nur Einblenden anbieten
            menu.addAction(" Einblenden (Safety-Mode)", self.unhide_selected)
            menu.addAction(" Permanent löschen", self.hard_delete_selected)
        
        elif not is_deleted:
            # Standard-Aktionen
            menu.addAction(" öffnen", self.open_selected_file)
            menu.addAction(" Im Explorer anzeigen", self.show_in_explorer)
            
            menu.addSeparator()
            
            # Kopieren/Einfgen/Umbenennen (NEU V13.2!)
            menu.addAction(" Kopieren", self.copy_selected)
            menu.addAction(" Umbenennen...", self.rename_selected)
            
            # Text kopieren für Textdateien (NEU V13.3!)
            if self.is_text_file(path):
                menu.addAction(" Text kopieren", self.copy_text_from_file)
            
            menu.addSeparator()
            
            # Favoriten
            if result.get('is_favorite'):
                menu.addAction(" Favorit entfernen", self.toggle_favorite)
            else:
                menu.addAction(" Als Favorit markieren", self.toggle_favorite)
            
            menu.addSeparator()
            
            # Version
            version_menu = menu.addMenu(" Version")
            version_menu.addAction(" Erhöhen (+)", self.version_increase)
            version_menu.addAction(" Verringern (-)", self.version_decrease)
            version_menu.addSeparator()
            version_menu.addAction(" Label setzen...", self.set_version_label)
            
            menu.addSeparator()
            
            # Collections
            coll_menu = menu.addMenu(" Zu Sammlung")
            collections = self.get_all_collections()
            for coll_id, coll_name in collections:
                coll_menu.addAction(coll_name, lambda cid=coll_id: self.add_to_collection(cid))
            
            menu.addSeparator()
            
            # PDF-spezifisch
            if is_pdf:
                pdf_menu = menu.addMenu(" PDF")
                pdf_menu.addAction(" Auszug erstellen...", self.create_pdf_excerpt)
                pdf_menu.addAction(" Verschlüsseln...", self.encrypt_pdf)
                pdf_menu.addAction(" Entschlüsseln", self.decrypt_pdf)
            
            # Python-spezifisch
            if path.lower().endswith('.py'):
                py_menu = menu.addMenu(" Python-Tools")
                py_menu.addAction(" In PythonBox öffnen", self.open_in_pythonbox)
                py_menu.addAction(" In Klassen zerlegen", self.split_into_classes)
                py_menu.addAction(" Encoding reparieren", self.fix_encoding)
            
            # SQLite-spezifisch
            if path.lower().endswith(('.db', '.sqlite', '.sqlite3')):
                menu.addSeparator()
                menu.addAction(" In SQLite Viewer öffnen", self.open_sqlite_viewer)
            
            menu.addSeparator()
            
            # Loeschen
            delete_mode = self.settings.get("delete_mode", "soft")
            if delete_mode == "soft":
                menu.addAction(" Loeschen (Papierkorb)", self.delete_selected)
            elif delete_mode == "hard":
                menu.addAction(" Loeschen (Permanent)", self.delete_selected)
            else:  # safety
                menu.addAction(" Ausblenden (Safety-Mode)", self.delete_selected)
            
            if is_deleted:
                menu.addAction(" Permanent löschen", self.delete_selected)
        
        else:
            # Deleted file actions
            menu.addAction(" Wiederherstellen", self.restore_selected)
            menu.addAction(" Permanent löschen", self.hard_delete_selected)
        
        menu.exec(self.result_tree.mapToGlobal(pos))
    
    def show_spawn_menu(self, pos):
        """Zeigt Spawn-Menü für Zwischenablage (NEU V13.2!)"""
        clipboard = QApplication.clipboard()
        text = clipboard.text()
        
        menu = QMenu(self)
        menu.setTitle(" Aus Zwischenablage erstellen")
        
        if not text:
            menu.addAction(" Zwischenablage leer").setEnabled(False)
        else:
            # Zeige Vorschau (erste 50 Zeichen)
            preview = text[:50] + ("..." if len(text) > 50 else "")
            menu.addAction(f" Vorschau: {preview}").setEnabled(False)
            menu.addSeparator()
            
            menu.addAction(" Als PDF spawnen", lambda: self.spawn_from_clipboard("pdf"))
            menu.addAction(" Als TXT spawnen", lambda: self.spawn_from_clipboard("txt"))
            menu.addAction(" Als Word spawnen", lambda: self.spawn_from_clipboard("docx"))
            menu.addAction(" Als RTF spawnen", lambda: self.spawn_from_clipboard("rtf"))
            menu.addAction(" Als ODT spawnen", lambda: self.spawn_from_clipboard("odt"))
        
        menu.exec(self.result_tree.mapToGlobal(pos))
    
    def show_batch_context_menu(self, pos, selected_items):
        """Zeigt Batch-Kontextmenü für Mehrfachauswahl (NEU V14.1)"""
        # Sammle Pfade und zähle Typen
        paths = []
        pdf_count = 0
        py_count = 0
        
        for item in selected_items:
            # Überspringe Gruppen-Header
            if item.childCount() > 0:
                continue
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if result and 'path' in result:
                path = result['path']
                paths.append(path)
                if path.lower().endswith('.pdf'):
                    pdf_count += 1
                elif path.lower().endswith('.py'):
                    py_count += 1
        
        if not paths:
            return
        
        menu = QMenu(self)
        
        # Header mit Anzahl
        header = menu.addAction(f" {len(paths)} Dateien ausgewählt")
        header.setEnabled(False)
        menu.addSeparator()
        
        # Batch-Operationen Untermenü
        batch_menu = menu.addMenu(" Batch-Operationen")
        batch_menu.addAction(" In Ordner kopieren...", lambda: self.batch_copy_to_folder(paths))
        batch_menu.addAction(" Auswahl exportieren (CSV)", lambda: self.batch_export_csv(paths))
        
        # PDF-Operationen wenn mindestens eine PDF dabei
        if pdf_count > 0:
            batch_menu.addSeparator()
            batch_menu.addAction(f" PDF verschlüsseln... ({pdf_count})", 
                               lambda: self.batch_encrypt_pdf(paths))
            batch_menu.addAction(f" PDF entschlüsseln ({pdf_count})", 
                               lambda: self.batch_decrypt_pdf(paths))
            batch_menu.addAction(f" Text extrahieren (OCR) ({pdf_count})", 
                               lambda: self.batch_extract_text(paths))
        
        # Python-Operationen
        if py_count > 0:
            batch_menu.addSeparator()
            batch_menu.addAction(f" Encoding reparieren ({py_count})",
                               lambda: self.batch_fix_encoding(paths))
        
        menu.addSeparator()
        
        # Auswahl-Aktionen
        menu.addAction(" Auswahl aufheben", self.clear_selection)
        menu.addAction(" Alle auswählen", self.select_all_items)
        
        menu.exec(self.result_tree.mapToGlobal(pos))
    
    def batch_copy_to_folder(self, paths):
        """Kopiert mehrere Dateien in einen Ordner (NEU V14.1)"""
        folder = QFileDialog.getExistingDirectory(
            self, "Zielordner wählen", str(Path.home())
        )
        if not folder:
            return
        
        copied = 0
        errors = []
        for path in paths:
            try:
                src = Path(path)
                dst = Path(folder) / src.name
                # Bei Duplikat: _1, _2 etc. anhängen
                if dst.exists():
                    base = dst.stem
                    ext = dst.suffix
                    counter = 1
                    while dst.exists():
                        dst = Path(folder) / f"{base}_{counter}{ext}"
                        counter += 1
                shutil.copy2(src, dst)
                copied += 1
            except Exception as e:
                errors.append(f"{Path(path).name}: {e}")
        
        msg = f"{copied} von {len(paths)} Dateien kopiert."
        if errors:
            msg += f"\n\nFehler ({len(errors)}):\n" + "\n".join(errors[:5])
            if len(errors) > 5:
                msg += f"\n... und {len(errors)-5} weitere"
        QMessageBox.information(self, "Batch-Kopieren", msg)
    
    def batch_export_csv(self, paths):
        """Exportiert Auswahl als CSV-Liste (NEU V14.1)"""
        csv_path, _ = QFileDialog.getSaveFileName(
            self, "CSV speichern", str(Path.home() / "auswahl.csv"), "CSV (*.csv)"
        )
        if not csv_path:
            return
        
        try:
            with open(csv_path, 'w', encoding='utf-8', newline='') as f:
                import csv
                writer = csv.writer(f, delimiter=';')
                writer.writerow(['Dateiname', 'Pfad', 'Größe (KB)', 'Geändert'])
                for path in paths:
                    p = Path(path)
                    size = p.stat().st_size / 1024 if p.exists() else 0
                    mtime = datetime.fromtimestamp(p.stat().st_mtime).strftime('%Y-%m-%d %H:%M') if p.exists() else ''
                    writer.writerow([p.name, str(p), f"{size:.1f}", mtime])
            QMessageBox.information(self, "Export", f"CSV exportiert:\n{csv_path}")
        except Exception as e:
            QMessageBox.warning(self, "Fehler", f"Export fehlgeschlagen:\n{e}")
    
    def batch_encrypt_pdf(self, paths):
        """Verschlüsselt mehrere PDFs (NEU V14.2) - Mit BatchDialog"""
        pdf_paths = [p for p in paths if p.lower().endswith('.pdf')]
        if not pdf_paths:
            QMessageBox.warning(self, "Keine PDFs", "Keine PDF-Dateien in der Auswahl!")
            return
        dialog = BatchDialog(pdf_paths, "pdf_encrypt", "PDF Verschlüsselung", self)
        dialog.exec()
    
    def batch_decrypt_pdf(self, paths):
        """Entschlüsselt mehrere PDFs (NEU V14.2) - Mit BatchDialog"""
        pdf_paths = [p for p in paths if p.lower().endswith('.pdf')]
        if not pdf_paths:
            QMessageBox.warning(self, "Keine PDFs", "Keine PDF-Dateien in der Auswahl!")
            return
        dialog = BatchDialog(pdf_paths, "pdf_decrypt", "PDF Entschlüsselung", self)
        dialog.exec()
    
    def batch_extract_text(self, paths):
        """Extrahiert Text aus mehreren PDFs (NEU V14.2) - Mit BatchDialog"""
        pdf_paths = [p for p in paths if p.lower().endswith('.pdf')]
        if not pdf_paths:
            QMessageBox.warning(self, "Keine PDFs", "Keine PDF-Dateien in der Auswahl!")
            return
        dialog = BatchDialog(pdf_paths, "pdf_extract_text", "Text-Extraktion", self)
        dialog.exec()
    
    def batch_fix_encoding(self, paths):
        """Repariert Encoding in mehreren Python-Dateien (NEU V14.1) - Platzhalter"""
        py_paths = [p for p in paths if p.lower().endswith('.py')]
        QMessageBox.information(
            self, "Batch Encoding-Reparatur",
            f"{len(py_paths)} Python-Dateien ausgewählt.\n\n"
            "Diese Funktion wird später implementiert."
        )
    
    def clear_selection(self):
        """Hebt Auswahl auf (NEU V14.1)"""
        self.result_tree.clearSelection()
    
    def select_all_items(self):
        """Wählt alle Items aus (NEU V14.1)"""
        self.result_tree.selectAll()
    
    def spawn_from_clipboard(self, format):
        """Erstellt Datei aus Zwischenablage (Repariert & Stabil)"""
        clipboard = QApplication.clipboard()
        text = clipboard.text()
        
        if not text:
            QMessageBox.warning(self, "Leer", "Zwischenablage ist leer!")
            return
        
        # Frage nach Dateinamen
        default_name = f"clipboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        filename, ok = QInputDialog.getText(
            self,
            "Dateiname",
            f"Name für neue .{format} Datei:",
            text=default_name
        )
        
        if not ok or not filename:
            return
        
        # Entferne Dateiendung falls vorhanden (um Dopplung zu vermeiden)
        if filename.lower().endswith(f".{format}"):
            filename = filename[:-len(format)-1]
        
        # --- ZIELORDNER ERMITTELN (FIX FÜR ABSTURZ) ---
        target_dir = None
        
        # 1. Versuche, den Ordner der aktuell ausgewählten Datei zu nutzen
        current_item = self.result_tree.currentItem()
        if current_item:
            data = current_item.data(0, Qt.ItemDataRole.UserRole)
            if data and 'path' in data:
                target_dir = os.path.dirname(data['path'])
        
        # 2. Fallback: Wenn nichts ausgewählt ist, frage den Benutzer
        if not target_dir or not os.path.exists(target_dir):
            target_dir = QFileDialog.getExistingDirectory(
                self, 
                "Zielordner für neue Datei auswählen",
                ""
            )
        
        if not target_dir:
            return # Abbruch durch Benutzer
            
        # Vollständiger Pfad
        output_path = os.path.join(target_dir, f"{filename}.{format}")
        
        # Prüfe ob Datei existiert
        if os.path.exists(output_path):
            reply = QMessageBox.question(
                self,
                "Datei existiert",
                f"Datei existiert bereits:\n{output_path}\n\nÜberschreiben?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.No:
                return
        
        try:
            # Erstelle Datei je nach Format
            if format == "txt":
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
            
            elif format == "pdf":
                # Einfaches PDF mit Text (benötigt reportlab)
                try:
                    from reportlab.lib.pagesizes import letter
                    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                    from reportlab.lib.styles import getSampleStyleSheet
                    
                    doc = SimpleDocTemplate(output_path, pagesize=letter)
                    styles = getSampleStyleSheet()
                    story = []
                    
                    # XML-Zeichen escapen (wichtig für reportlab!)
                    safe_text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    
                    # Text in Paragraphen aufteilen
                    for para in safe_text.split('\n\n'):
                        if para.strip():
                            # Zeilenumbrüche innerhalb von Absätzen erhalten
                            formatted_para = para.replace('\n', '<br/>')
                            story.append(Paragraph(formatted_para, styles['Normal']))
                            story.append(Spacer(1, 12))
                    
                    doc.build(story)
                except ImportError:
                    QMessageBox.warning(
                        self,
                        "Modul fehlt",
                        "reportlab nicht installiert.\nVerwende TXT stattdessen."
                    )
                    # Fallback auf TXT
                    output_path = output_path.replace('.pdf', '.txt')
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
            
            elif format == "docx":
                # Word-Dokument
                try:
                    from docx import Document
                    doc = Document()
                    # Text zeilenweise hinzufügen
                    for para in text.split('\n'):
                        doc.add_paragraph(para)
                    doc.save(output_path)
                except ImportError:
                    QMessageBox.warning(
                        self,
                        "Modul fehlt",
                        "python-docx nicht installiert.\nVerwende TXT stattdessen."
                    )
                    output_path = output_path.replace('.docx', '.txt')
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
            
            elif format == "rtf":
                # Einfaches RTF
                rtf_content = "{\\rtf1\\ansi\\deff0\n"
                for line in text.split('\n'):
                    # RTF-Syntax escapen
                    safe_line = line.replace('\\', '\\\\').replace('{', '\\{').replace('}', '\\}')
                    rtf_content += safe_line + "\\par\n"
                rtf_content += "}"
                
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(rtf_content)
            
            elif format == "odt":
                # ODT benötigt odfpy
                try:
                    from odf.opendocument import OpenDocumentText
                    from odf.text import P
                    
                    textdoc = OpenDocumentText()
                    for para in text.split('\n'):
                        p = P(text=para)
                        textdoc.text.addElement(p)
                    textdoc.save(output_path)
                except ImportError:
                    QMessageBox.warning(
                        self,
                        "Modul fehlt",
                        "odfpy nicht installiert.\nVerwende TXT stattdessen."
                    )
                    output_path = output_path.replace('.odt', '.txt')
                    with open(output_path, 'w', encoding='utf-8') as f:
                        f.write(text)
            
            QMessageBox.information(
                self,
                "Erfolgreich",
                f"✅ Datei erstellt:\n{os.path.basename(output_path)}"
            )
            
            # Aktualisiere Ansicht (Re-Indexierung/Suche)
            self.perform_search()
            
        except Exception as e:
            import traceback
            traceback.print_exc()
            QMessageBox.critical(
                self,
                "Fehler",
                f"Fehler beim Erstellen der Datei:\n{str(e)}"
            )
            
    def show_group_context_menu(self, pos, group_item):
        """Zeigt Kontext-Menfr Gruppen"""
        menu = QMenu(self)
        
        # Expand/Collapse
        if group_item.isExpanded():
            menu.addAction("Zuklappen", lambda: group_item.setExpanded(False))
        else:
            menu.addAction(" Aufklappen", lambda: group_item.setExpanded(True))
        
        menu.addSeparator()
        
        # Alle Kinder auswählen
        def select_all_children():
            self.result_tree.clearSelection()
            for i in range(group_item.childCount()):
                child = group_item.child(i)
                child.setSelected(True)
        
        menu.addAction(" Alle Versionen auswählen", select_all_children)
        
        menu.addSeparator()
        
        # Toggle für alle Gruppen (NEU V13!)
        menu.addAction("Alle Gruppen aufklappen", self.expand_all_groups)
        menu.addAction("Alle Gruppen zuklappen", self.collapse_all_groups)
        
        menu.addSeparator()
        
        # Alte Versionen löschen (alle außer neueste)
        def delete_old_versions():
            reply = QMessageBox.question(
                self,
                "Alte Versionen löschen",
                f"Alle Versionen außer der neuesten löschen?\n({group_item.childCount() - 1} Dateien)",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            
            if reply == QMessageBox.StandardButton.Yes:
                delete_mode = self.settings.get("delete_mode", "soft")
                
                # überspringe erstes Kind (neueste Version)
                for i in range(1, group_item.childCount()):
                    child = group_item.child(i)
                    result = child.data(0, Qt.ItemDataRole.UserRole)
                    if result:
                        vid = result['id']
                        db_path = result['db']
                        
                        db = ConnectionDB(db_path)
                        if delete_mode == "soft":
                            db.soft_delete_version(vid)
                        else:
                            db.hard_delete_version(vid)
                        db.close()
                
                self.perform_search()
        
        menu.addAction("Alte Versionen löschen", delete_old_versions)
        
        menu.exec(self.result_tree.mapToGlobal(pos))

    def expand_all_groups(self):
        """Klappt alle Gruppen auf (NEU V13!)"""
        root = self.result_tree.invisibleRootItem()
        for i in range(root.childCount()):
            item = root.child(i)
            if item.childCount() > 0:
                item.setExpanded(True)
    
    def collapse_all_groups(self):
        """Klappt alle Gruppen zu (NEU V13!)"""
        root = self.result_tree.invisibleRootItem()
        for i in range(root.childCount()):
            item = root.child(i)
            if item.childCount() > 0:
                item.setExpanded(False)
    
    def get_selected_results(self):
        """Hilfsfunktion: Holt ausgewählte Dateien (keine Gruppen-Header)"""
        items = self.result_tree.selectedItems()
        results = []
        
        for item in items:
            # überspringe Gruppen-Header
            if item.childCount() > 0:
                continue
            
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if result:
                results.append(result)
        
        return results
    
    # ========== PDF ACTIONS ==========
    
    def encrypt_pdf(self):
        """Verschlüsselt ausgewählte PDFs"""
        results = self.get_selected_results()
        if not results:
            return
        
        paths = [r['path'] for r in results]
        
        # Dialog
        dialog = PDFPasswordDialog(paths, mode="encrypt", settings=self.settings, parent=self)
        if dialog.exec() != QDialog.DialogCode.Accepted:
            return
        
        password = dialog.password
        
        # Progress
        progress = QProgressBar(self)
        progress.setMaximum(len(paths))
        self.statusBar().addWidget(progress)
        
        success_count = 0
        
        for idx, path in enumerate(paths):
            try:
                # Output path
                base = os.path.splitext(path)[0]
                output = f"{base}_encrypted.pdf"
                
                # Encrypt
                PDFUtils.encrypt_pdf(path, output, password)
                success_count += 1
                
            except Exception as e:
                print(f"Encryption error {path}: {e}")
            
            progress.setValue(idx + 1)
        
        self.statusBar().removeWidget(progress)
        
        QMessageBox.information(
            self,
            "Verschlüsselung abgeschlossen",
            f" {success_count} von {len(paths)} Dateien verschlüsselt"
        )
        
        self.perform_search()
    
    def decrypt_pdf(self):
        """Entschlüsselt ausgewählte PDFs"""
        results = self.get_selected_results()
        if not results:
            return
        
        paths = [r['path'] for r in results]
        
        # Dialog
        dialog = PDFPasswordDialog(paths, mode="decrypt", settings=self.settings, parent=self)
        if dialog.exec() != QDialog.DialogCode.Accepted:
            return
        
        password = dialog.password
        
        success_count = 0
        
        for path in paths:
            try:
                base = os.path.splitext(path)[0]
                output = f"{base}_decrypted.pdf"
                
                PDFUtils.decrypt_pdf(path, output, password)
                success_count += 1
                
            except Exception as e:
                QMessageBox.warning(self, "Fehler", f"Entschlüsselung fehlgeschlagen:\n{str(e)}")
                break
        
        QMessageBox.information(
            self,
            "Entschlüsselung abgeschlossen",
            f" {success_count} von {len(paths)} Dateien entschlüsselt"
        )
        
        self.perform_search()
    
    def create_pdf_excerpt(self):
        """Erstellt PDF-Auszug"""
        results = self.get_selected_results()
        if not results:
            return
        
        path = results[0]['path']
        
        dialog = PDFExcerptDialog(path, parent=self)
        dialog.exec()
        
        self.perform_search()
    
    def remove_pdf_text(self):
        """Entfernt Text aus PDF"""
        results = self.get_selected_results()
        if not results:
            return
        
        for result in results:
            path = result['path']
            
            try:
                base = os.path.splitext(path)[0]
                output = f"{base}_blank.pdf"
                
                PDFUtils.remove_text_from_pdf(path, output)
                
                QMessageBox.information(self, "Erfolg", f" Text entfernt:\n{output}")
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Fehler:\n{str(e)}")
        
        self.perform_search()
    
    def apply_pdf_ocr(self):
        """Wendet OCR auf PDF an"""
        results = self.get_selected_results()
        if not results:
            return
        
        lang = self.settings.get("ocr_language", "deu")
        
        for result in results:
            path = result['path']
            
            try:
                base = os.path.splitext(path)[0]
                output = f"{base}_ocr.pdf"
                
                PDFUtils.apply_ocr_to_pdf(path, output, lang)
                
                QMessageBox.information(self, "Erfolg", f" OCR angewendet:\n{output}")
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"OCR fehlgeschlagen:\n{str(e)}")
        
        self.perform_search()
    
    # ========== STANDARD ACTIONS ==========
    
    def open_selected_file(self):
        """Öffnet ausgewählte Datei"""
        results = self.get_selected_results()
        if not results:
            return
        
        path = results[0]['path']
        
        try:
            if sys.platform == 'win32':
                os.startfile(path)
            elif sys.platform == 'darwin':
                subprocess.run(['open', path])
            else:
                subprocess.run(['xdg-open', path])
        except Exception as e:
            QMessageBox.warning(self, "Fehler", f"Datei konnte nicht geöffnet werden:\n{str(e)}")
    
    def open_on_double_click(self, item, column):
        """Öffnet Datei bei Doppelklick (NEU V13.1!)"""
        # Wenn Gruppen-Header -> Aufklappen/Zuklappen statt öffnen
        if item.childCount() > 0:
            item.setExpanded(not item.isExpanded())
            return
        
        # Datei öffnen
        result = item.data(0, Qt.ItemDataRole.UserRole)
        if result and 'path' in result:
            path = result['path']
            try:
                if sys.platform == 'win32':
                    os.startfile(path)
                elif sys.platform == 'darwin':
                    subprocess.run(['open', path])
                else:
                    subprocess.run(['xdg-open', path])
            except Exception as e:
                QMessageBox.warning(self, "Fehler", f"Datei konnte nicht geöffnet werden:\n{str(e)}")
    
    def show_in_explorer(self):
        """Zeigt Datei im Explorer"""
        results = self.get_selected_results()
        if not results:
            return
        
        path = results[0]['path']
        
        try:
            if sys.platform == 'win32':
                subprocess.run(['explorer', '/select,', path])
            elif sys.platform == 'darwin':
                subprocess.run(['open', '-R', path])
            else:
                subprocess.run(['xdg-open', os.path.dirname(path)])
        except Exception as e:
            QMessageBox.warning(self, "Fehler", f"Explorer konnte nicht geöffnet werden:\n{str(e)}")
    
    def toggle_favorite(self):
        """Favoriten-Status umschalten"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            is_fav = result.get('is_favorite', False)
            db.set_favorite(vid, not is_fav)
            db.close()
        
        self.perform_search()
    
    def version_increase(self):
        """Version erhhen"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.swap_version_index(vid, +1)
            db.close()
        
        self.perform_search()
    
    def version_decrease(self):
        """Version verringern"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.swap_version_index(vid, -1)
            db.close()
        
        self.perform_search()
    
    def set_version_label(self):
        """Setzt Version-Label"""
        results = self.get_selected_results()
        if not results:
            return
        
        labels = ["", "Original", "V1", "V2", "V3", "V4", "V5", "Draft", "Final", "Review"]
        
        label, ok = QInputDialog.getItem(
            self,
            "Version Label",
            "Label auswählen:",
            labels,
            0,
            False
        )
        
        if ok:
            for result in results:
                vid = result['id']
                db_path = result['db']
                
                db = ConnectionDB(db_path)
                db.set_version_label(vid, label)
                db.close()
            
            self.perform_search()
    
    def add_to_collection(self, coll_id):
        """Fgt zu Sammlung hinzu"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.add_to_collection(coll_id, vid)
            db.close()
    
    def delete_selected(self):
        """Löscht ausgewählte Dateien"""
        results = self.get_selected_results()
        if not results:
            return
        
        mode = self.settings.get("delete_mode", "soft")
        
        # Angepasste Besttigungsmeldung
        if mode == "safety":
            mode_text = "Ausblenden (Safety-Mode - Keine Dateisystem-nderung)"
        elif mode == "soft":
            mode_text = "Soft-Delete (Papierkorb)"
        else:
            mode_text = "Permanent Loeschen"
        
        reply = QMessageBox.question(
            self,
            "Loeschen bestätigen",
            f"{mode_text}: {len(results)} Datei(en)?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            for result in results:
                vid = result['id']
                db_path = result['db']
                
                db = ConnectionDB(db_path)
                if mode == "soft":
                    db.soft_delete_version(vid)
                elif mode == "hard":
                    db.hard_delete_version(vid)
                else:  # safety
                    db.safety_hide_version(vid)
                db.close()
            
            self.perform_search()
    
    def restore_selected(self):
        """Stellt gelöschte Dateien wieder her"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.restore_version(vid)
            db.close()
        
        self.perform_search()
    
    
    def unhide_selected(self):
        """Blendet versteckte Dateien wieder ein (Safety-Mode) - NEU V13.2!"""
        results = self.get_selected_results()
        if not results:
            return
        
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.safety_unhide_version(vid)
            db.close()
        
        QMessageBox.information(
            self,
            "Eingeblendet",
            f"{len(results)} Datei(en) wieder sichtbar gemacht."
        )
        self.perform_search()
        """Stellt gelschte Dateien wieder her"""
        results = self.get_selected_results()
        for result in results:
            vid = result['id']
            db_path = result['db']
            
            db = ConnectionDB(db_path)
            db.restore_version(vid)
            db.close()
        
        self.perform_search()
    def copy_selected(self):
        """Kopiert ausgewaehlte Dateien in System-Zwischenablage (NEU V13.2!)"""
        results = self.get_selected_results()
        if not results:
            return
        
        # Sammle Dateipfade
        paths = [result['path'] for result in results]
        
        # Für Windows: Verwende spezielle API
        if sys.platform == 'win32':
            try:
                import win32clipboard
                import win32con
                
                # Format für Dateipfade
                win32clipboard.OpenClipboard()
                win32clipboard.EmptyClipboard()
                
                # CF_HDROP Format für Dateien
                files_data = '\0'.join(paths) + '\0\0'
                win32clipboard.SetClipboardData(win32con.CF_HDROP, files_data.encode('utf-16-le'))
                win32clipboard.CloseClipboard()
                
                QMessageBox.information(
                    self,
                    "Kopiert",
                    f"{len(paths)} Datei(en) in Zwischenablage kopiert.\nKann jetzt eingefgt werden."
                )
            except ImportError:
                # Fallback: Pfade als Text
                clipboard = QApplication.clipboard()
                clipboard.setText('\n'.join(paths))
                QMessageBox.information(
                    self,
                    "Kopiert (Text)",
                    f"{len(paths)} Dateipfad(e) als Text kopiert."
                )
        else:
            # macOS/Linux: Verwende QMimeData
            clipboard = QApplication.clipboard()
            mime_data = QMimeData()
            
            # Setze URLs
            from PyQt6.QtCore import QUrl
            urls = [QUrl.fromLocalFile(path) for path in paths]
            mime_data.setUrls(urls)
            
            # Auch als Text
            mime_data.setText('\n'.join(paths))
            
            clipboard.setMimeData(mime_data)
            
            QMessageBox.information(
                self,
                "Kopiert",
                f"{len(paths)} Datei(en) in Zwischenablage kopiert.\nKann jetzt eingefgt werden."
            )
            vid = result['id']
            db_path = result['db']
            
    
    def rename_selected(self):
        """Benennt ausgewaehlte Datei um (NEU V13.2!)"""
        results = self.get_selected_results()
        if not results or len(results) != 1:
            QMessageBox.information(
                self,
                "Hinweis",
                "Bitte genau EINE Datei zum Umbenennen auswählen."
            )
            return
        
        result = results[0]
        old_path = result['path']
        old_name = os.path.basename(old_path)
        old_dir = os.path.dirname(old_path)
        
        # Frage nach neuem Namen
        new_name, ok = QInputDialog.getText(
            self,
            "Umbenennen",
            "Neuer Dateiname:",
            text=old_name
        )
        
        if not ok or not new_name or new_name == old_name:
            return
        
        # Prfe Einstellung: Dateisystem oder nur DB?
        rename_filesystem = self.settings.get("rename_in_filesystem", True)
        
        if rename_filesystem:
            # Echte Umbenennung im Dateisystem
            new_path = os.path.join(old_dir, new_name)
            
            # Prfe ob Ziel existiert
            if os.path.exists(new_path):
                QMessageBox.warning(
                    self,
                    "Datei existiert",
                    f"Eine Datei mit dem Namen '{new_name}' existiert bereits."
                )
                return
            
            try:
                os.rename(old_path, new_path)
                
                # Update DB
                vid = result['id']
                db_path = result['db']
                db = ConnectionDB(db_path)
                db.conn.execute(
                    "UPDATE versions SET path=?, name=? WHERE id=?",
                    (new_path, new_name, vid)
                )
                db.conn.commit()
                db.close()
                
                QMessageBox.information(
                    self,
                    "Umbenannt",
                    f"Datei umbenannt:\n{old_name}  {new_name}"
                )
                
                # Re-Indexierung triggern
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(
                    self,
                    "Fehler",
                    f"Umbenennen fehlgeschlagen:\n{str(e)}"
                )
        else:
            # Nur lokale Umbenennung (display_name in DB)
            vid = result['id']
            db_path = result['db']
            
            # Prfe ob display_name Spalte existiert
            db = ConnectionDB(db_path)
            
            # Fge display_name Spalte hinzu falls nicht vorhanden
            try:
                cur = db.conn.cursor()
                cur.execute("PRAGMA table_info(versions)")
                columns = [row[1] for row in cur.fetchall()]
                
                if 'display_name' not in columns:
                    db.conn.execute("ALTER TABLE versions ADD COLUMN display_name TEXT")
                    db.conn.commit()
                
                # Setze display_name
                db.conn.execute(
                    "UPDATE versions SET display_name=? WHERE id=?",
                    (new_name, vid)
                )
                db.conn.commit()
                db.close()
                
                QMessageBox.information(
                    self,
                    "Umbenannt (Lokal)",
                    f"Anzeigename geändert:\n{old_name}  {new_name}\n\n(Datei auf Festplatte unverändert)"
                )
                
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(
                    self,
                    "Fehler",
                    f"Lokale Umbenennung fehlgeschlagen:\n{str(e)}"
                )
                if db:
                    db.close()
            db = ConnectionDB(db_path)
            db.restore_version(vid)
            db.close()
        
        self.perform_search()
    
    def hard_delete_selected(self):
        """Löscht permanent"""
        results = self.get_selected_results()
        if not results:
            return
        
        reply = QMessageBox.warning(
            self,
            "Permanent löschen",
            f" {len(results)} Datei(en) PERMANENT löschen?\nDies kann nicht rückgängig gemacht werden!",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            for result in results:
                vid = result['id']
                db_path = result['db']
                
                db = ConnectionDB(db_path)
                db.hard_delete_version(vid)
                db.close()
            
            self.perform_search()
    
    
    def is_text_file(self, path):
        """Prüft ob Datei eine Textdatei ist (NEU V13.3!)"""
        text_extensions = [
            '.txt', '.md', '.log', '.py', '.js', '.ts', '.jsx', '.tsx',
            '.html', '.htm', '.css', '.json', '.xml', '.yaml', '.yml',
            '.csv', '.tsv', '.ini', '.conf', '.cfg', '.sh', '.bash',
            '.c', '.cpp', '.h', '.hpp', '.java', '.cs', '.php', '.rb',
            '.go', '.rs', '.swift', '.kt', '.sql', '.r', '.m', '.pl'
        ]
        return any(path.lower().endswith(ext) for ext in text_extensions)
    
    def copy_text_from_file(self):
        """Kopiert gesamten Text aus Datei in Zwischenablage (NEU V13.3!)"""
        results = self.get_selected_results()
        if not results or len(results) != 1:
            QMessageBox.information(
                self,
                "Hinweis",
                "Bitte genau EINE Textdatei auswählen."
            )
            return
        
        path = results[0]['path']
        
        if not self.is_text_file(path):
            QMessageBox.warning(
                self,
                "Keine Textdatei",
                "Diese Datei ist keine Textdatei."
            )
            return
        
        try:
            # Versuche verschiedene Encodings
            encodings = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
            text = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(path, 'r', encoding=encoding) as f:
                        text = f.read()
                    used_encoding = encoding
                    break
                except (UnicodeDecodeError, LookupError):
                    continue
            
            if text is None:
                # Letzte Chance: Binary lesen
                with open(path, 'rb') as f:
                    raw = f.read()
                text = raw.decode('utf-8', errors='replace')
                used_encoding = "utf-8 (with errors)"
            
            # Kopiere in Zwischenablage
            clipboard = QApplication.clipboard()
            clipboard.setText(text)
            
            # Zeige Erfolg
            lines = text.count('\n') + 1
            chars = len(text)
            
            QMessageBox.information(
                self,
                "Text kopiert",
                f"Text erfolgreich kopiert:\n\n"
                f" Datei: {os.path.basename(path)}\n"
                f" Zeilen: {lines:,}\n"
                f" Zeichen: {chars:,}\n"
                f" Encoding: {used_encoding}"
            )
            
        except Exception as e:
            QMessageBox.critical(
                self,
                "Fehler",
                f"Fehler beim Lesen der Datei:\n{str(e)}"
            )
        """Löscht permanent"""
        results = self.get_selected_results()
        if not results:
            return
        
        reply = QMessageBox.warning(
            self,
            "Permanent löschen",
            f" {len(results)} Datei(en) PERMANENT löschen?\nDies kann nicht rückgängig gemacht werden!",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            for result in results:
                vid = result['id']
                db_path = result['db']
                
                db = ConnectionDB(db_path)
                db.hard_delete_version(vid)
                db.close()
            
            self.perform_search()
    
    # ========== COLLECTIONS ==========
    
    def load_collections(self):
        """Ldt Collections"""
        self.collection_list.clear()
        
        for db_path in self.manager.dbs:
            if not os.path.exists(db_path):
                continue
            
            db = ConnectionDB(db_path)
            collections = db.get_collections()
            db.close()
            
            for coll_id, name, desc in collections:
                item = QListWidgetItem(name)
                item.setData(Qt.ItemDataRole.UserRole, coll_id)
                item.setToolTip(desc or "Keine Beschreibung")
                self.collection_list.addItem(item)
    
    def get_all_collections(self):
        """Gibt alle Collections zurück"""
        collections = []
        
        for db_path in self.manager.dbs:
            if not os.path.exists(db_path):
                continue
            
            db = ConnectionDB(db_path)
            colls = db.get_collections()
            db.close()
            
            collections.extend(colls)
        
        return [(cid, name) for cid, name, desc in collections]
    
    def create_new_collection(self):
        """Erstellt neue Collection"""
        name, ok = QInputDialog.getText(self, "Neue Sammlung", "Name:")
        
        if ok and name:
            # In erste DB einfgen
            if self.manager.dbs:
                db = ConnectionDB(self.manager.dbs[0])
                success = db.add_collection(name)
                db.close()
                
                if success:
                    self.load_collections()
                else:
                    QMessageBox.warning(self, "Fehler", "Sammlung existiert bereits!")
    
    def show_collection_menu(self, pos):
        """Zeigt Collection-Menü"""
        item = self.collection_list.itemAt(pos)
        if not item:
            return
        
        menu = QMenu(self)
        menu.addAction(" Als PDF exportieren", lambda: self.export_collection_list(item))
        menu.addSeparator()
        menu.addAction(" Loeschen", lambda: self.delete_collection(item))
        menu.exec(self.collection_list.mapToGlobal(pos))
    
    def delete_collection(self, item):
        """Löscht Collection"""
        coll_id = item.data(Qt.ItemDataRole.UserRole)
        
        reply = QMessageBox.question(
            self,
            "Löschen bestätigen",
            f"Sammlung '{item.text()}' löschen?",
            QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
        )
        
        if reply == QMessageBox.StandardButton.Yes:
            for db_path in self.manager.dbs:
                if not os.path.exists(db_path):
                    continue
                
                db = ConnectionDB(db_path)
                db.remove_collection(coll_id)
                db.close()
            
            self.load_collections()
    
    def on_collection_selected(self):
        """Collection wurde ausgewählt"""
        self.perform_search()
    
    # ========================================================================
    # ANONYMIZATION FUNCTIONS
    # ========================================================================
    
    def show_anonymization_settings(self):
        pass  # TODO: Implement anonymization settings dialog
    
    def export_collection_list(self, collection_item):
        """Exportiert Sammlung als PDF (NEU V13.3!)"""
        coll_id = collection_item.data(Qt.ItemDataRole.UserRole)
        coll_name = collection_item.text()
        
        # Sammle Dateien aus allen DBs
        all_files = []
        
        for db_path in self.manager.dbs:
            if not os.path.exists(db_path):
                continue
            
            try:
                db = ConnectionDB(db_path)
                
                # Hole Dateien der Collection
                query = """
                    SELECT v.path, v.name, v.mtime, f.size, f.category
                    FROM file_tags ft
                    JOIN versions v ON ft.version_id = v.id
                    JOIN files f ON v.file_id = f.id
                    WHERE ft.collection_id = ? AND v.is_deleted = 0
                    ORDER BY v.path
                """
                
                conn = sqlite3.connect(db_path)
                conn.row_factory = sqlite3.Row
                rows = conn.execute(query, (coll_id,)).fetchall()
                conn.close()
                
                all_files.extend(rows)
                db.close()
                
            except Exception as e:
                print(f"Fehler beim Lesen von {db_path}: {e}")
                continue
        
        if not all_files:
            QMessageBox.information(
                self,
                "Keine Dateien",
                f"Sammlung '{coll_name}' enthaelt keine Dateien."
            )
            return
        
        # PDF-Pfad wählen
        default_name = f"Sammlung_{coll_name}_{datetime.now().strftime('%Y%m%d')}.pdf"
        safe_name = "".join(c for c in default_name if c.isalnum() or c in ('_', '-', '.'))
        
        pdf_path, _ = QFileDialog.getSaveFileName(
            self,
            "Sammlung exportieren",
            safe_name,
            "PDF-Dateien (*.pdf)"
        )
        
        if not pdf_path:
            return
        
        try:
            if not HAS_REPORTLAB:
                QMessageBox.warning(
                    self,
                    "Bibliothek fehlt",
                    "reportlab ist nicht installiert.\n\nInstalliere mit: pip install reportlab"
                )
                return
            
            # PDF erstellen
            doc = SimpleDocTemplate(pdf_path, pagesize=A4)
            story = []
            styles = getSampleStyleSheet()
            
            # Titel
            title = Paragraph(f"<b>Sammlung: {coll_name}</b>", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 0.5*cm))
            
            # Info
            info = Paragraph(
                f"Erstellt: {datetime.now().strftime('%d.%m.%Y %H:%M')}<br/>"
                f"Dateien: {len(all_files)}",
                styles['Normal']
            )
            story.append(info)
            story.append(Spacer(1, 1*cm))
            
            # Tabelle mit Dateien
            table_data = [['Dateiname', 'Pfad', 'Größe', 'Kategorie']]
            
            for row in all_files:
                # Größe formatieren
                size_bytes = row['size']
                if size_bytes < 1024:
                    size_str = f"{size_bytes} B"
                elif size_bytes < 1024**2:
                    size_str = f"{size_bytes/1024:.1f} KB"
                elif size_bytes < 1024**3:
                    size_str = f"{size_bytes/(1024**2):.1f} MB"
                else:
                    size_str = f"{size_bytes/(1024**3):.2f} GB"
                
                table_data.append([
                    row['name'],
                    row['path'],
                    size_str,
                    row['category'] or '-'
                ])
            
            # Tabelle erstellen
            
            table = Table(table_data, colWidths=[5*cm, 8*cm, 2*cm, 3*cm])
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('FONTSIZE', (0, 1), (-1, -1), 8),
            ]))
            
            story.append(table)
            
            # PDF bauen
            doc.build(story)
            
            QMessageBox.information(
                self,
                "Export erfolgreich",
                f" Sammlung exportiert:\n{pdf_path}\n\n{len(all_files)} Dateien"
            )
        
        except Exception as e:
            QMessageBox.critical(
                self,
                "Fehler",
                f"Export fehlgeschlagen:\n{str(e)}\n\nTraceback:\n{traceback.format_exc()}"
            )
        """Zeigt Anonymisierungs-Einstellungen Dialog"""
        dialog = AnonymizationSettingsDialog(self.settings, self)
        dialog.exec()
    
    def anonymize_file(self):
        """Anonymisiert ausgewählte Datei(en) mit Platzhalter"""
        selected_paths = self.get_selected_paths()
        if not selected_paths:
            QMessageBox.warning(self, "Keine Auswahl", "Bitte wählen Sie mindestens eine Datei aus")
            return
        
        # Prfe ob Blacklist leer ist
        blacklist = self.settings.get("anonymization_blacklist", [])
        if not blacklist:
            reply = QMessageBox.question(
                self,
                "Blacklist leer",
                "Die Blacklist ist leer. Mchten Sie jetzt Begriffe hinzufügen?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.Yes:
                self.show_anonymization_settings()
                return
            else:
                return
        
        # Starte Anonymisierung
        whitelist = self.settings.get("anonymization_whitelist", [])
        placeholder = self.settings.get("anonymization_placeholder", "[-----]")
        
        # Progress Dialog
        progress = QProgressBar()
        progress.setWindowTitle("Anonymisierung")
        progress.setRange(0, len(selected_paths))
        progress.show()
        
        # Worker starten
        self.anon_worker = AnonymizationWorker(
            selected_paths, blacklist, whitelist, placeholder, mode="anonymize"
        )
        self.anon_worker.progress.connect(progress.setValue)
        self.anon_worker.log_message.connect(lambda msg: print(msg))
        self.anon_worker.finished.connect(lambda: [
            progress.close(),
            QMessageBox.information(self, "Fertig", "Anonymisierung abgeschlossen"),
            self.perform_search()  # Refresh
        ])
        self.anon_worker.start()
    
    def redact_pdf_file(self):
        """Schwärzt ausgewählte PDF(s)"""
        selected_paths = self.get_selected_paths()
        if not selected_paths:
            QMessageBox.warning(self, "Keine Auswahl", "Bitte wählen Sie mindestens eine Datei aus")
            return
        
        # Prfe ob Blacklist leer ist
        blacklist = self.settings.get("anonymization_blacklist", [])
        if not blacklist:
            reply = QMessageBox.question(
                self,
                "Blacklist leer",
                "Die Blacklist ist leer. Mchten Sie jetzt Begriffe hinzufügen?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.Yes:
                self.show_anonymization_settings()
                return
            else:
                return
        
        # Prfe ob fitz installiert ist
        if not HAS_FITZ:
            QMessageBox.critical(
                self,
                "Fehler",
                "PyMuPDF (fitz) ist nicht installiert.\n\nInstallieren Sie es mit:\npip install PyMuPDF"
            )
            return
        
        # Starte Schwärzung
        whitelist = self.settings.get("anonymization_whitelist", [])
        
        # Progress Dialog
        progress = QProgressBar()
        progress.setWindowTitle("PDF-Schwärzung")
        progress.setRange(0, len(selected_paths))
        progress.show()
        
        # Worker starten
        self.redact_worker = AnonymizationWorker(
            selected_paths, blacklist, whitelist, mode="redact"
        )
        self.redact_worker.progress.connect(progress.setValue)
        self.redact_worker.log_message.connect(lambda msg: print(msg))
        self.redact_worker.finished.connect(lambda: [
            progress.close(),
            QMessageBox.information(self, "Fertig", "Schwärzung abgeschlossen"),
            self.perform_search()  # Refresh
        ])
        self.redact_worker.start()
    
    def get_selected_paths(self):
        """Gibt Pfade aller ausgewählten Dateien zurück"""
        paths = []
        for item in self.result_tree.selectedItems():
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if result and 'path' in result:
                paths.append(result['path'])
        return paths
    
    # ========================================================================
    # PYTHON-TOOLS FUNCTIONS
    # ========================================================================
    
    def split_python_classes_txt(self):
        """Zerlegt Python-Datei in Klassen (.txt Format)"""
        self._split_python_classes(as_py=False)
    
    def split_python_classes_py(self):
        """Zerlegt Python-Datei in Klassen (.py Format)"""
        self._split_python_classes(as_py=True)
    
    def _split_python_classes(self, as_py=False):
        """Zerlegt Python-Datei in einzelne Klassen"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        
        try:
            import ast
            from datetime import datetime
            
            # Ausgabeordner erstellen
            base = os.path.basename(path)
            name, _ = os.path.splitext(base)
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            folder = os.path.dirname(path)
            outdir = os.path.join(folder, f"pyCutter_{name}_{timestamp}")
            os.makedirs(outdir, exist_ok=True)
            
            # Datei lesen und parsen
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                source = f.read()
            
            tree = ast.parse(source)
            
            # Klassen und Funktionen sammeln
            classes = [node for node in tree.body if isinstance(node, ast.ClassDef)]
            functions = [node for node in tree.body if isinstance(node, ast.FunctionDef)]
            imports = [node for node in tree.body if isinstance(node, (ast.Import, ast.ImportFrom))]
            
            ext = '.py' if as_py else '.txt'
            
            # Klassen speichern
            for cls in classes:
                start_line = cls.lineno - 1
                end_line = max(getattr(cls, "end_lineno", start_line), start_line)
                code = "\n".join(source.splitlines()[start_line:end_line])
                
                with open(os.path.join(outdir, f"{cls.name}{ext}"), "w", encoding="utf-8") as out:
                    out.write(code)
            
            # Hilfsfunktionen + Imports
            helper_lines = []
            lines = source.splitlines()
            
            for imp in imports:
                start, end = imp.lineno - 1, getattr(imp, "end_lineno", imp.lineno) - 1
                helper_lines.extend(lines[start:end+1])
            
            for func in functions:
                start, end = func.lineno - 1, getattr(func, "end_lineno", func.lineno) - 1
                helper_lines.extend(lines[start:end+1])
            
            # Restlicher Code
            occupied = set()
            for node in classes + functions + imports:
                occupied.update(range(node.lineno - 1, getattr(node, "end_lineno", node.lineno)))
            for i, line in enumerate(lines):
                if i not in occupied and line.strip():
                    helper_lines.append(line)
            
            if helper_lines:
                with open(os.path.join(outdir, f"Hilfsfunktionen{ext}"), "w", encoding="utf-8") as out:
                    out.write("\n".join(helper_lines))
            
            QMessageBox.information(
                self,
                "Erfolg",
                f"{len(classes)} Klassen extrahiert!\n\nOrdner: {outdir}"
            )
            
            self.perform_search()  # Refresh
        
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Zerlegung fehlgeschlagen:\n{str(e)}")
    
    def fix_python_encoding(self):
        """Repariert Encoding-Probleme in Python-Datei"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        
        # Prfe ob ftfy verfgbar ist
        try:
            from ftfy import fix_text
        except ImportError:
            QMessageBox.warning(
                self,
                "Modul fehlt",
                "Das Modul 'ftfy' ist nicht installiert.\n\nInstallieren Sie es mit:\npip install ftfy"
            )
            return
        
        try:
            # Lese Datei
            with open(path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            
            # Fixe Encoding
            fixed_content = fix_text(content)
            
            # Erstelle Backup
            backup_path = path + ".bak"
            shutil.copy(path, backup_path)
            
            # berschreibe Original
            with open(path, 'w', encoding='utf-8') as f:
                f.write(fixed_content)
            
            QMessageBox.information(
                self,
                "Erfolg",
                f"Encoding repariert!\n\nBackup: {backup_path}"
            )
            
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Encoding-Reparatur fehlgeschlagen:\n{str(e)}")
    
    def check_python_indentation(self):
        """Prüft Python-Datei auf Einrckungsfehler"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        
        try:
            import re
            
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                lines = f.readlines()
            
            errors = []
            for i, line in enumerate(lines):
                stripped = line.strip()
                indent_level = len(line) - len(line.lstrip())
                
                # Struktur ohne ':'
                if re.match(r"^(def|if|elif|else|for|while|try|except|class)\b", stripped) and not stripped.endswith(":"):
                    errors.append(f"Zeile {i+1}: Struktur ohne ':'")
                
                # Return außerhalb Block
                if stripped.startswith("return") and indent_level == 0:
                    errors.append(f"Zeile {i+1}: 'return' außerhalb von Block")
                
                # Tab/Space Mix
                if "\t" in line and " " in line[:line.find("\t")] if "\t" in line else False:
                    errors.append(f"Zeile {i+1}: Mischung aus Tab & Leerzeichen")
            
            if errors:
                # Zeige Fehler
                error_text = "\n".join(errors)
                msg = QMessageBox()
                msg.setIcon(QMessageBox.Icon.Warning)
                msg.setWindowTitle("Einrckungsfehler gefunden")
                msg.setText(f"{len(errors)} Fehler gefunden:")
                msg.setDetailedText(error_text)
                msg.exec()
            else:
                QMessageBox.information(self, "Erfolg", " Keine Einrckungsfehler gefunden!")
        
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Prüfung fehlgeschlagen:\n{str(e)}")
    
    def analyze_python_methods(self):
        """öffnet MethodenAnalyser3 als externes Tool"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        analyzer_path = os.path.join(os.path.dirname(__file__), "MethodenAnalyser3.py")
        
        # Prfe ob Tool existiert
        if not os.path.exists(analyzer_path):
            # Versuche im Projektordner
            analyzer_path = "/mnt/project/MethodenAnalyser3.py"
            if not os.path.exists(analyzer_path):
                QMessageBox.warning(
                    self,
                    "Tool nicht gefunden",
                    "MethodenAnalyser3.py wurde nicht gefunden.\n\nLegen Sie es neben die Hauptdatei."
                )
                return
        
        try:
            # Starte als separater Prozess
            subprocess.Popen([sys.executable, analyzer_path, path])
            QMessageBox.information(
                self,
                "MethodenAnalyser gestartet",
                f"Analyse läuft...\n\nDatei: {os.path.basename(path)}"
            )
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")
    
    def compile_to_exe(self):
        """öffnet Kompilator als externes Tool"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        compiler_path = os.path.join(os.path.dirname(__file__), "Kompilator.py")
        
        # Prfe ob Tool existiert
        if not os.path.exists(compiler_path):
            compiler_path = "/mnt/project/Kompilator.py"
            if not os.path.exists(compiler_path):
                QMessageBox.warning(
                    self,
                    "Tool nicht gefunden",
                    "Kompilator.py wurde nicht gefunden.\n\nLegen Sie es neben die Hauptdatei."
                )
                return
        
        try:
            # Starte als separater Prozess
            subprocess.Popen([sys.executable, compiler_path])
            QMessageBox.information(
                self,
                "Kompilator gestartet",
                f"Sie knnen nun die Datei auswählen:\n{os.path.basename(path)}"
            )
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")
    
    def open_sqlite_viewer(self):
        """öffnet SQLite Viewer für Datenbank"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        viewer_path = os.path.join(os.path.dirname(__file__), "SQLiteViewer.py")
        
        # Prfe ob Tool existiert
        if not os.path.exists(viewer_path):
            viewer_path = "/mnt/project/SQLiteViewer.py"
            if not os.path.exists(viewer_path):
                QMessageBox.warning(
                    self,
                    "Tool nicht gefunden",
                    "SQLiteViewer.py wurde nicht gefunden.\n\nLegen Sie es neben die Hauptdatei."
                )
                return
        
        try:
            # Starte als separater Prozess
            subprocess.Popen([sys.executable, viewer_path, path])
            QMessageBox.information(
                self,
                "SQLite Viewer gestartet",
                f"Datenbank geöffnet:\n{os.path.basename(path)}"
            )
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")

    def open_in_pythonbox(self):
        """öffnet .py-Datei in PythonBox (NEU V13!)"""
        selected = self.result_tree.currentItem()
        if not selected:
            return
        
        result = selected.data(0, Qt.ItemDataRole.UserRole)
        if not result:
            return
        
        path = result['path']
        
        # Hole PythonBox-Pfad aus Einstellungen
        pythonbox_path = self.settings.get("pythonbox_path", "")
        
        # Fallback: Versuche gängige Pfade
        if not pythonbox_path or not os.path.exists(pythonbox_path):
            possible_paths = [
                os.path.join(os.path.dirname(__file__), "PythonBox.py"),
                "/mnt/project/PythonBox.py",
                os.path.join(os.path.expanduser("~"), "PythonBox.py")
            ]
            
            for p in possible_paths:
                if os.path.exists(p):
                    pythonbox_path = p
                    break
        
        if not pythonbox_path or not os.path.exists(pythonbox_path):
            reply = QMessageBox.question(
                self,
                "PythonBox nicht gefunden",
                "PythonBox.py wurde nicht gefunden.\n\n"
                "Mchten Sie den Pfad jetzt in den Einstellungen festlegen?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            
            if reply == QMessageBox.StandardButton.Yes:
                # ffne Einstellungen
                parent = self.parent()
                while parent and not isinstance(parent, QMainWindow):
                    parent = parent.parent()
                if parent and hasattr(parent, 'show_settings'):
                    parent.show_settings()
            return
        
        try:
            # Starte PythonBox mit der Datei
            subprocess.Popen([sys.executable, pythonbox_path, path])
            QMessageBox.information(
                self,
                "PythonBox gestartet",
                f"PythonBox wurde geöffnet mit:\n{os.path.basename(path)}"
            )
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")
    
    # ========================================================================
    # NEUE DOKUMENTTYPEN (PHASE 8)
    # ========================================================================
    
    def create_material_reference(self):
        """Erstellt neuen Materialverweis"""
        dialog = MaterialReferenceDialog(self)
        if dialog.exec() == QDialog.DialogCode.Accepted:
            data = dialog.result_data
            
            # Dateiname generieren
            filename = f"Materialverweis_{data['bezeichnung']}.material.txt"
            # Ungltige Zeichen entfernen
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-'))
            
            # Speicherpfad wählen
            save_path, _ = QFileDialog.getSaveFileName(
                self,
                "Materialverweis speichern",
                filename,
                "Material-Dateien (*.material.txt);;Alle Dateien (*)"
            )
            
            if not save_path:
                return
            
            try:
                # Als Text-Datei speichern
                content = f"""Materialverweis: {data['bezeichnung']}
Erstellt: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
================================================================================

Bezeichnung:     {data['bezeichnung']}
Raum:            {data['raum']}
Regal:           {data['regal']}
Bereich:         {data['bereich']}
Inventarnummer:  {data['inventarnummer']}
Tags:            {data['tags']}

Beschreibung:
{data['beschreibung']}
"""
                
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                
                # In Datenbank indexieren
                # TODO: Spezielle Indexierung für Materialverweise
                
                QMessageBox.information(
                    self,
                    "Erfolg",
                    f" Materialverweis gespeichert:\n{os.path.basename(save_path)}"
                )
                
                # Refresh
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Speichern fehlgeschlagen:\n{str(e)}")
    
    def create_prompt_file(self):
        """Erstellt neue Prompt-Datei oder neue Version"""
        dialog = PromptFileDialog(self)
        if dialog.exec() == QDialog.DialogCode.Accepted:
            data = dialog.result_data
            
            # Dateiname generieren
            filename = f"{data['name']}.prompt"
            # Ungltige Zeichen entfernen
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-'))
            
            # Speicherpfad wählen
            save_path, _ = QFileDialog.getSaveFileName(
                self,
                "Prompt-Datei speichern",
                filename,
                "Prompt-Dateien (*.prompt);;Alle Dateien (*)"
            )
            
            if not save_path:
                return
            
            try:
                # Prfen ob Datei bereits existiert (für Versionierung)
                versions = []
                if os.path.exists(save_path):
                    # Existierende Versionen laden
                    try:
                        with open(save_path, 'r', encoding='utf-8') as f:
                            import json
                            versions = json.load(f)
                            if not isinstance(versions, list):
                                versions = [versions]
                    except:
                        versions = []
                
                # Neue Version hinzufügen
                versions.append(data)
                
                # Als JSON speichern
                with open(save_path, 'w', encoding='utf-8') as f:
                    json.dump(versions, f, indent=2, ensure_ascii=False)
                
                QMessageBox.information(
                    self,
                    "Erfolg",
                    f" Prompt-Datei gespeichert:\n{os.path.basename(save_path)}\n\nVersion: {data['version']}"
                )
                
                # Refresh
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Speichern fehlgeschlagen:\n{str(e)}")
    
    def create_internet_resource(self):
        """Erstellt neue Internetressource"""
        dialog = InternetResourceDialog(self)
        if dialog.exec() == QDialog.DialogCode.Accepted:
            data = dialog.result_data
            
            # Dateiname generieren
            filename = f"{data['bezeichnung']}.url"
            # Ungltige Zeichen entfernen
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-'))
            
            # Speicherpfad wählen
            save_path, _ = QFileDialog.getSaveFileName(
                self,
                "Internetressource speichern",
                filename,
                "URL-Verknüpfungen (*.url);;Alle Dateien (*)"
            )
            
            if not save_path:
                return
            
            try:
                # Als Windows .url Datei speichern
                content = f"""[InternetShortcut]
URL={data['adresse']}
IconIndex=0

[Metadata]
Bezeichnung={data['bezeichnung']}
Anbieter={data['anbieter']}
Tags={data['tags']}
Beschreibung={data['beschreibung']}
Erstellt={datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
"""
                
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                
                QMessageBox.information(
                    self,
                    "Erfolg",
                    f" Internetressource gespeichert:\n{os.path.basename(save_path)}"
                )
                
                # Refresh
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Speichern fehlgeschlagen:\n{str(e)}")
    
    def create_literature_reference(self):
        """Erstellt neuen Literaturverweis"""
        dialog = LiteratureReferenceDialog(self)
        if dialog.exec() == QDialog.DialogCode.Accepted:
            data = dialog.result_data
            
            # Dateiname generieren
            filename = f"Literaturverweis_{data['titel']}.txt"
            # Ungltige Zeichen entfernen
            filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-'))
            
            # Speicherpfad wählen
            save_path, _ = QFileDialog.getSaveFileName(
                self,
                "Literaturverweis speichern",
                filename,
                "Text-Dateien (*.txt);;Alle Dateien (*)"
            )
            
            if not save_path:
                return
            
            try:
                # Als strukturierte Text-Datei speichern
                content = f"""Literaturverweis: {data['titel']}
Erstellt: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
================================================================================

Titel:      {data['titel']}
Jahr:       {data['jahr']}
Autoren:    {data['autoren']}
Verlag:     {data['verlag']}
ISBN:       {data['isbn']}

Physischer Standort:
  Raum:     {data['raum']}
  Regal:    {data['regal']}

Online:     {data['internet']}
Tags:       {data['tags']}

Beschreibung/Notizen:
{data['beschreibung']}
"""
                
                with open(save_path, 'w', encoding='utf-8') as f:
                    f.write(content)
                
                QMessageBox.information(
                    self,
                    "Erfolg",
                    f" Literaturverweis gespeichert:\n{os.path.basename(save_path)}"
                )
                
                # Refresh
                self.perform_search()
                
            except Exception as e:
                QMessageBox.critical(self, "Fehler", f"Speichern fehlgeschlagen:\n{str(e)}")
    
    def import_browser_favorites(self):
        """Importiert Browser-Favoriten als Internetressourcen"""
        # HTML-Datei auswählen
        html_path, _ = QFileDialog.getOpenFileName(
            self,
            "Browser-Favoriten importieren",
            "",
            "HTML-Dateien (*.html *.htm);;Alle Dateien (*)"
        )
        
        if not html_path:
            return
        
        try:
            # Parse Bookmarks
            bookmarks = parse_browser_bookmarks(html_path)
            
            if not bookmarks:
                QMessageBox.warning(
                    self,
                    "Keine Favoriten",
                    "Keine Favoriten in der Datei gefunden.\n\nStellen Sie sicher, dass es eine Browser-Export-Datei ist."
                )
                return
            
            # Zielordner auswählen
            output_dir = QFileDialog.getExistingDirectory(
                self,
                "Zielordner für Internetressourcen auswählen"
            )
            
            if not output_dir:
                return
            
            # Progress Dialog
            progress = QMessageBox(self)
            progress.setWindowTitle("Import läuft...")
            progress.setText(f"Importiere {len(bookmarks)} Favoriten...")
            progress.setStandardButtons(QMessageBox.StandardButton.NoButton)
            progress.show()
            QApplication.processEvents()
            
            # Erstelle .url Dateien
            success_count = 0
            for bookmark in bookmarks:
                try:
                    title = bookmark.get('title', 'Unbenannt')
                    url = bookmark.get('url', '')
                    
                    if not url:
                        continue
                    
                    # Dateiname generieren
                    filename = f"{title}.url"
                    # Ungltige Zeichen entfernen
                    filename = "".join(c for c in filename if c.isalnum() or c in (' ', '.', '_', '-'))
                    
                    # Pfad
                    file_path = os.path.join(output_dir, filename)
                    
                    # URL-Datei erstellen
                    content = f"""[InternetShortcut]
URL={url}
IconIndex=0

[Metadata]
Bezeichnung={title}
Importiert={datetime.now().strftime('%Y-%m-%d %H:%M:%S')}
Quelle=Browser-Favoriten
"""
                    
                    with open(file_path, 'w', encoding='utf-8') as f:
                        f.write(content)
                    
                    success_count += 1
                
                except Exception as e:
                    continue
            
            progress.close()
            
            QMessageBox.information(
                self,
                "Import abgeschlossen",
                f" {success_count} von {len(bookmarks)} Favoriten importiert!\n\n"
                f"Ordner: {output_dir}"
            )
            
            # Refresh
            self.perform_search()
        
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Import fehlgeschlagen:\n{str(e)}")


# ============================================================================
# 8B. VERBINDUNGEN-WIDGET (Themen-basierte Verwaltung)
# ============================================================================



    # ============================================================================
    # PHASE 3-7: POOLING, KONVERTIERUNG, MERKEN (VOLLSTäNDIG)
    # ============================================================================
    
    def handle_tree_drop(self, event):
        """Verarbeitet Drag & Drop zwischen Dateien"""
        target_item = self.result_tree.itemAt(event.pos())
        if not target_item:
            event.ignore()
            return
        
        source_items = self.result_tree.selectedItems()
        if not source_items:
            event.ignore()
            return
        
        source_files = []
        for item in source_items:
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if result and 'path' in result:
                source_files.append(result['path'])
        
        target_result = target_item.data(0, Qt.ItemDataRole.UserRole)
        if not target_result or 'path' not in target_result:
            event.ignore()
            return
        
        target_file = target_result['path']
        
        if target_file in source_files:
            QMessageBox.warning(self, "Fehler", "Kann Datei nicht auf sich selbst droppen")
            event.ignore()
            return
        
        self.pool_files(source_files, target_file)
        event.accept()
    
    def pool_files(self, source_files, target_file):
        """Hauptlogik für Datei-Pooling"""
        if not source_files:
            return
        
        target_ext = os.path.splitext(target_file)[1].lower()
        target_base = os.path.splitext(target_file)[0]
        
        pool_count = len(source_files) + 1
        
        output_format_mode = self.settings.get("pooling_target_format", "target")
        if output_format_mode == "target":
            output_ext = target_ext
        elif output_format_mode == "pdf":
            output_ext = ".pdf"
        elif output_format_mode == "txt":
            output_ext = ".txt"
        elif output_format_mode == "docx":
            output_ext = ".docx"
        elif output_format_mode == "pptx":
            output_ext = ".pptx"
        else:
            output_ext = target_ext
        
        output_file = f"{target_base}_{pool_count}_pooled{output_ext}"
        
        if os.path.exists(output_file):
            reply = QMessageBox.question(
                self, "Datei existiert",
                f"Datei existiert bereits:\n{os.path.basename(output_file)}\n\nberschreiben?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.No:
                return
        
        sort_order = self.settings.get("pooling_sort_order", "source_first")
        if sort_order == "target_first":
            all_files = [target_file] + source_files
        else:
            all_files = source_files + [target_file]
        
        try:
            if output_ext in ['.txt', '.md', '.log']:
                self._pool_text(all_files, output_file)
            elif output_ext == '.pdf' and HAS_PDF:
                self._pool_to_pdf(all_files, output_file)
            elif output_ext == '.docx' and HAS_DOCX:
                self._pool_to_docx(all_files, output_file)
            elif output_ext in ['.ppt', '.pptx']:
                self._pool_to_pptx(all_files, output_file)
            elif output_ext in ['.mp3', '.wav', '.flac', '.aac', '.ogg', '.m4a']:
                self._pool_audio(all_files, output_file)
            elif output_ext in ['.mp4', '.mkv', '.avi', '.mov', '.wmv']:
                self._pool_video(all_files, output_file)
            else:
                QMessageBox.warning(self, "Format nicht unterstützt",
                    f"Pooling für Format '{output_ext}' noch nicht implementiert")
                return
            
            QMessageBox.information(self, "Pooling erfolgreich",
                f" Datei erstellt:\n{os.path.basename(output_file)}")
            self.perform_search()
        
        except Exception as e:
            QMessageBox.critical(self, "Pooling fehlgeschlagen",
                f"Fehler beim Poolen:\n\n{str(e)}")
    
    def _pool_text(self, files, output_file):
        """Kombiniert Textdateien"""
        combined_content = []
        separator = self.settings.get("pooling_text_separator", "\n\n--- NEXT FILE ---\n\n")
        
        for filepath in files:
            try:
                with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                    combined_content.append(f"=== {os.path.basename(filepath)} ===\n")
                    combined_content.append(content)
            except Exception as e:
                combined_content.append(f"[FEHLER beim Lesen von {os.path.basename(filepath)}: {e}]\n")
        
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(separator.join(combined_content))
    
    def _pool_to_pdf(self, files, output_file):
        """Kombiniert Dateien zu PDF"""
        if not HAS_PDF:
            raise Exception("PyPDF2 nicht installiert")
        
        writer = PdfWriter()
        
        for filepath in files:
            ext = os.path.splitext(filepath)[1].lower()
            
            try:
                if ext == '.pdf':
                    reader = PdfReader(filepath)
                    for page in reader.pages:
                        writer.add_page(page)
                else:
                    # Konvertiere zu PDF falls möglich
                    temp_pdf = f"{filepath}_temp.pdf"
                    if ext in ['.txt', '.md', '.log']:
                        # Einfache Text2PDF
                        from reportlab.pdfgen import canvas
                        from reportlab.lib.pagesizes import letter
                        
                        c = canvas.Canvas(temp_pdf, pagesize=letter)
                        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                            text = f.read()
                        
                        y = 750
                        for line in text.split('\n'):
                            if y < 50:
                                c.showPage()
                                y = 750
                            c.drawString(50, y, line[:80])
                            y -= 15
                        c.save()
                        
                        reader = PdfReader(temp_pdf)
                        for page in reader.pages:
                            writer.add_page(page)
                        os.remove(temp_pdf)
            
            except Exception as e:
                print(f"Fehler bei {filepath}: {e}")
        
        with open(output_file, 'wb') as f:
            writer.write(f)
    
    def _pool_to_docx(self, files, output_file):
        """Kombiniert zu Word-Dokument"""
        if not HAS_DOCX:
            raise Exception("python-docx nicht installiert")
        
        doc = docx.Document()
        
        for filepath in files:
            ext = os.path.splitext(filepath)[1].lower()
            
            doc.add_heading(os.path.basename(filepath), level=1)
            
            try:
                if ext == '.docx':
                    source_doc = docx.Document(filepath)
                    for element in source_doc.element.body:
                        doc.element.body.append(element)
                
                elif ext in ['.txt', '.md', '.log', '.py']:
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content = f.read()
                    doc.add_paragraph(content)
                
                doc.add_page_break()
            
            except Exception as e:
                doc.add_paragraph(f"[FEHLER: {e}]")
        
        doc.save(output_file)
    
    def _pool_to_pptx(self, files, output_file):
        """Kombiniert zu PowerPoint"""
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise Exception("python-pptx nicht installiert")
        
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        for filepath in files:
            ext = os.path.splitext(filepath)[1].lower()
            
            try:
                if ext in ['.ppt', '.pptx']:
                    source_prs = Presentation(filepath)
                    for slide in source_prs.slides:
                        prs.slides.add_slide(slide.slide_layout)
                
                elif ext in ['.txt', '.md', '.log']:
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    title.text = os.path.basename(filepath)
                    
                    content = slide.placeholders[1]
                    with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                        content.text = f.read()[:500]
                
                elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                    slide = prs.slides.add_slide(prs.slide_layouts[6])
                    slide.shapes.add_picture(filepath, Inches(1), Inches(1), 
                                           width=Inches(8))
            
            except Exception as e:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = f"Fehler: {os.path.basename(filepath)}"
        
        prs.save(output_file)
    
    def _pool_audio(self, files, output_file):
        """Erstellt Audio-Playlist (M3U)"""
        gap_seconds = self.settings.get("pooling_audio_gap", 3)
        
        with open(output_file.replace(os.path.splitext(output_file)[1], '.m3u'), 
                 'w', encoding='utf-8') as f:
            f.write("#EXTM3U\n")
            for filepath in files:
                f.write(f"#EXTINF:-1,{os.path.basename(filepath)}\n")
                f.write(f"{filepath}\n")
    
    def _pool_video(self, files, output_file):
        """Erstellt Video-Playlist (M3U)"""
        gap_seconds = self.settings.get("pooling_video_gap", 1)
        
        with open(output_file.replace(os.path.splitext(output_file)[1], '.m3u'), 
                 'w', encoding='utf-8') as f:
            f.write("#EXTM3U\n")
            for filepath in files:
                f.write(f"#EXTINF:-1,{os.path.basename(filepath)}\n")
                f.write(f"{filepath}\n")
    
    def convert_text_file(self, target_format):
        """Konvertiert Textdateien zwischen Formaten"""
        selected_items = self.result_tree.selectedItems()
        if not selected_items:
            return
        
        success_count = 0
        
        for item in selected_items:
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if not result or 'path' not in result:
                continue
            
            filepath = result['path']
            ext = os.path.splitext(filepath)[1].lower()
            
            if ext not in ['.txt', '.docx', '.rtf', '.md', '.log', '.py', '.odt']:
                continue
            
            try:
                content = self._read_text_content(filepath, ext)
                
                base = os.path.splitext(filepath)[0]
                output_file = f"{base}_{ext[1:]}2{target_format}.{target_format}"
                
                self._write_text_content(content, output_file, target_format)
                success_count += 1
            
            except Exception as e:
                QMessageBox.warning(self, "Fehler", 
                    f"Konvertierung fehlgeschlagen für {os.path.basename(filepath)}:\n{e}")
        
        if success_count > 0:
            QMessageBox.information(self, "Erfolgreich", 
                f"{success_count} Datei(en) konvertiert")
            self.perform_search()
    
    def _read_text_content(self, filepath, ext):
        """Liest Text aus verschiedenen Formaten"""
        if ext == '.docx':
            if not HAS_DOCX:
                raise Exception("python-docx nicht installiert")
            doc = docx.Document(filepath)
            return '\n'.join([para.text for para in doc.paragraphs])
        
        elif ext == '.rtf':
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                rtf_content = f.read()
            import re
            text = re.sub(r'\{[^}]*\}', '', rtf_content)
            text = re.sub(r'\[a-z]+', '', text)
            return text.strip()
        
        else:  # txt, md, log, py, etc.
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()
    
    def _write_text_content(self, content, output_file, target_format):
        """Schreibt Text in verschiedene Formate"""
        if target_format == 'docx':
            if not HAS_DOCX:
                raise Exception("python-docx nicht installiert")
            doc = docx.Document()
            doc.add_paragraph(content)
            doc.save(output_file)
        
        elif target_format == 'rtf':
            rtf_content = r"{\rtf1\ansi\deff0 {\fonttbl {\f0 Arial;}}" + "\n"
            rtf_content += r"\f0\fs24 " + content.replace('\n', '\\par ') + "}"
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(rtf_content)
        
        elif target_format == 'pdf':
            try:
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                
                c = canvas.Canvas(output_file, pagesize=letter)
                y = 750
                for line in content.split('\n'):
                    if y < 50:
                        c.showPage()
                        y = 750
                    c.drawString(50, y, line[:80])
                    y -= 15
                c.save()
            except ImportError:
                raise Exception("reportlab nicht installiert")
        
        else:  # txt, md, etc.
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(content)
    
    def convert_image_file(self, target_format):
        """Konvertiert Bilddateien"""
        selected_items = self.result_tree.selectedItems()
        if not selected_items:
            return
        
        success_count = 0
        
        for item in selected_items:
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if not result or 'path' not in result:
                continue
            
            filepath = result['path']
            ext = os.path.splitext(filepath)[1].lower()
            
            if ext not in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', '.ico']:
                continue
            
            try:
                from PIL import Image
                
                img = Image.open(filepath)
                
                base = os.path.splitext(filepath)[0]
                output_file = f"{base}_{ext[1:]}2{target_format}.{target_format}"
                
                if target_format == 'jpg' or target_format == 'jpeg':
                    if img.mode in ('RGBA', 'LA', 'P'):
                        background = Image.new('RGB', img.size, (255, 255, 255))
                        if img.mode == 'P':
                            img = img.convert('RGBA')
                        background.paste(img, mask=img.split()[-1])
                        img = background
                    img.save(output_file, 'JPEG', quality=95)
                
                elif target_format == 'ico':
                    sizes = [(16,16), (32,32), (48,48), (256,256)]
                    img.save(output_file, format='ICO', sizes=sizes)
                
                else:
                    img.save(output_file, target_format.upper())
                
                success_count += 1
            
            except Exception as e:
                QMessageBox.warning(self, "Fehler",
                    f"Konvertierung fehlgeschlagen für {os.path.basename(filepath)}:\n{e}")
        
        if success_count > 0:
            QMessageBox.information(self, "Erfolgreich",
                f"{success_count} Bild(er) konvertiert")
            self.perform_search()
    
    def convert_ppt_to_pdf(self):
        """Konvertiert PowerPoint zu PDF"""
        selected_items = self.result_tree.selectedItems()
        if not selected_items:
            return
        
        for item in selected_items:
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if not result or 'path' not in result:
                continue
            
            filepath = result['path']
            ext = os.path.splitext(filepath)[1].lower()
            
            if ext not in ['.ppt', '.pptx']:
                continue
            
            try:
                # Verwendet LibreOffice/OpenOffice für Konvertierung
                base = os.path.splitext(filepath)[0]
                output_file = f"{base}_ppt2pdf.pdf"
                
                subprocess.run([
                    'soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', os.path.dirname(filepath),
                    filepath
                ], check=True)
                
                QMessageBox.information(self, "Erfolgreich",
                    f"PDF erstellt:\n{os.path.basename(output_file)}")
                self.perform_search()
            
            except Exception as e:
                QMessageBox.warning(self, "Fehler",
                    f"Konvertierung fehlgeschlagen:\n{e}")
    
    def remember_files(self):
        """Merkt ausgewählte Dateien"""
        selected_items = self.result_tree.selectedItems()
        if not selected_items:
            return
        
        for item in selected_items:
            result = item.data(0, Qt.ItemDataRole.UserRole)
            if result and 'path' in result:
                filepath = result['path']
                if filepath not in self.remembered_files:
                    self.remembered_files.append(filepath)
        
        self.update_remembered_count()
        QMessageBox.information(self, "Gemerkt",
            f"{len(self.remembered_files)} Datei(en) in Sammlung")
    
    def show_remembered_files(self):
        """Zeigt Dialog mit gemerkten Dateien"""
        if not self.remembered_files:
            QMessageBox.information(self, "Sammlung leer",
                "Keine Dateien gemerkt")
            return
        
        dialog = QDialog(self)
        dialog.setWindowTitle(f"Gemerkte Dateien ({len(self.remembered_files)})")
        dialog.resize(600, 400)
        
        layout = QVBoxLayout(dialog)
        
        list_widget = QListWidget()
        for filepath in self.remembered_files:
            list_widget.addItem(os.path.basename(filepath))
        layout.addWidget(list_widget)
        
        btn_layout = QHBoxLayout()
        
        btn_copy = QPushButton(" Alle kopieren")
        btn_copy.clicked.connect(lambda: self.copy_remembered_files() or dialog.accept())
        
        btn_pool = QPushButton(" Zusammenfhren")
        btn_pool.clicked.connect(lambda: self.pool_remembered_files() or dialog.accept())
        
        btn_clear = QPushButton("Leeren")
        btn_clear.clicked.connect(lambda: self.clear_remembered_files() or dialog.accept())
        
        btn_close = QPushButton("Schließen")
        btn_close.clicked.connect(dialog.accept)
        
        btn_layout.addWidget(btn_copy)
        btn_layout.addWidget(btn_pool)
        btn_layout.addWidget(btn_clear)
        btn_layout.addWidget(btn_close)
        
        layout.addLayout(btn_layout)
        dialog.exec()
    
    def copy_remembered_files(self):
        """Kopiert gemerkte Dateien"""
        if not self.remembered_files:
            return
        
        target_dir = QFileDialog.getExistingDirectory(self, "Zielordner wählen")
        if not target_dir:
            return
        
        success_count = 0
        for filepath in self.remembered_files:
            try:
                shutil.copy2(filepath, target_dir)
                success_count += 1
            except Exception as e:
                print(f"Fehler bei {filepath}: {e}")
        
        QMessageBox.information(self, "Kopiert",
            f"{success_count} Datei(en) nach\n{target_dir}")
    
    def pool_remembered_files(self):
        """Fhrt gemerkte Dateien zusammen"""
        if len(self.remembered_files) < 2:
            QMessageBox.warning(self, "Zu wenige Dateien",
                "Mindestens 2 Dateien zum Pooling ntig")
            return
        
        # Nimm erste Datei als Target
        target = self.remembered_files[0]
        sources = self.remembered_files[1:]
        
        self.pool_files(sources, target)
    
    def clear_remembered_files(self):
        """Leert Sammlung"""
        self.remembered_files = []
        self.update_remembered_count()
        QMessageBox.information(self, "Geleert", "Sammlung wurde geleert")
    
    def update_remembered_count(self):
        """Aktualisiert Counter im Button"""
        if hasattr(self, 'btn_remembered'):
            count = len(self.remembered_files)
            self.btn_remembered.setText(f" Gemerkte Dateien ({count})")


class ConnectionsWidget(QWidget):
    """Widget zur Verwaltung von Verbindungen/Themen mit Indizierung"""
    
    def __init__(self, parent=None):
        super().__init__(parent)
        
        # Config Manager
        config_path = os.path.join(
            os.path.expanduser("~"),
            ".profiler_suite",
            "connections.json"
        )
        self.cfg = ConnectionConfigManager(config_path)
        self.worker = None
        
        self.init_ui()
        self.populate_list()
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Toolbar
        toolbar = QHBoxLayout()
        
        btn_new = QPushButton(" Neue Verbindung")
        btn_new.clicked.connect(self.add_connection)
        
        btn_refresh = QPushButton(" Aktualisieren")
        btn_refresh.clicked.connect(self.populate_list)
        
        toolbar.addWidget(btn_new)
        toolbar.addWidget(btn_refresh)
        toolbar.addStretch()
        
        layout.addLayout(toolbar)
        
        # Liste der Verbindungen
        self.list_widget = QListWidget()
        self.list_widget.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.list_widget.customContextMenuRequested.connect(self.show_context_menu)
        layout.addWidget(self.list_widget)
        
        # Controls
        ctrl_layout = QHBoxLayout()
        
        self.btn_index = QPushButton(" Indizieren")
        self.btn_index.clicked.connect(self.start_indexing)
        
        self.btn_export = QPushButton(" Dateiliste exportieren")
        self.btn_export.clicked.connect(self.export_file_list)
        
        self.btn_pause = QPushButton(" Pause")
        self.btn_pause.setEnabled(False)
        self.btn_pause.clicked.connect(self.toggle_pause)
        
        self.btn_stop = QPushButton(" Stop")
        self.btn_stop.setEnabled(False)
        self.btn_stop.clicked.connect(self.stop_worker)
        
        ctrl_layout.addWidget(self.btn_index)
        ctrl_layout.addWidget(self.btn_export)
        ctrl_layout.addWidget(self.btn_pause)
        ctrl_layout.addWidget(self.btn_stop)
        
        layout.addLayout(ctrl_layout)
        
        # Status
        self.lbl_status = QLabel("Bereit")
        self.progress = QProgressBar()
        
        layout.addWidget(self.lbl_status)
        # Status für enabled Verbindungen
        self.status_label = QLabel("Bereit")
        self.status_label.setStyleSheet("color: #888; font-size: 11px; padding: 5px;")
        layout.addWidget(self.status_label)
        
        layout.addWidget(self.progress)
    
    def populate_list(self):
        """Ldt alle Verbindungen"""
        self.list_widget.clear()
        
        enabled_count = 0
        total_count = len(self.cfg.list_connections())
        
        for conn in self.cfg.list_connections():
            sources_count = len(conn.get("sources", []))
            
            # Icons basierend auf Status
            enabled = conn.get("enabled", True)
            icon = "" if enabled else ""
            auto_update = conn.get("auto_update", False)
            auto_icon = " " if auto_update else ""
            
            if enabled:
                enabled_count += 1
            
            item_text = f"{icon} {conn['name']}{auto_icon} ({sources_count} Ordner)"
            
            item = QListWidgetItem(item_text)
            item.setData(Qt.ItemDataRole.UserRole, conn)
            self.list_widget.addItem(item)
        
        # Update Status-Label
        if hasattr(self, 'status_label'):
            self.status_label.setText(f"{enabled_count}/{total_count} Verbindungen aktiv")
    
    def add_connection(self):
        """Neue Verbindung erstellen"""
        dlg = ConnectionDialog(self)
        if dlg.exec():
            result = dlg.get_result()
            if result:
                self.cfg.add_or_update_connection(result)
                self.populate_list()
    
    def show_context_menu(self, pos):
        """Kontextmenfr Verbindungen"""
        item = self.list_widget.itemAt(pos)
        if not item:
            return
        
        menu = QMenu()
        
        conn = item.data(Qt.ItemDataRole.UserRole)
        enabled = conn.get("enabled", True)
        auto_update = conn.get("auto_update", False)
        
        # Toggle Enabled
        act_toggle = menu.addAction(" Deaktivieren" if enabled else " Aktivieren")
        
        # Toggle Auto-Update
        act_auto = menu.addAction(" Auto-Update: AUS" if auto_update else " Auto-Update: EIN")
        
        menu.addSeparator()
        
        act_edit = menu.addAction(" Bearbeiten")
        act_delete = menu.addAction("Löschen")
        menu.addSeparator()
        act_open_db = menu.addAction("Datenbank öffnen")
        
        action = menu.exec(self.list_widget.viewport().mapToGlobal(pos))
        conn = item.data(Qt.ItemDataRole.UserRole)
        
        # Handle Toggle Enabled
        if action == act_toggle:
            new_enabled = not enabled
            self.cfg.toggle_connection(conn["id"], new_enabled)
            self.populate_list()
            status = "aktiviert" if new_enabled else "deaktiviert"
            self.lbl_status.setText(f"Verbindung '{conn['name']}' {status}")
            return
        
        # Handle Toggle Auto-Update
        elif action == act_auto:
            new_auto = not auto_update
            self.cfg.toggle_auto_update(conn["id"], new_auto)
            self.populate_list()
            status = "aktiviert" if new_auto else "deaktiviert"
            self.lbl_status.setText(f"Auto-Update für '{conn['name']}' {status}")
            return
        
        if action == act_delete:
            reply = QMessageBox.question(
                self,
                "Verbindung löschen",
                f"Wirklich '{conn['name']}' löschen?\n\nDie Datenbank-Datei bleibt erhalten.",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            if reply == QMessageBox.StandardButton.Yes:
                self.cfg.remove_connection(conn["id"])
                self.populate_list()
        
        elif action == act_edit:
            dlg = ConnectionDialog(self, conn)
            if dlg.exec():
                result = dlg.get_result()
                if result:
                    self.cfg.add_or_update_connection(result)
                    self.populate_list()
        
        elif action == act_open_db:
            db_path = conn.get("db_path", "")
            if os.path.exists(db_path):
                # SQLite Viewer starten (falls vorhanden)
                viewer_path = find_tool_path("SQLiteViewer.py")
                if viewer_path and os.path.exists(viewer_path):
                    try:
                        subprocess.Popen([sys.executable, viewer_path, db_path])
                    except:
                        pass
                else:
                    QMessageBox.information(
                        self,
                        "DB-Pfad",
                        f"Datenbank:\n{db_path}\n\n(SQLiteViewer nicht gefunden)"
                    )
            else:
                QMessageBox.warning(self, "Nicht gefunden", "Datenbank existiert noch nicht.")
    
    def start_indexing(self):
        """Startet Indizierung der ausgewählten Verbindung"""
        item = self.list_widget.currentItem()
        if not item:
            QMessageBox.warning(self, "Fehler", "Bitte erst eine Verbindung auswählen.")
            return
        
        conn = item.data(Qt.ItemDataRole.UserRole)
        
        # Worker starten
        self.worker = SyncWorker(conn, mode="index")
        self.worker.signals.status.connect(self.lbl_status.setText)
        self.worker.signals.progress.connect(self.on_progress)
        self.worker.signals.finished.connect(self.on_worker_finished)
        
        self.btn_index.setEnabled(False)
        self.btn_pause.setEnabled(True)
        self.btn_stop.setEnabled(True)
        
        self.worker.start()
    
    def toggle_pause(self):
        """Pausiert/Fortsetzt Worker"""
        if self.worker:
            self.worker.is_paused = not self.worker.is_paused
            if self.worker.is_paused:
                self.btn_pause.setText("Fortsetzen")
            else:
                self.btn_pause.setText(" Pause")
    
    def stop_worker(self):
        """Stoppt Worker"""
        if self.worker:
            self.worker.is_killed = True
            self.lbl_status.setText("Wird gestoppt...")
    
    def on_progress(self, value, text):
        """Progress Update"""
        self.progress.setValue(value)
        if text:
            self.lbl_status.setText(text)
    
    def on_worker_finished(self):
        """Worker fertig"""
        self.btn_index.setEnabled(True)
        self.btn_pause.setEnabled(False)
        self.btn_stop.setEnabled(False)
        self.btn_pause.setText(" Pause")
        self.progress.setValue(100)
    
    def export_file_list(self):
            """Exportiert Dateiliste als PDF"""
            item = self.list_widget.currentItem()
            if not item:
                QMessageBox.warning(self, "Fehler", "Bitte erst eine Verbindung auswählen.")
                return
            
            conn = item.data(Qt.ItemDataRole.UserRole)
            db_path = conn.get("db_path", "")
            
            if not os.path.exists(db_path):
                QMessageBox.warning(
                    self,
                    "Keine Datenbank",
                    "Datenbank existiert noch nicht.\nBitte erst indizieren."
                )
                return
            
            # PDF-Pfad wählen
            default_name = f"Dateiliste_{conn['name']}_{datetime.now().strftime('%Y%m%d')}.pdf"
            safe_name = "".join(c for c in default_name if c.isalnum() or c in ('_', '-', '.'))
            
            pdf_path, _ = QFileDialog.getSaveFileName(
                self,
                "Dateiliste exportieren",
                safe_name,
                "PDF-Dateien (*.pdf)"
            )
            
            if not pdf_path:
                return
            
            try:
                if not HAS_REPORTLAB:
                    QMessageBox.warning(
                        self,
                        "Bibliothek fehlt",
                        "reportlab ist nicht installiert.\n\nInstalliere mit: pip install reportlab"
                    )
                    return
                
                # Daten aus DB laden
                db_conn = sqlite3.connect(db_path)
                db_conn.row_factory = sqlite3.Row
                
                query = """
                    SELECT v.path, v.name, v.mtime, f.size
                    FROM versions v
                    JOIN files f ON v.file_id = f.id
                    WHERE v.is_deleted = 0
                    ORDER BY v.path
                """
                
                rows = db_conn.execute(query).fetchall()
                db_conn.close()
                
                # PDF erstellen
                doc = SimpleDocTemplate(pdf_path, pagesize=A4)
                story = []
                styles = getSampleStyleSheet()
                
                # Titel
                title = Paragraph(f"<b>Dateiliste: {conn['name']}</b>", styles['Title'])
                story.append(title)
                story.append(Spacer(1, 0.5*cm))
                
                # Info
                info = Paragraph(
                    f"Erstellt: {datetime.now().strftime('%d.%m.%Y %H:%M')}<br/>"
                    f"Quellordner: {len(conn.get('sources', []))}<br/>"
                    f"Dateien: {len(rows)}",
                    styles['Normal']
                )
                story.append(info)
                story.append(Spacer(1, 1*cm))
                
                # Dateiliste
                for row in rows:
                    # XML-Zeichen im Pfad escapen, da Paragraph XML-Tags erwartet
                    safe_path = str(row['path']).replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    path_text = Paragraph(f"<b>{safe_path}</b>", styles['Normal'])
                    story.append(path_text)
                    story.append(Spacer(1, 0.2*cm))
                
                doc.build(story)
                
                QMessageBox.information(
                    self,
                    "Export erfolgreich",
                    f"✅ Dateiliste exportiert:\n{pdf_path}\n\n{len(rows)} Dateien"
                )
            
            except Exception as e:
                import traceback
                traceback.print_exc()
                QMessageBox.critical(self, "Fehler", f"Export fehlgeschlagen:\n{str(e)}")


# ============================================================================
# 8.5. AUTO-SYNC WATCHDOG (V14.3)
# ============================================================================

class AutoSyncHandler(FileSystemEventHandler if HAS_WATCHDOG else object):
    """Handler für Dateisystem-Events mit Auto-Sync Logik"""
    
    def __init__(self, target_folder, extensions=None, callback=None):
        if HAS_WATCHDOG:
            super().__init__()
        self.target = target_folder
        self.extensions = extensions or ['.pdf', '.docx', '.xlsx', '.txt', '.jpg', '.png']
        self.callback = callback
        self.synced_files = set()  # Duplikat-Schutz
    
    def on_created(self, event):
        if event.is_directory:
            return
        if self._should_sync(event.src_path):
            self._sync_file(event.src_path)
    
    def on_modified(self, event):
        if event.is_directory:
            return
        if self._should_sync(event.src_path):
            self._sync_file(event.src_path)
    
    def _should_sync(self, path):
        return any(path.lower().endswith(ext) for ext in self.extensions)
    
    def _sync_file(self, src_path):
        """Kopiert Datei zum Zielordner mit Duplikat-Handling"""
        try:
            filename = os.path.basename(src_path)
            dest_path = os.path.join(self.target, filename)
            
            # Duplikat-Check mit Timestamp
            file_key = f"{src_path}_{os.path.getmtime(src_path)}"
            if file_key in self.synced_files:
                return
            self.synced_files.add(file_key)
            
            # Falls Ziel existiert: Umbenennen
            if os.path.exists(dest_path):
                name, ext = os.path.splitext(filename)
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                dest_path = os.path.join(self.target, f"{name}_{timestamp}{ext}")
            
            shutil.copy2(src_path, dest_path)
            
            if self.callback:
                self.callback(f"✅ Synchronisiert: {filename}")
                
        except Exception as e:
            if self.callback:
                self.callback(f"❌ Fehler: {e}")


class AutoSyncManager(QObject):
    """Manager für mehrere Watch-Ordner mit UI-Integration"""
    
    status_changed = pyqtSignal(str)
    file_synced = pyqtSignal(str)
    
    def __init__(self, parent=None):
        super().__init__(parent)
        self.observers = {}  # source_path -> Observer
        self.handlers = {}   # source_path -> AutoSyncHandler
        self.stats = {"synced": 0, "errors": 0}
    
    def start_watch(self, source, target, extensions=None):
        """Startet Überwachung eines Ordners"""
        if not HAS_WATCHDOG:
            self.status_changed.emit("❌ watchdog nicht installiert!")
            return False
        
        if source in self.observers:
            self.status_changed.emit(f"⚠️ {source} wird bereits überwacht")
            return False
        
        if not os.path.isdir(source):
            self.status_changed.emit(f"❌ Quellordner existiert nicht: {source}")
            return False
        
        if not os.path.isdir(target):
            try:
                os.makedirs(target, exist_ok=True)
            except Exception as e:
                self.status_changed.emit(f"❌ Zielordner Fehler: {e}")
                return False
        
        handler = AutoSyncHandler(target, extensions, self._on_file_synced)
        observer = Observer()
        observer.schedule(handler, source, recursive=False)
        observer.start()
        
        self.observers[source] = observer
        self.handlers[source] = handler
        self.status_changed.emit(f"✅ Überwachung gestartet: {source}")
        return True
    
    def stop_watch(self, source=None):
        """Stoppt Überwachung (alle oder spezifisch)"""
        if source:
            if source in self.observers:
                self.observers[source].stop()
                self.observers[source].join()
                del self.observers[source]
                del self.handlers[source]
                self.status_changed.emit(f"⏹️ Überwachung gestoppt: {source}")
        else:
            for obs in self.observers.values():
                obs.stop()
                obs.join()
            self.observers.clear()
            self.handlers.clear()
            self.status_changed.emit("⏹️ Alle Überwachungen gestoppt")
    
    def _on_file_synced(self, message):
        self.stats["synced"] += 1
        self.file_synced.emit(message)
    
    def get_status(self):
        """Gibt Status aller Watches zurück"""
        return {src: obs.is_alive() for src, obs in self.observers.items()}
    
    def get_stats(self):
        return self.stats.copy()


class AutoSyncWidget(QWidget):
    """UI-Widget für Auto-Sync Konfiguration und Status"""
    
    def __init__(self, settings, parent=None):
        super().__init__(parent)
        self.settings = settings
        self.manager = AutoSyncManager(self)
        self.init_ui()
        self.load_config()
        
        # Signale verbinden
        self.manager.status_changed.connect(self.log_message)
        self.manager.file_synced.connect(self.log_message)
    
    def init_ui(self):
        layout = QVBoxLayout(self)
        
        # Status-Gruppe
        status_group = QGroupBox("Auto-Sync Status")
        status_layout = QVBoxLayout(status_group)
        
        # Watchdog-Prüfung
        if HAS_WATCHDOG:
            self.status_label = QLabel("✅ Watchdog verfügbar - Bereit")
        else:
            self.status_label = QLabel("❌ watchdog nicht installiert!\nInstallation: pip install watchdog")
        status_layout.addWidget(self.status_label)
        
        layout.addWidget(status_group)
        
        # Ordner-Konfiguration
        folder_group = QGroupBox("Ordner-Konfiguration")
        folder_layout = QFormLayout(folder_group)
        
        # Quellordner
        source_row = QHBoxLayout()
        self.source_edit = QLineEdit()
        self.source_edit.setPlaceholderText("Quellordner auswählen...")
        source_btn = QPushButton("📁")
        source_btn.setFixedWidth(40)
        source_btn.clicked.connect(self.select_source)
        source_row.addWidget(self.source_edit)
        source_row.addWidget(source_btn)
        folder_layout.addRow("Quelle:", source_row)
        
        # Zielordner
        target_row = QHBoxLayout()
        self.target_edit = QLineEdit()
        self.target_edit.setPlaceholderText("Zielordner auswählen...")
        target_btn = QPushButton("📁")
        target_btn.setFixedWidth(40)
        target_btn.clicked.connect(self.select_target)
        target_row.addWidget(self.target_edit)
        target_row.addWidget(target_btn)
        folder_layout.addRow("Ziel:", target_row)
        
        # Extensions
        self.ext_edit = QLineEdit(".pdf, .docx, .xlsx, .txt, .jpg, .png")
        folder_layout.addRow("Extensions:", self.ext_edit)
        
        layout.addWidget(folder_group)
        
        # Buttons
        btn_layout = QHBoxLayout()
        self.start_btn = QPushButton("▶️ Starten")
        self.start_btn.clicked.connect(self.start_sync)
        self.stop_btn = QPushButton("⏹️ Stoppen")
        self.stop_btn.clicked.connect(self.stop_sync)
        self.stop_btn.setEnabled(False)
        btn_layout.addWidget(self.start_btn)
        btn_layout.addWidget(self.stop_btn)
        layout.addLayout(btn_layout)
        
        # Log
        log_group = QGroupBox("Aktivitäts-Log")
        log_layout = QVBoxLayout(log_group)
        self.log_text = QTextEdit()
        self.log_text.setReadOnly(True)
        self.log_text.setMaximumHeight(200)
        log_layout.addWidget(self.log_text)
        layout.addWidget(log_group)
        
        # Statistiken
        stats_group = QGroupBox("Statistiken")
        stats_layout = QHBoxLayout(stats_group)
        self.stats_label = QLabel("Synchronisiert: 0 | Fehler: 0")
        stats_layout.addWidget(self.stats_label)
        layout.addWidget(stats_group)
        
        layout.addStretch()
    
    def select_source(self):
        folder = QFileDialog.getExistingDirectory(self, "Quellordner wählen")
        if folder:
            self.source_edit.setText(folder)
    
    def select_target(self):
        folder = QFileDialog.getExistingDirectory(self, "Zielordner wählen")
        if folder:
            self.target_edit.setText(folder)
    
    def start_sync(self):
        source = self.source_edit.text().strip()
        target = self.target_edit.text().strip()
        
        if not source or not target:
            QMessageBox.warning(self, "Fehler", "Bitte Quell- und Zielordner angeben!")
            return
        
        ext_text = self.ext_edit.text().strip()
        extensions = [e.strip() for e in ext_text.split(",") if e.strip()]
        
        if self.manager.start_watch(source, target, extensions):
            self.start_btn.setEnabled(False)
            self.stop_btn.setEnabled(True)
            self.status_label.setText(f"🟢 Aktiv: Überwache {source}")
            self.save_config()
    
    def stop_sync(self):
        self.manager.stop_watch()
        self.start_btn.setEnabled(True)
        self.stop_btn.setEnabled(False)
        self.status_label.setText("⏹️ Gestoppt")
        self.update_stats()
    
    def log_message(self, message):
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.log_text.append(f"[{timestamp}] {message}")
        self.update_stats()
    
    def update_stats(self):
        stats = self.manager.get_stats()
        self.stats_label.setText(f"Synchronisiert: {stats['synced']} | Fehler: {stats['errors']}")
    
    def save_config(self):
        """Speichert Auto-Sync Konfiguration"""
        config = {
            "source": self.source_edit.text(),
            "target": self.target_edit.text(),
            "extensions": self.ext_edit.text()
        }
        self.settings.set("autosync_config", config)
    
    def load_config(self):
        """Lädt Auto-Sync Konfiguration"""
        config = self.settings.get("autosync_config", {})
        if config:
            self.source_edit.setText(config.get("source", ""))
            self.target_edit.setText(config.get("target", ""))
            self.ext_edit.setText(config.get("extensions", ".pdf, .docx, .xlsx, .txt, .jpg, .png"))


# ============================================================================
# 9. MAIN WINDOW
# ============================================================================

class UnifiedMainWindow(QMainWindow):
    """Haupt-Fenster mit System-Tray"""
    
    def __init__(self):
        super().__init__()
        
        self.setWindowTitle("ProFiler Suite V14.3 - Auto-Sync Watchdog")
        self.resize(1400, 900)
        
        # Managers
        self.search_config = SearchConfigManager()
        self.sync_config = SyncConfigManager(SYNC_CONFIG_PATH)
        self.settings = SettingsManager()
        
        # Auto-Cleanup
        #if self.settings.get("auto_cleanup_enabled", True):
            #self.perform_auto_cleanup()
        
        # UI
        self.init_ui()
        self.init_tray()
        self.apply_dark_theme()
    
    def init_ui(self):
        """Initialisiert UI"""
        central = QWidget()
        self.setCentralWidget(central)
        layout = QVBoxLayout(central)
        
        # Tabs
        tabs = QTabWidget()
        
        # Tab 1: Search
        self.search_widget = SearchWidgetHybrid(self.search_config, self.settings)
        tabs.addTab(self.search_widget, "Suche & Explorer")
        
        # Tab 2: Verbindungen/Themen
        self.connections_widget = ConnectionsWidget(self)
        tabs.addTab(self.connections_widget, "Verbindungen & Themen")
        
        # Tab 3: Auto-Sync (NEU V14.3)
        self.autosync_widget = AutoSyncWidget(self.settings, self)
        tabs.addTab(self.autosync_widget, "🔄 Auto-Sync")
        
        layout.addWidget(tabs)
        
        # Statusbar
        self.statusBar().showMessage("Bereit")
        
        # Menu
        menubar = self.menuBar()
        
        file_menu = menubar.addMenu("Datei")
        file_menu.addAction("🌐 Browser-Favoriten importieren...", self.import_browser_favorites)
        file_menu.addSeparator()
        file_menu.addAction("⚙️ Einstellungen", self.show_settings)
        file_menu.addSeparator()
        file_menu.addAction("❌ Beenden", self.close)
        
        
        tools_menu = menubar.addMenu("Tools")
        tools_menu.addAction("🚦 Datenschutzampel starten...", self.start_datenschutzampel)
        tools_menu.addAction("FormConstructor öffnen...", self.launch_form_constructor)

        help_menu = menubar.addMenu("Hilfe")
        help_menu.addAction("über", self.show_about)
    
    def init_tray(self):
        """System-Tray"""
        self.tray_icon = QSystemTrayIcon(self)
        self.tray_icon.setIcon(self.style().standardIcon(QStyle.StandardPixmap.SP_FileIcon))
        
        tray_menu = QMenu()
        tray_menu.addAction("Anzeigen", self.show)
        tray_menu.addAction("❌ Beenden", QApplication.quit)
        
        self.tray_icon.setContextMenu(tray_menu)
        self.tray_icon.activated.connect(self.on_tray_activated)
        self.tray_icon.show()
    
    def on_tray_activated(self, reason):
        """Tray-Icon Aktivierung"""
        if reason == QSystemTrayIcon.ActivationReason.DoubleClick:
            self.show()
            self.activateWindow()
    
    def closeEvent(self, event):
        """Minimiert zu Tray"""
        event.ignore()
        self.hide()
        self.tray_icon.showMessage(
            "ProFiler Suite",
            "Luft im Hintergrund. Doppelklick zum öffnen.",
            QSystemTrayIcon.MessageIcon.Information,
            2000
        )
    
    def import_browser_favorites(self):
        """Delegiert Browser-Favoriten Import an SearchWidget"""
        self.search_widget.import_browser_favorites()
    
    def show_settings(self):
        """Zeigt Einstellungen"""
        dialog = SettingsDialog(self.settings, self)
        dialog.exec()
    
    def show_about(self):
        """Über-Dialog"""
        QMessageBox.about(
            self,
            "ProFiler Suite V13.3",
            "ProFiler Suite V13.3 - Enhanced\n\n"
            "Features:\n"
            "• PDF-Verschlüsselung/Entschlüsselung\n"
            "• PDF-Auszüge erstellen\n"
            "• OCR-Texterkennung\n"
            "• Text aus PDF entfernen\n"
            "• Masterpasswort-Management\n\n"
            "© 2024"
        )
        
    def launch_form_constructor(self):
        """Startet FormConstructor (NEU V13!)"""
        # Hole Pfad aus Einstellungen
        fc_path = self.settings.get("formconstructor_path", "")
        
        # Fallback: Versuche gängige Pfade
        if not fc_path or not os.path.exists(fc_path):
            possible_paths = [
                os.path.join(os.path.dirname(__file__), "FormConstructor_V1_5.py"),
                "/mnt/project/FormConstructor_V1_5.py",
                os.path.join(os.path.dirname(__file__), "FormConstructor.py")
            ]
            
            for p in possible_paths:
                if os.path.exists(p):
                    fc_path = p
                    break
        
        if not fc_path or not os.path.exists(fc_path):
            reply = QMessageBox.question(
                self,
                "FormConstructor nicht gefunden",
                "FormConstructor_V1_5.py wurde nicht gefunden.\n\n"
                "Mchten Sie den Pfad jetzt in den Einstellungen festlegen?",
                QMessageBox.StandardButton.Yes | QMessageBox.StandardButton.No
            )
            
            if reply == QMessageBox.StandardButton.Yes:
                self.show_settings()
            return
        
        try:
            subprocess.Popen([sys.executable, fc_path])
            QMessageBox.information(
                self,
                "FormConstructor gestartet",
                "FormConstructor wurde erfolgreich geöffnet!"
            )
        except Exception as e:
            QMessageBox.critical(self, "Fehler", f"Start fehlgeschlagen:\n{str(e)}")
    
    def start_datenschutzampel(self):
        """Startet die Datenschutzampel als separate Anwendung"""
        # 1. Suche im aktuellen Verzeichnis
        current_dir = os.path.dirname(os.path.abspath(__file__))
        ampel_path = os.path.join(current_dir, "ProFiler_Datenschutzampel.py")
        
        # 2. Alternative Pfade prüfen, falls nicht gefunden
        if not os.path.exists(ampel_path):
            fallback_path = "/mnt/user-data/outputs/ProFiler_Datenschutzampel.py"
            # Nur überschreiben, wenn der Fallback tatsächlich existiert
            if os.path.exists(fallback_path):
                ampel_path = fallback_path
        
        # 3. Ausführen oder Fehler melden
        if os.path.exists(ampel_path):
            try:
                subprocess.Popen([sys.executable, ampel_path])
                QMessageBox.information(
                    self,
                    "Datenschutzampel gestartet",
                    "🚦 Die Datenschutzampel wurde als separate Anwendung gestartet.\n\n"
                    "Sie läuft jetzt im System-Tray parallel zu ProFiler."
                )
            except Exception as e:
                QMessageBox.critical(
                    self,
                    "Fehler",
                    f"Konnte Datenschutzampel nicht starten:\n{str(e)}"
                )
        else:
            QMessageBox.warning(
                self,
                "Nicht gefunden",
                "ProFiler_Datenschutzampel.py wurde nicht gefunden.\n\n"
                f"Gesucht in:\n{ampel_path}\n\n"
                "Bitte stelle sicher, dass die Datei existiert."
            )
    

    def perform_auto_cleanup(self):
        """Automatisches Aufrumen"""
        days = self.settings.get("trash_retention_days", 30)
        
        if days <= 0:
            return
        
        total_deleted = 0
        
        for db_path in self.search_config.dbs:
            if not os.path.exists(db_path):
                continue
            
            db = ConnectionDB(db_path)
            deleted = db.cleanup_old_deleted(days)
            db.close()
            
            total_deleted += deleted
        
        if total_deleted > 0:
            self.statusBar().showMessage(f"✅ Auto-Cleanup: {total_deleted} alte Dateien entfernt", 5000)
    
    def apply_dark_theme(self):
        """Dark Theme"""
        app = QApplication.instance()
        app.setStyle("Fusion")
        
        palette = QPalette()
        palette.setColor(QPalette.ColorRole.Window, QColor(53, 53, 53))
        palette.setColor(QPalette.ColorRole.WindowText, Qt.GlobalColor.white)
        palette.setColor(QPalette.ColorRole.Base, QColor(35, 35, 35))
        palette.setColor(QPalette.ColorRole.AlternateBase, QColor(53, 53, 53))
        palette.setColor(QPalette.ColorRole.ToolTipBase, QColor(25, 25, 25))
        palette.setColor(QPalette.ColorRole.ToolTipText, Qt.GlobalColor.white)
        palette.setColor(QPalette.ColorRole.Text, Qt.GlobalColor.white)
        palette.setColor(QPalette.ColorRole.Button, QColor(53, 53, 53))
        palette.setColor(QPalette.ColorRole.ButtonText, Qt.GlobalColor.white)
        palette.setColor(QPalette.ColorRole.Link, QColor(42, 130, 218))
        palette.setColor(QPalette.ColorRole.Highlight, QColor(42, 130, 218))
        palette.setColor(QPalette.ColorRole.HighlightedText, Qt.GlobalColor.black)
        
        app.setPalette(palette)


# ============================================================================
# 10. MAIN ENTRY POINT
# ============================================================================

def main():
    app = QApplication(sys.argv)
    app.setApplicationName("ProFiler Suite V9")
    
    # Check Dependencies
    warnings = []
    
    if not HAS_PDF:
        warnings.append("⚠️ PyPDF2 nicht installiert - PDF-Features deaktiviert")
    
    if not HAS_OCR:
        warnings.append("⚠️ pytesseract nicht installiert - OCR deaktiviert")
    
    if not HAS_PDF2IMAGE:
        warnings.append("⚠️ pdf2image nicht installiert - Text-Removal deaktiviert")
    
    if not HAS_DOCX:
        warnings.append("⚠️ python-docx nicht installiert - Word-Features deaktiviert")
    
    if warnings:
        msg = QMessageBox()
        msg.setIcon(QMessageBox.Icon.Warning)
        msg.setWindowTitle("Fehlende Abhängigkeiten")
        msg.setText("Einige Features sind nicht verfgbar:\n\n" + "\n".join(warnings))
        msg.setInformativeText("Installation:\npip install PyPDF2 pytesseract pdf2image python-docx Pillow")
        msg.exec()
    
    window = UnifiedMainWindow()
    window.show()
    
    sys.exit(app.exec())


if __name__ == "__main__":
    main()